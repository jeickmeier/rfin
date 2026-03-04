"""Generate a PowerPoint presentation from the PIK Structural Credit notebook.

Restructured for a mixed audience (quants + non-quants) with 'so what' callouts at
the top of each content slide. Methodology (Merton + Hazard Rate models) upfront in
the main deck. Academic references as footnotes and a dedicated References slide.
Core narrative: the PIK premium and what drives it.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette ──────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)   # bright blue
ACCENT_GOLD = RGBColor(0xE8, 0xA8, 0x38)   # gold
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY    = RGBColor(0x99, 0x99, 0x99)
TABLE_HDR   = RGBColor(0x00, 0x50, 0x8A)   # darker blue for table headers
TABLE_ALT   = RGBColor(0xF0, 0xF4, 0xF8)   # light blue-gray for alt rows
TABLE_BDR   = RGBColor(0xDD, 0xDD, 0xDD)
RED         = RGBColor(0xCC, 0x33, 0x33)
GREEN       = RGBColor(0x33, 0x99, 0x66)
BLACK       = RGBColor(0x00, 0x00, 0x00)
CALLOUT_BG  = RGBColor(0xFF, 0xF5, 0xE0)   # light gold for callout boxes
DARK_TEXT   = RGBColor(0x33, 0x33, 0x33)
FOOTNOTE_CLR = RGBColor(0x88, 0x88, 0x88)  # lighter gray for footnotes

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ───────────────────────────────────────────────────

def add_dark_bg(slide):
    """Fill slide with dark background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=WHITE, spacing=1.2, font_name="Calibri"):
    """Add textbox with multiple lines (list of (text, bold, color_override))."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, bold, clr = item, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(font_size * spacing * 0.3)
    return tf


def add_table(slide, left, top, width, height, headers, rows, col_widths=None,
              font_size=11, header_font_size=12):
    """Add a formatted table to the slide."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                          Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(header_font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = BLACK
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT

    return table


def add_accent_bar(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_callout_box(slide, text, top=0.9, left=0.5, width=12.3, height=0.55):
    """Add a gold 'so what' callout box below the title on white-bg slides."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CALLOUT_BG
    shape.line.fill.background()
    add_accent_bar(slide, left, top, 0.06, height, ACCENT_GOLD)
    txBox = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.05),
        Inches(width - 0.3), Inches(height - 0.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    p.font.bold = True
    p.font.italic = True
    p.font.name = "Calibri"
    return tf


def add_dark_callout(slide, text, top=0.9):
    """Add a callout on dark background slides."""
    add_accent_bar(slide, 0.5, top, 0.06, 0.45, ACCENT_GOLD)
    txBox = slide.shapes.add_textbox(
        Inches(0.7), Inches(top + 0.05),
        Inches(12), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_GOLD
    p.font.bold = True
    p.font.italic = True
    p.font.name = "Calibri"
    return tf


def add_footnote(slide, text, bottom=7.15, color=FOOTNOTE_CLR):
    """Add a footnote citation at the bottom of a slide."""
    add_textbox(slide, 0.5, bottom, 12.3, 0.3, text,
                font_size=9, color=color, font_name="Calibri")


# ════════════════════════════════════════════════════════════════════════
#                        MAIN DECK
# ════════════════════════════════════════════════════════════════════════


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08)
add_accent_bar(slide, 0, 7.42, 13.333, 0.08)

add_textbox(slide, 1.5, 1.5, 10, 1.5,
            "PIK Structural Credit Pricing",
            font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, 4, 3.0, 5.333, 0.04, ACCENT_GOLD)

add_textbox(slide, 1.5, 3.3, 10, 1.0,
            "How Much Extra Spread Should Investors Demand for PIK Risk?",
            font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 4.3, 10, 0.5,
            "Structural Credit + Monte Carlo Analysis",
            font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_multi_text(slide, 1.5, 5.5, 10, 1.2, [
    ("Cash-Pay  |  Full PIK  |  PIK Toggle  |  9 Issuers (BB+ to CCC-)", False, ACCENT_GOLD),
    ("5-Year 8.5% Semi-Annual Bond  |  25,000 MC Paths  |  Calibrated to Market", False, MED_GRAY),
], font_size=14)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda / Overview
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08)

add_textbox(slide, 0.5, 0.3, 5, 0.6, "Agenda", font_size=32, color=ACCENT_GOLD, bold=True)

sections = [
    ("1.", "What Are PIK Bonds?", "How cash-pay, full PIK, and PIK toggle structures work"),
    ("2.", "The Credit Spectrum", "9 issuers from BB+ to CCC- \u2014 how credit quality shapes PIK risk"),
    ("3.", "Modelling Framework", "Merton structural model [1,2] and reduced-form hazard rate model [3,4]"),
    ("4.", "The PIK Feedback Loop", "Why deferred coupons create a self-reinforcing leverage spiral [5,6]"),
    ("5.", "The PIK Premium", "How much extra spread should investors demand? (+58bp to +628bp)"),
    ("6.", "What Drives the Premium?", "Sensitivity to coverage, volatility, and model assumptions"),
    ("7.", "Borrower Optionality", "Toggle exercise and its cost to lenders [8]"),
    ("8.", "Conclusions & Implications", "Key findings for investors and risk managers"),
]

for i, (num, title, desc) in enumerate(sections):
    y = 1.2 + i * 0.74
    add_textbox(slide, 1.0, y, 0.5, 0.5, num, font_size=20, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, 1.5, y, 5, 0.35, title, font_size=20, color=WHITE, bold=True)
    add_textbox(slide, 1.5, y + 0.32, 10, 0.35, desc, font_size=13, color=LIGHT_GRAY)

add_footnote(slide, "[1] Merton (1974)  [2] Black & Cox (1976)  [3] Duffie & Singleton (1999)  "
             "[4] Jarrow & Turnbull (1995)  [5] Leland (1994)  [6] Collin-Dufresne & Goldstein (2001)  "
             "[7] Glasserman (2003)  [8] Longstaff & Schwartz (2001)  \u2014  see References slide",
             color=MED_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3: Bond Structures
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 8, 0.6, "What Are PIK Bonds?",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "PIK bonds defer cash coupons into notional growth. If the issuer defaults while notional "
    "has grown, investors lose more. The question: how much extra spread compensates for this risk?")

add_textbox(slide, 0.5, 1.6, 12, 0.4,
            "Stylized 5-year 8.5% semi-annual fixed-coupon bond (4.5% risk-free + 400bp credit spread)",
            font_size=15, color=MED_GRAY)

add_table(slide, 0.5, 2.1, 6.5, 1.5,
          ["Structure", "Coupon Mechanics", "PIK Schedule"],
          [
              ["Cash-Pay", "All coupons paid in cash at each date", "None"],
              ["Full PIK", "All coupons accreted to notional", "Full PIK"],
              ["PIK Toggle", "Borrower decides cash vs PIK", "State-dependent"],
          ],
          col_widths=[1.5, 3.0, 2.0])

# Key parameters box
add_textbox(slide, 7.5, 2.1, 5.3, 0.4, "Key Parameters", font_size=18,
            color=DARK_BG, bold=True)
params = [
    "Risk-free rate: 4.5%",
    "Coupon spread: 400bp (total coupon = 8.5%)",
    "Maturity: 5 years",
    "Notional: 100",
    "Monte Carlo paths: 25,000",
    "Antithetic variates + Brownian bridge [7]",
    "Seed: 42 (reproducible)",
]
add_multi_text(slide, 7.5, 2.5, 5.3, 3.0, [
    (f"\u2022 {p}", False, DARK_TEXT) for p in params
], font_size=14, color=DARK_TEXT)

# PIK mechanics note
add_textbox(slide, 0.5, 4.3, 12, 0.4, "PIK Mechanics in Detail",
            font_size=18, color=DARK_BG, bold=True)
add_multi_text(slide, 0.5, 4.7, 12, 2.5, [
    ("\u2022 Cash-Pay: investor receives coupon cash each period. Notional stays at par.", False, DARK_TEXT),
    ("\u2022 Full PIK: no cash coupons \u2014 each coupon accretes to notional. At maturity, N(T) = 100 \u00d7 (1 + 4.25%)^10 = 151.6", False, DARK_TEXT),
    ("\u2022 PIK Toggle: borrower chooses each period whether to pay cash or PIK. Creates adverse selection risk [8].", False, DARK_TEXT),
], font_size=14)

add_footnote(slide, "[7] Glasserman (2003), Monte Carlo Methods in Financial Engineering  "
             "[8] Longstaff & Schwartz (2001), Valuing American Options by Simulation")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4: Issuer Profiles
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 8, 0.6, "The Credit Spectrum",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "We test PIK economics across 9 issuers from BB+ to CCC-. The same PIK structure behaves "
    "very differently across credit quality \u2014 a BB+ PIK barely matters, but a CCC- PIK premium can exceed 600bp.")

add_table(slide, 0.5, 1.7, 12.3, 3.5,
          ["Issuer", "Assets (V)", "Vol (\u03c3)", "Mkt Spread", "h\u2080 (cal)", "h\u2248s/(1-R)", "R\u2080"],
          [
              ["BB+",  "200", "20%", "150bp", "2.77%", "3.00%", "50%"],
              ["BB",   "190", "20%", "175bp", "2.99%", "3.18%", "45%"],
              ["BB-",  "180", "25%", "200bp", "3.18%", "3.33%", "40%"],
              ["B+",   "170", "25%", "250bp", "3.72%", "3.85%", "35%"],
              ["B",    "150", "30%", "300bp", "4.34%", "4.44%", "32%"],
              ["B-",   "130", "35%", "400bp", "5.65%", "5.71%", "30%"],
              ["CCC+", "120", "35%", "550bp", "7.35%", "7.33%", "25%"],
              ["CCC",  "115", "40%", "700bp", "8.85%", "8.75%", "20%"],
              ["CCC-", "110", "40%", "975bp", "11.71%", "11.47%", "15%"],
          ],
          col_widths=[1.2, 1.3, 1.0, 1.5, 1.5, 1.5, 1.0])

# PD table
add_textbox(slide, 0.5, 5.4, 8, 0.4, "Market-Implied Default Probabilities: PD(t) = 1 - exp(-\u03bbt)",
            font_size=16, color=DARK_BG, bold=True)

add_table(slide, 0.5, 5.8, 12.3, 1.2,
          ["Issuer", "\u03bb (bp)", "PD(1Y)", "PD(3Y)", "PD(5Y)", "PD(7Y)", "PD(10Y)"],
          [
              ["BB+",  "277", "2.73%", "7.98%", "12.94%", "17.63%", "24.20%"],
              ["B",    "434", "4.25%", "12.21%", "19.51%", "26.20%", "35.21%"],
              ["CCC-", "1171", "11.05%", "29.63%", "44.32%", "55.95%", "69.00%"],
          ],
          col_widths=[1.2, 1.3, 1.5, 1.5, 1.5, 1.5, 1.5],
          font_size=10, header_font_size=10)

add_footnote(slide, "Hazard rates calibrated via reduced-form model [3,4]. "
             "Credit triangle approximation: \u03bb \u2248 s / (1 - R), see Duffie & Singleton (1999).")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 5: Merton (1974) Structural Credit Model
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "Merton (1974) Structural Credit Model",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "The Merton model treats the company\u2019s equity as a call option on its assets. Default "
    "occurs when assets fall below the debt barrier. This structural approach naturally captures "
    "how PIK-inflated debt changes default dynamics \u2014 unlike hazard-rate models.", height=0.6)

# Left column: model description
add_multi_text(slide, 0.5, 1.7, 6, 4.5, [
    ("Asset Dynamics (GBM under risk-neutral measure) [1]:", True, DARK_BG),
    ("dV = (r - q) \u00b7 V \u00b7 dt + \u03c3v \u00b7 V \u00b7 dW", False, ACCENT_BLUE),
    ("", False, DARK_BG),
    ("Distance-to-Default:", True, DARK_BG),
    ("DD = [ln(V/B) + (r - q - \u00bd\u03c3\u00b2) \u00b7 T] / (\u03c3 \u00b7 \u221aT)", False, ACCENT_BLUE),
    ("", False, DARK_BG),
    ("Default Probability:", True, DARK_BG),
    ("PD(T) = N(-DD)  where N(\u00b7) is the standard normal CDF", False, ACCENT_BLUE),
    ("", False, DARK_BG),
    ("Implied Credit Spread:", True, DARK_BG),
    ("s = -ln(1 - PD \u00b7 LGD) / T,  LGD = 1 - R", False, ACCENT_BLUE),
    ("", False, DARK_BG),
    ("Terminal vs First-Passage [2]:", True, DARK_BG),
    ("\u2022 Terminal: default only at maturity (Merton 1974)", False, DARK_TEXT),
    ("\u2022 First-passage: default at any time V < B (Black & Cox 1976)", False, DARK_TEXT),
    ("\u2022 First-passage defaults are 1.5-2x higher than terminal rates", False, RED),
], font_size=13)

# Right column: analytical results
add_textbox(slide, 7.0, 1.7, 6, 0.4, "Analytical Credit Metrics",
            font_size=18, color=DARK_BG, bold=True)

add_table(slide, 7.0, 2.1, 5.8, 3.5,
          ["Issuer", "Barrier", "DD", "PD(5Y)", "Impl Spread"],
          [
              ["BB+",  "136.8", "1.13", "12.94%", "134 bp"],
              ["BB",   "132.5", "1.08", "13.90%", "159 bp"],
              ["BB-",  "107.3", "1.05", "14.70%", "185 bp"],
              ["B+",   "106.8", "0.95", "16.99%", "234 bp"],
              ["B",    "84.3",  "0.86", "19.51%", "282 bp"],
              ["B-",   "70.0",  "0.69", "24.60%", "378 bp"],
              ["CCC+", "74.7",  "0.50", "30.76%", "525 bp"],
              ["CCC",  "69.6",  "0.37", "35.75%", "674 bp"],
              ["CCC-", "81.3",  "0.14", "44.32%", "946 bp"],
          ],
          col_widths=[0.9, 1.0, 0.8, 1.0, 1.2])

add_multi_text(slide, 7.0, 5.8, 5.8, 1.0, [
    ("Barriers calibrated to match market-implied PD", False, MED_GRAY),
    ("via from_target_pd (terminal basis)", False, MED_GRAY),
], font_size=11)

add_footnote(slide, "[1] Merton, R.C. (1974), On the Pricing of Corporate Debt, Journal of Finance 29(2)  "
             "[2] Black, F. & Cox, J.C. (1976), Valuing Corporate Securities, Journal of Finance 31(2)")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 6: Reduced-Form Hazard Rate Model
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "Reduced-Form Hazard Rate Model",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "The hazard rate model prices bonds using a constant default intensity \u03bb calibrated to market "
    "spreads. It captures the timing and notional effects of PIK but cannot capture the leverage "
    "feedback loop \u2014 it serves as our baseline for measuring what the structural model adds.",
    height=0.65)

# Left: model formulation
add_accent_bar(slide, 0.5, 1.7, 5.8, 0.04, ACCENT_BLUE)
add_textbox(slide, 0.5, 1.8, 5.8, 0.4, "Model Formulation",
            font_size=20, color=DARK_BG, bold=True)

add_multi_text(slide, 0.5, 2.3, 5.8, 3.0, [
    ("Survival function:", True, DARK_BG),
    ("S(t) = exp(-\u03bb \u00b7 t)", False, ACCENT_BLUE),
    ("", False, DARK_BG),
    ("Bond pricing with default [3,4]:", True, DARK_BG),
    ("PV = \u03a3 c\u1d62 \u00b7 D(t\u1d62) \u00b7 S(t\u1d62)  +  N \u00b7 D(T) \u00b7 S(T)", False, ACCENT_BLUE),
    ("    + R \u00b7 N \u00b7 \u03a3 D(t\u1d62) \u00b7 [S(t\u2099\u208b\u2081) - S(t\u1d62)]", False, ACCENT_BLUE),
    ("", False, DARK_BG),
    ("Credit triangle approximation:", True, DARK_BG),
    ("\u03bb\u2080 \u2248 s / (1 - R)", False, ACCENT_BLUE),
    ("", False, DARK_BG),
    ("\u2022 \u03bb is FLAT \u2014 does not respond to leverage changes", False, RED),
    ("\u2022 Captures PIK timing: PIK shifts notional to maturity", False, DARK_TEXT),
    ("\u2022 Misses feedback: PIK \u2192 higher leverage \u2192 higher \u03bb", False, DARK_TEXT),
], font_size=13)

# Right: comparison setup
add_accent_bar(slide, 7.0, 1.7, 5.8, 0.04, ACCENT_GOLD)
add_textbox(slide, 7.0, 1.8, 5.8, 0.4, "Two Models, Two Purposes",
            font_size=20, color=DARK_BG, bold=True)

add_multi_text(slide, 7.0, 2.3, 5.8, 3.5, [
    ("Hazard Rate (Reduced-Form):", True, ACCENT_BLUE),
    ("\u2022 Constant \u03bb calibrated from market", False, DARK_TEXT),
    ("\u2022 Tractable, closed-form pricing", False, DARK_TEXT),
    ("\u2022 PIK premium from timing effect only", False, DARK_TEXT),
    ("", False, DARK_BG),
    ("Merton MC (Structural + Feedback):", True, ACCENT_GOLD),
    ("\u2022 Asset dynamics drive default", False, DARK_TEXT),
    ("\u2022 Endogenous \u03bb(L) rises with leverage [5]", False, DARK_TEXT),
    ("\u2022 Dynamic recovery R(N) falls with PIK accrual", False, DARK_TEXT),
    ("\u2022 Captures the feedback loop", False, DARK_TEXT),
    ("", False, DARK_BG),
    ("Gap between models = value of the feedback loop", True, RED),
], font_size=13)

# Bottom: what the gap tells us
add_multi_text(slide, 0.5, 5.5, 12, 1.5, [
    ("The difference between the two models isolates the structural feedback premium:", True, DARK_BG),
    ("At CCC-: HR premium = +220bp, MC premium = +628bp \u2192 feedback adds +408bp (\u2248 2/3 of total PIK premium)", False, RED),
], font_size=14)

add_footnote(slide, "[3] Duffie, D. & Singleton, K.J. (1999), Modeling Term Structures of Defaultable Bonds, RFS 12(4)  "
             "[4] Jarrow, R.A. & Turnbull, S.M. (1995), Pricing Derivatives on Financial Securities Subject to Credit Risk, JoF 50(1)")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 7: The PIK Feedback Loop
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "The PIK Feedback Loop",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "PIK creates a vicious cycle: unpaid coupons inflate the debt, pushing leverage higher, "
    "which raises default probability AND lowers recovery. This self-reinforcing spiral is the "
    "core mechanism driving the PIK premium.", height=0.65)

# Left: Endogenous Hazard
add_accent_bar(slide, 0.5, 1.7, 5.8, 0.04, ACCENT_BLUE)
add_textbox(slide, 0.5, 1.8, 5.8, 0.4, "Endogenous Hazard (Power-Law) [5,6]",
            font_size=20, color=DARK_BG, bold=True)

add_multi_text(slide, 0.5, 2.3, 5.8, 3.0, [
    ("\u03bb(L) = \u03bb\u2080 \u00b7 (L / L\u2080)\u1d5d", True, ACCENT_BLUE),
    ("", False, DARK_BG),
    ("\u2022 \u03bb\u2080 = base hazard, calibrated from market spreads", False, DARK_TEXT),
    ("\u2022 L\u2080 = base leverage at origination (B/V)", False, DARK_TEXT),
    ("\u2022 \u03b2 = exponent controlling feedback speed", False, DARK_TEXT),
    ("", False, DARK_BG),
    ("As PIK inflates notional and leverage rises,", False, DARK_TEXT),
    ("hazard accelerates non-linearly.", False, DARK_TEXT),
], font_size=14)

# Right: Dynamic Recovery
add_accent_bar(slide, 7.0, 1.7, 5.8, 0.04, ACCENT_GOLD)
add_textbox(slide, 7.0, 1.8, 5.8, 0.4, "Dynamic Recovery (Floored Inverse)",
            font_size=20, color=DARK_BG, bold=True)

add_multi_text(slide, 7.0, 2.3, 5.8, 3.0, [
    ("R(N) = max(floor, R\u2080 \u00b7 N\u2080 / N)", True, ACCENT_GOLD),
    ("", False, DARK_BG),
    ("\u2022 R\u2080 = base recovery at par notional N\u2080", False, DARK_TEXT),
    ("\u2022 floor = minimum recovery rate (10%)", False, DARK_TEXT),
    ("\u2022 \u03b1 variant: R(N) = R\u2080 \u00b7 (N\u2080/N)^\u03b1", False, DARK_TEXT),
    ("", False, DARK_BG),
    ("PIK-inflated notional means recovery < base rate.", False, DARK_TEXT),
    ("Primary driver of full-PIK premium.", True, DARK_TEXT),
], font_size=14)

# Feedback demo table
add_textbox(slide, 0.5, 5.0, 8, 0.4, "Feedback Demonstration: B+ Issuer (barrier = 106.8, base leverage = 0.63)",
            font_size=16, color=DARK_BG, bold=True)

add_table(slide, 0.5, 5.4, 8, 1.6,
          ["LTV", "Hazard", "Notional", "Recovery"],
          [
              ["0.50", "2.36%", "79.6", "35.00%"],
              ["0.70", "4.62%", "111.4", "31.41%"],
              ["0.90", "7.64%", "143.3", "24.43%"],
              ["1.00", "9.44%", "159.2", "21.99%"],
              ["1.10", "11.42%", "175.1", "19.99%"],
          ],
          col_widths=[1.5, 1.5, 1.5, 1.5],
          font_size=10, header_font_size=10)

add_footnote(slide, "[5] Leland, H.E. (1994), Corporate Debt Value, Bond Covenants, and Optimal Capital Structure, JoF 49(4)  "
             "[6] Collin-Dufresne, P. & Goldstein, R.S. (2001), Do Credit Spreads Reflect Stationary Leverage Ratios?, JoF 56(5)")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8: MC-Calibrated Breakeven Spreads — THE key slide
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.6, "The PIK Premium: How Much Extra Spread?",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "After calibrating to market, investors should demand +58bp (BB+) to +628bp (CCC-) "
    "additional spread for full PIK structures. The premium scales exponentially with credit risk.",
    height=0.55)

add_table(slide, 0.5, 1.6, 12.3, 3.8,
          ["Issuer", "Mkt Spread", "Cash* Z", "PIK* Z", "Toggle* Z", "PIK-Cash*", "Toggle-Cash*"],
          [
              ["BB+",  "150bp", "149bp", "206bp", "150bp", "+58bp", "+1bp"],
              ["BB",   "175bp", "174bp", "241bp", "176bp", "+68bp", "+2bp"],
              ["BB-",  "200bp", "199bp", "275bp", "228bp", "+76bp", "+29bp"],
              ["B+",   "250bp", "253bp", "358bp", "295bp", "+105bp", "+43bp"],
              ["B",    "300bp", "302bp", "438bp", "424bp", "+136bp", "+122bp"],
              ["B-",   "400bp", "407bp", "616bp", "621bp", "+209bp", "+214bp"],
              ["CCC+", "550bp", "548bp", "848bp", "845bp", "+300bp", "+297bp"],
              ["CCC",  "700bp", "698bp", "1096bp", "1095bp", "+398bp", "+397bp"],
              ["CCC-", "975bp", "981bp", "1609bp", "1586bp", "+628bp", "+606bp"],
          ],
          col_widths=[1.2, 1.5, 1.5, 1.5, 1.5, 1.8, 1.8])

add_multi_text(slide, 0.5, 5.7, 12, 1.3, [
    ("* = MC-calibrated: cash Z-spread matches market by construction", False, MED_GRAY),
    ("Calibrated PIK premiums are wider: +628bp at CCC- (vs +409bp uncalibrated)", True, RED),
    ("Uncalibrated barriers inflate cash-pay spreads disproportionately \u2014 see Appendix H", False, DARK_TEXT),
], font_size=14)

add_footnote(slide, "Z-spreads computed on cash-equivalent basis for cross-structure comparability. "
             "MC calibration: bisect on barrier to match cash Z-spread to market. See [7] for MC methodology.")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 9: Hazard-Rate vs Merton MC
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.6, "Why Simple Models Underestimate PIK Risk",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "Traditional hazard-rate models [3,4] miss 2/3 of the PIK premium for stressed issuers. "
    "A CCC- PIK bond is mispriced by ~400bp if you ignore the leverage spiral [5,6]. "
    "Structural models with feedback are essential for accurate pricing.")

add_table(slide, 0.5, 1.6, 12.3, 3.8,
          ["Issuer", "\u03bb (bp)", "HR Cash Z", "HR PIK Z", "HR PIK-Cash", "MC Cash* Z", "MC PIK* Z", "MC PIK-Cash*"],
          [
              ["BB+",  "277", "149bp", "197bp", "+48bp", "149bp", "206bp", "+58bp"],
              ["BB",   "299", "174bp", "227bp", "+53bp", "174bp", "241bp", "+68bp"],
              ["BB-",  "318", "199bp", "257bp", "+58bp", "199bp", "275bp", "+76bp"],
              ["B+",   "372", "249bp", "316bp", "+67bp", "253bp", "358bp", "+105bp"],
              ["B",    "434", "299bp", "375bp", "+76bp", "302bp", "438bp", "+136bp"],
              ["B-",   "565", "398bp", "493bp", "+95bp", "407bp", "616bp", "+209bp"],
              ["CCC+", "735", "548bp", "673bp", "+125bp", "548bp", "848bp", "+300bp"],
              ["CCC",  "885", "698bp", "855bp", "+157bp", "698bp", "1096bp", "+398bp"],
              ["CCC-", "1171", "973bp", "1193bp", "+220bp", "981bp", "1609bp", "+628bp"],
          ],
          col_widths=[1.0, 0.9, 1.3, 1.3, 1.5, 1.3, 1.3, 1.8])

add_multi_text(slide, 0.5, 5.7, 12, 1.3, [
    ("At CCC-: Hazard-rate PIK premium = +220bp vs Structural MC = +628bp", True, RED),
    ("The structural feedback (endogenous hazard + dynamic recovery) accounts for nearly 2/3 of the PIK premium", True, DARK_BG),
], font_size=14)

add_footnote(slide, "HR = reduced-form hazard rate model [3,4] with flat \u03bb. "
             "MC = Merton structural model [1,2] with endogenous hazard [5] and dynamic recovery, "
             "calibrated to market.")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 10: Sensitivity — Asset Coverage
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "What Drives the Premium: Asset Coverage",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "Asset coverage is the single most important driver. The premium drops from +372bp at 1.1x "
    "coverage to just +26bp at 2.15x. Better collateral dramatically reduces PIK risk.")

add_textbox(slide, 0.5, 1.6, 12, 0.4,
            "Sweep V from 110 to 220 (coverage 1.10x to 2.15x), fixed barrier at par, \u03c3=30%, h\u2080=6%",
            font_size=15, color=MED_GRAY)

add_table(slide, 0.5, 2.0, 12.3, 3.3,
          ["Coverage", "Cash Price", "PIK Price", "Discount", "Cash Z", "PIK Z", "PIK Premium"],
          [
              ["1.10x", "92.98", "79.81", "+13.17", "557bp", "929bp", "+372bp"],
              ["1.25x", "97.05", "86.67", "+10.38", "454bp", "728bp", "+273bp"],
              ["1.40x", "100.42", "92.34", "+8.08", "372bp", "574bp", "+202bp"],
              ["1.55x", "103.33", "97.23", "+6.10", "304bp", "450bp", "+146bp"],
              ["1.70x", "105.66", "101.16", "+4.50", "251bp", "355bp", "+104bp"],
              ["1.85x", "107.64", "104.49", "+3.15", "207bp", "277bp", "+71bp"],
              ["2.00x", "109.19", "107.10", "+2.09", "173bp", "219bp", "+46bp"],
              ["2.15x", "110.48", "109.28", "+1.20", "145bp", "171bp", "+26bp"],
          ],
          col_widths=[1.3, 1.5, 1.5, 1.5, 1.5, 1.5, 1.8])

add_multi_text(slide, 0.5, 5.6, 12, 1.3, [
    ("Key Insight: PIK premium is strongly non-linear in coverage", True, DARK_BG),
    ("\u2022 At 1.10x coverage: +372bp PIK premium (asset barely covers debt)", False, RED),
    ("\u2022 At 2.15x coverage: +26bp PIK premium (defaults are rare, PIK accrual harmless)", False, GREEN),
    ("\u2022 The relationship is convex \u2014 premiums accelerate as coverage drops below 1.5x", False, DARK_TEXT),
], font_size=14)

add_footnote(slide, "Structural framework per Merton (1974) [1] with first-passage default [2] and "
             "endogenous hazard [5]. Coverage = V/B at origination.")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 11: Sensitivity — Asset Volatility
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "What Drives the Premium: Asset Volatility",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "Higher volatility amplifies PIK risk exponentially \u2014 at 45% vol the PIK premium is 9x larger "
    "than at 20% vol. Volatile sectors (e.g. tech, energy) face outsized PIK costs.")

add_textbox(slide, 0.5, 1.6, 12, 0.4,
            "Fixed 1.40x coverage (V=140, barrier=par), sweep \u03c3 from 20% to 45%",
            font_size=15, color=MED_GRAY)

add_table(slide, 0.5, 2.0, 12.3, 2.5,
          ["Vol (\u03c3)", "Cash Price", "PIK Price", "Discount", "Cash Z", "PIK Z", "PIK Premium"],
          [
              ["20%", "109.17", "107.07", "+2.10", "173bp", "219bp", "+46bp"],
              ["25%", "104.58", "99.34", "+5.24", "275bp", "398bp", "+123bp"],
              ["30%", "100.42", "92.34", "+8.08", "372bp", "574bp", "+202bp"],
              ["35%", "96.90", "86.42", "+10.48", "458bp", "735bp", "+277bp"],
              ["40%", "93.81", "81.20", "+12.61", "536bp", "887bp", "+351bp"],
              ["45%", "91.06", "76.58", "+14.48", "608bp", "1031bp", "+424bp"],
          ],
          col_widths=[1.3, 1.5, 1.5, 1.5, 1.5, 1.5, 1.8])

add_multi_text(slide, 0.5, 4.8, 12, 2.0, [
    ("Volatility amplifies PIK risk through three channels:", True, DARK_BG),
    ("", False, DARK_BG),
    ("1. More frequent barrier crossings \u2192 higher default rate [2]", False, DARK_TEXT),
    ("2. More extreme PIK paths \u2192 wider leverage swings \u2192 deeper recovery erosion", False, DARK_TEXT),
    ("3. Non-linear amplification \u2192 both default frequency AND loss severity worsen simultaneously", False, DARK_TEXT),
    ("", False, DARK_BG),
    ("PIK premium at 20% vol: +46bp.  At 45% vol: +424bp  (\u00d79.2x increase)", True, RED),
], font_size=14)

add_footnote(slide, "GBM asset dynamics [1] with first-passage barrier crossing [2] and "
             "Brownian bridge correction [7].")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 12: Sensitivity — Feedback Loop Speed
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 12, 0.6,
            "What Drives the Premium: Feedback Speed (\u03b2 and \u03b1)",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "The speed of the feedback loop is a key modelling assumption. How fast hazard rises with "
    "leverage (\u03b2) and how fast recovery erodes with notional growth (\u03b1) can swing the premium by 100bp+. "
    "Getting these parameters right is critical.")

# Beta sensitivity
add_textbox(slide, 0.5, 1.6, 6, 0.4,
            "Endogenous Hazard Speed (\u03b2) [5,6] \u2192 Drives Toggle Premium",
            font_size=16, color=ACCENT_BLUE, bold=True)

add_table(slide, 0.5, 2.0, 6.0, 2.2,
          ["Issuer", "Metric", "\u03b2=0", "\u03b2=1", "\u03b2=2", "\u03b2=4"],
          [
              ["BB-", "Toggle-Cash", "+0bp", "+2bp", "+29bp", "+66bp"],
              ["B",   "Toggle-Cash", "+0bp", "+53bp", "+122bp", "+145bp"],
              ["CCC", "Toggle-Cash", "+398bp", "+393bp", "+397bp", "+399bp"],
          ],
          col_widths=[0.7, 1.3, 0.8, 0.8, 0.8, 0.8],
          font_size=10, header_font_size=10)

# Alpha sensitivity
add_textbox(slide, 6.8, 1.6, 6, 0.4,
            "Recovery Decay Speed (\u03b1) \u2192 Drives PIK Premium",
            font_size=16, color=ACCENT_GOLD, bold=True)

add_table(slide, 6.8, 2.0, 6.0, 2.2,
          ["Issuer", "Metric", "\u03b1=0", "\u03b1=0.5", "\u03b1=1.0", "\u03b1=2.0"],
          [
              ["BB-", "PIK-Cash", "+16bp", "+48bp", "+76bp", "+118bp"],
              ["B",   "PIK-Cash", "+66bp", "+104bp", "+136bp", "+186bp"],
              ["CCC", "PIK-Cash", "+294bp", "+351bp", "+398bp", "+473bp"],
          ],
          col_widths=[0.7, 1.3, 0.8, 0.8, 0.8, 0.8],
          font_size=10, header_font_size=10)

add_multi_text(slide, 0.5, 4.4, 12, 2.5, [
    ("Two distinct feedback channels:", True, DARK_BG),
    ("", False, DARK_BG),
    ("\u03b2 (Hazard Speed) [5,6]:", True, ACCENT_BLUE),
    ("\u2022 Controls how quickly hazard ramps as leverage rises", False, DARK_TEXT),
    ("\u2022 Drives the toggle exercise decision (PIK when hazard is high)", False, DARK_TEXT),
    ("\u2022 At \u03b2=0, toggle sees flat hazard and rarely PIKs; at \u03b2=4, toggle aggressively PIKs", False, DARK_TEXT),
    ("", False, DARK_BG),
    ("\u03b1 (Recovery Decay):", True, ACCENT_GOLD),
    ("\u2022 Controls how fast recovery erodes as PIK inflates notional", False, DARK_TEXT),
    ("\u2022 Directly affects loss severity \u2014 primary driver of full-PIK premium", False, DARK_TEXT),
    ("\u2022 At \u03b1=0, recovery is constant; at \u03b1=2.0, recovery collapses quickly", False, DARK_TEXT),
    ("", False, DARK_BG),
    ("Both effects are strongly non-linear in credit quality", True, RED),
], font_size=13)

add_footnote(slide, "[5] Leland (1994): endogenous default barrier and leverage-dependent hazard.  "
             "[6] Collin-Dufresne & Goldstein (2001): stationary leverage ratios and credit spread dynamics.")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 13: Toggle Strategy Comparison
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "Borrower Optionality: Toggle Strategies",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "Borrowers exercise the toggle at the worst time for lenders \u2014 they PIK when the company is "
    "most stressed, seeding the leverage spiral on already-distressed paths. Optimal exercise [8] "
    "produces spreads nearly as wide as full PIK (adverse selection).")

add_textbox(slide, 0.5, 1.6, 10, 0.3, "B+ Issuer Results",
            font_size=16, color=MED_GRAY)

add_table(slide, 0.5, 1.9, 12.3, 2.5,
          ["Strategy", "Price", "Z-Spread", "E[Loss]", "Def Rate", "PIK%", "Terminal Notional"],
          [
              ["Cash-Pay", "108.21", "194bp", "7.89%", "16.77%", "0.0%", "100.0"],
              ["Full PIK", "105.46", "255bp", "10.24%", "16.77%", "100.0%", "151.6"],
              ["Threshold (h > 10%)", "106.51", "232bp", "9.34%", "16.77%", "6.6%", "100.9"],
              ["Stochastic (sigmoid)", "106.59", "230bp", "9.27%", "16.77%", "15.5%", "105.2"],
              ["Optimal (nested MC) [8]", "105.99", "244bp", "9.78%", "16.77%", "14.6%", "104.1"],
          ],
          col_widths=[2.2, 1.2, 1.4, 1.4, 1.4, 1.2, 2.2])

# Strategy descriptions
add_multi_text(slide, 0.5, 4.7, 12, 2.3, [
    ("Toggle Exercise Strategies:", True, DARK_BG),
    ("", False, DARK_BG),
    ("Threshold:  Deterministic rule \u2014 PIK when hazard rate > 10%", False, DARK_TEXT),
    ("Stochastic: Smooth sigmoid \u2014 P(PIK) = \u03c3(a + b \u00b7 x), avoids hard threshold discontinuity", False, DARK_TEXT),
    ("Optimal:    Forward-looking nested MC [8] \u2014 maximizes equity value (most expensive to compute)", False, DARK_TEXT),
    ("", False, DARK_BG),
    ("Insight: Optimal exercise produces spreads close to full PIK (244bp vs 255bp)", True, RED),
    ("The borrower rationally chooses PIK when it's worst for the lender (adverse selection)", False, DARK_TEXT),
], font_size=14)

add_footnote(slide, "[8] Longstaff, F.A. & Schwartz, E.S. (2001), Valuing American Options by Simulation: "
             "A Simple Least-Squares Approach, RFS 14(1). Adapted for toggle exercise as nested MC.")


# ════════════════════════════════════════════════════════════════════════
# SLIDE 14: Key Conclusions (1 of 2)
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_GOLD)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "Conclusions (1/2)",
            font_size=28, color=WHITE, bold=True)

add_dark_callout(slide,
    "Five key takeaways: PIK risk is non-linear, feedback-driven, and underestimated by standard models.")

conclusions_1 = [
    ("1", "PIK premium is non-linear in LTV",
     "BB+ (~50% LTV): PIK barely matters. CCC (~87% LTV): PIK premium can exceed 500bp."),
    ("2", "The feedback loop drives PIK risk [5,6]",
     "Endogenous hazard + dynamic recovery create a self-reinforcing spiral that dramatically worsens PIK economics."),
    ("3", "Simple hazard rates miss the story [3,4]",
     "Flat \u03bb captures timing/notional effects but underestimates premium by 2-3x for stressed issuers. CCC-: +220bp (HR) vs +628bp (MC)."),
    ("4", "Toggle is not free [8]",
     "Borrower PIKs precisely when credit deteriorates \u2014 negative convexity for the investor. Toggle premium sits between cash and full PIK."),
    ("5", "Volatility amplifies PIK risk",
     "Higher \u03c3 increases both default frequency AND loss severity [1,2]. At 40% vol, PIK premium >350bp; at 20% vol, only ~46bp."),
]

for i, (num, title, desc) in enumerate(conclusions_1):
    y = 1.6 + i * 1.1
    add_textbox(slide, 0.8, y, 0.5, 0.4, num, font_size=24, color=ACCENT_GOLD, bold=True)
    add_textbox(slide, 1.3, y, 11, 0.35, title, font_size=18, color=WHITE, bold=True)
    add_textbox(slide, 1.3, y + 0.4, 11, 0.6, desc, font_size=14, color=LIGHT_GRAY)

add_footnote(slide, "Reference numbers: see References slide. [1] Merton (1974)  [2] Black & Cox (1976)  "
             "[3] Duffie & Singleton (1999)  [5] Leland (1994)", color=MED_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 15: Key Conclusions (2 of 2)
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_GOLD)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "Conclusions (2/2)",
            font_size=28, color=WHITE, bold=True)

add_dark_callout(slide,
    "Model robustness, duration implications, and measurement considerations for practitioners.")

conclusions_2 = [
    ("6", "MC-to-market calibration eliminates fudge factors",
     "Calibrating the barrier directly to market Z-spread ensures cash base case matches by construction. PIK differentials are computed consistently."),
    ("7", "Toggle can be worse than full PIK (adverse selection)",
     "Toggle concentrates PIK on the worst paths, seeding the leverage spiral on already-stressed paths. Can produce larger investor losses than mandatory full PIK."),
    ("8", "PIK bonds carry ~0.75yr longer modified duration",
     "Full-PIK ModDur ~4.8yr vs ~4.1yr for cash-pay (~18% extension). Greater interest rate sensitivity and higher DV01 per dollar of notional."),
    ("9", "Metrics use the right cashflow basis",
     "Spread/yield metrics: cash-equivalent basis for cross-structure comparability. Risk metrics (duration, DV01): actual PIK cashflows for true rate sensitivity."),
    ("10", "Calibration parameter choice does not affect PIK premiums",
     "Full-PIK premium is identical under barrier or vol calibration. Toggle shows small differences. Model is robust to calibration method."),
]

for i, (num, title, desc) in enumerate(conclusions_2):
    y = 1.6 + i * 1.1
    add_textbox(slide, 0.8, y, 0.5, 0.4, num, font_size=24, color=ACCENT_GOLD, bold=True)
    add_textbox(slide, 1.3, y, 11, 0.35, title, font_size=18, color=WHITE, bold=True)
    add_textbox(slide, 1.3, y + 0.4, 11, 0.6, desc, font_size=14, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 16: Summary / Takeaway
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_GOLD)
add_accent_bar(slide, 0, 7.42, 13.333, 0.08, ACCENT_GOLD)

add_textbox(slide, 1.5, 1.0, 10, 1.0,
            "Key Takeaway",
            font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, 4, 2.0, 5.333, 0.04, ACCENT_GOLD)

add_textbox(slide, 1.5, 2.5, 10, 2.5,
            "PIK risk is fundamentally a leverage spiral risk.\n"
            "Structural models with endogenous hazard and dynamic recovery\n"
            "reveal premiums 2-3x larger than simple hazard-rate approaches,\n"
            "especially for stressed issuers where it matters most.",
            font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, 4, 4.8, 5.333, 0.04, ACCENT_GOLD)

# Summary stats in a row
stats = [
    ("BB+ PIK Premium", "+58bp"),
    ("CCC- PIK Premium", "+628bp"),
    ("Max Vol Impact", "9.2x"),
    ("Duration Extension", "+0.75yr"),
]
for i, (label, val) in enumerate(stats):
    x = 1.0 + i * 3.1
    add_textbox(slide, x, 5.2, 2.8, 0.8, val,
                font_size=36, color=ACCENT_GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, 5.9, 2.8, 0.5, label,
                font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 17: References
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, ACCENT_BLUE)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "References",
            font_size=28, color=DARK_BG, bold=True)

refs = [
    ("[1]", "Merton, R.C. (1974)",
     "\"On the Pricing of Corporate Debt: The Risk Structure of Interest Rates.\" "
     "Journal of Finance, 29(2), 449-470."),

    ("[2]", "Black, F. & Cox, J.C. (1976)",
     "\"Valuing Corporate Securities: Some Effects of Bond Indenture Provisions.\" "
     "Journal of Finance, 31(2), 351-367."),

    ("[3]", "Duffie, D. & Singleton, K.J. (1999)",
     "\"Modeling Term Structures of Defaultable Bonds.\" "
     "Review of Financial Studies, 12(4), 687-720."),

    ("[4]", "Jarrow, R.A. & Turnbull, S.M. (1995)",
     "\"Pricing Derivatives on Financial Securities Subject to Credit Risk.\" "
     "Journal of Finance, 50(1), 53-85."),

    ("[5]", "Leland, H.E. (1994)",
     "\"Corporate Debt Value, Bond Covenants, and Optimal Capital Structure.\" "
     "Journal of Finance, 49(4), 1213-1252."),

    ("[6]", "Collin-Dufresne, P. & Goldstein, R.S. (2001)",
     "\"Do Credit Spreads Reflect Stationary Leverage Ratios?\" "
     "Journal of Finance, 56(5), 1929-1957."),

    ("[7]", "Glasserman, P. (2003)",
     "Monte Carlo Methods in Financial Engineering. "
     "Springer. (Brownian bridge barrier correction, antithetic variates, variance reduction.)"),

    ("[8]", "Longstaff, F.A. & Schwartz, E.S. (2001)",
     "\"Valuing American Options by Simulation: A Simple Least-Squares Approach.\" "
     "Review of Financial Studies, 14(1), 113-147."),

    ("[9]", "Leland, H.E. & Toft, K.B. (1996)",
     "\"Optimal Capital Structure, Endogenous Bankruptcy, and the Term Structure of Credit Spreads.\" "
     "Journal of Finance, 51(3), 987-1019."),
]

for i, (num, author, title) in enumerate(refs):
    y = 1.1 + i * 0.66
    add_textbox(slide, 0.5, y, 0.6, 0.3, num, font_size=12, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, 1.1, y, 3.0, 0.3, author, font_size=12, color=DARK_BG, bold=True)
    add_textbox(slide, 4.1, y, 9.0, 0.55, title, font_size=11, color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════
#                        APPENDIX
# ════════════════════════════════════════════════════════════════════════


# ════════════════════════════════════════════════════════════════════════
# APPENDIX A: Monte Carlo Pricing Engine
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, MED_GRAY)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "Appendix A: Monte Carlo Pricing Engine",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "The MC engine prices each bond by simulating 25,000 asset paths, checking for barrier "
    "crossings via Brownian bridge [7], processing coupons (cash or PIK), and computing "
    "Z-spreads on a cash-equivalent basis.")

steps = [
    ("1. Asset Evolution", "GBM [1]: V(t+\u0394t) = V(t) \u00b7 exp[(\u03bc - \u00bd\u03c3\u00b2)\u0394t + \u03c3\u221a\u0394t \u00b7 Z]"),
    ("2. First-Passage Default", "Brownian bridge barrier-crossing correction [2,7]. Crossing prob: p = exp(-2x\u2080x\u2081 / \u03c3\u00b2\u0394t)"),
    ("3. Coupon Processing", "Cash \u2192 pay coupon; PIK \u2192 accrete to notional; Toggle \u2192 evaluate exercise model [8]"),
    ("4. Default Recovery", "R(N_current) via DynamicRecoverySpec, discounted to valuation date"),
    ("5. Maturity Payoff", "Surviving paths: discount terminal notional N(T). For PIK: N(T) = N\u2080 \u00b7 (1 + c/f)\u207f"),
]

for i, (title, desc) in enumerate(steps):
    y = 1.7 + i * 0.9
    add_textbox(slide, 0.5, y, 2.5, 0.35, title, font_size=16, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, 3.0, y, 10, 0.7, desc, font_size=14, color=DARK_TEXT)

# Variance reduction
add_accent_bar(slide, 0.5, 6.2, 12.3, 0.04, ACCENT_BLUE)
add_textbox(slide, 0.5, 6.3, 12, 0.3, "Variance Reduction & Z-Spread Computation",
            font_size=16, color=DARK_BG, bold=True)
add_multi_text(slide, 0.5, 6.6, 12, 0.8, [
    ("\u2022 Antithetic variates [7]: negate normal draws for half the paths (halves variance)", False, DARK_TEXT),
    ("\u2022 Common random numbers (CRN): per-path deterministic RNG for calibration consistency", False, DARK_TEXT),
    ("\u2022 Z-spread solved as: \u03a3 CF\u1d62 \u00b7 D(t\u1d62) \u00b7 exp(-z\u00b7t\u1d62) = MC PV  (cash-equivalent basis)", False, DARK_TEXT),
], font_size=12)

add_footnote(slide, "[1] Merton (1974)  [2] Black & Cox (1976)  [7] Glasserman (2003)  [8] Longstaff & Schwartz (2001)")


# ════════════════════════════════════════════════════════════════════════
# APPENDIX B: Detailed Results — BB+ to BB-
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, MED_GRAY)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "Appendix B: Detailed Results \u2014 BB+ to BB-",
            font_size=28, color=DARK_BG, bold=True)

for idx, (issuer, data) in enumerate([
    ("BB+", [
        ["Cash-Pay",   "111.91", "115bp", "5.73%", "4.11", "4.75%", "12.88%", "0.0%"],
        ["Full PIK",   "110.62", "142bp", "6.01%", "4.86", "5.84%", "12.88%", "100.0%"],
        ["PIK Toggle", "111.85", "116bp", "5.74%", "4.11", "4.80%", "12.88%", "0.2%"],
    ]),
    ("BB", [
        ["Cash-Pay",   "110.92", "136bp", "5.95%", "4.10", "5.58%", "13.86%", "0.0%"],
        ["Full PIK",   "109.27", "171bp", "6.31%", "4.86", "6.99%", "13.86%", "100.0%"],
        ["PIK Toggle", "110.82", "138bp", "5.97%", "4.10", "5.67%", "13.86%", "0.3%"],
    ]),
    ("BB-", [
        ["Cash-Pay",   "109.97", "156bp", "6.16%", "4.10", "6.39%", "14.64%", "0.0%"],
        ["Full PIK",   "108.02", "198bp", "6.59%", "4.85", "8.05%", "14.64%", "100.0%"],
        ["PIK Toggle", "108.74", "183bp", "6.43%", "4.08", "7.45%", "14.64%", "4.5%"],
    ]),
]):
    y = 1.1 + idx * 1.9
    add_textbox(slide, 0.5, y, 2, 0.4, issuer, font_size=18, color=ACCENT_BLUE, bold=True)
    add_table(slide, 0.5, y + 0.35, 12.3, 1.2,
              ["Structure", "Price", "Z-Spread", "YTM", "ModDur", "E[Loss]", "Def Rate", "PIK%"],
              data,
              col_widths=[1.8, 1.2, 1.2, 1.2, 1.2, 1.2, 1.2, 1.2],
              font_size=11)


# ════════════════════════════════════════════════════════════════════════
# APPENDIX C: Detailed Results — B+ to B-
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, MED_GRAY)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "Appendix C: Detailed Results \u2014 B+ to B-",
            font_size=28, color=DARK_BG, bold=True)

for idx, (issuer, data) in enumerate([
    ("B+", [
        ["Cash-Pay",   "108.21", "194bp", "6.55%", "4.08", "7.89%", "16.77%", "0.0%"],
        ["Full PIK",   "105.46", "255bp", "7.18%", "4.84", "10.24%", "16.77%", "100.0%"],
        ["PIK Toggle", "106.51", "232bp", "6.94%", "4.07", "9.34%", "16.77%", "6.6%"],
    ]),
    ("B", [
        ["Cash-Pay",   "106.39", "235bp", "6.97%", "4.06", "9.44%", "19.36%", "0.0%"],
        ["Full PIK",   "102.66", "320bp", "7.85%", "4.83", "12.62%", "19.36%", "100.0%"],
        ["PIK Toggle", "102.32", "327bp", "7.93%", "4.03", "12.91%", "19.36%", "30.7%"],
    ]),
    ("B-", [
        ["Cash-Pay",   "102.97", "312bp", "7.78%", "4.03", "12.35%", "24.48%", "0.0%"],
        ["Full PIK",    "97.30", "448bp", "9.19%", "4.80", "17.18%", "24.48%", "100.0%"],
        ["PIK Toggle",  "96.73", "462bp", "9.34%", "3.98", "17.67%", "24.48%", "63.7%"],
    ]),
]):
    y = 1.1 + idx * 1.9
    add_textbox(slide, 0.5, y, 2, 0.4, issuer, font_size=18, color=ACCENT_BLUE, bold=True)
    add_table(slide, 0.5, y + 0.35, 12.3, 1.2,
              ["Structure", "Price", "Z-Spread", "YTM", "ModDur", "E[Loss]", "Def Rate", "PIK%"],
              data,
              col_widths=[1.8, 1.2, 1.2, 1.2, 1.2, 1.2, 1.2, 1.2],
              font_size=11)


# ════════════════════════════════════════════════════════════════════════
# APPENDIX D: Detailed Results — CCC+ to CCC-
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, MED_GRAY)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "Appendix D: Detailed Results \u2014 CCC+ to CCC-",
            font_size=28, color=DARK_BG, bold=True)

for idx, (issuer, data) in enumerate([
    ("CCC+", [
        ["Cash-Pay",    "97.98", "431bp", "9.01%", "3.99", "16.60%", "30.82%", "0.0%"],
        ["Full PIK",    "89.91", "638bp", "11.19%", "4.77", "23.47%", "30.82%", "100.0%"],
        ["PIK Toggle",  "89.59", "647bp", "11.28%", "3.90", "23.74%", "30.82%", "69.2%"],
    ]),
    ("CCC", [
        ["Cash-Pay",    "93.44", "545bp", "10.21%", "3.94", "20.47%", "35.74%", "0.0%"],
        ["Full PIK",    "83.51", "818bp", "13.10%", "4.73", "28.92%", "35.74%", "100.0%"],
        ["PIK Toggle",  "83.27", "825bp", "13.18%", "3.83", "29.12%", "35.74%", "79.3%"],
    ]),
    ("CCC-", [
        ["Cash-Pay",    "85.93", "748bp", "12.36%", "3.86", "26.85%", "44.27%", "0.0%"],
        ["Full PIK",    "72.78", "1158bp", "16.75%", "4.67", "38.05%", "44.27%", "100.0%"],
        ["PIK Toggle",  "72.83", "1156bp", "16.73%", "3.70", "38.01%", "44.27%", "78.3%"],
    ]),
]):
    y = 1.1 + idx * 1.9
    add_textbox(slide, 0.5, y, 2, 0.4, issuer, font_size=18, color=ACCENT_BLUE, bold=True)
    add_table(slide, 0.5, y + 0.35, 12.3, 1.2,
              ["Structure", "Price", "Z-Spread", "YTM", "ModDur", "E[Loss]", "Def Rate", "PIK%"],
              data,
              col_widths=[1.8, 1.2, 1.2, 1.2, 1.2, 1.2, 1.2, 1.2],
              font_size=11)


# ════════════════════════════════════════════════════════════════════════
# APPENDIX E: Duration Analysis
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, MED_GRAY)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "Appendix E: Modified Duration Across Structures",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "PIK bonds have ~0.75yr longer duration than cash-pay (~18% more rate sensitivity). "
    "This is because PIK concentrates all value at maturity, extending the weighted-average cashflow timing.")

add_table(slide, 0.5, 1.7, 8, 3.8,
          ["Issuer", "Cash ModDur", "PIK ModDur", "Toggle ModDur", "PIK - Cash", "Toggle - Cash"],
          [
              ["BB+",  "4.11", "4.86", "4.11", "+0.75", "0.00"],
              ["BB",   "4.10", "4.86", "4.10", "+0.76", "0.00"],
              ["BB-",  "4.10", "4.85", "4.08", "+0.75", "-0.02"],
              ["B+",   "4.08", "4.84", "4.07", "+0.76", "-0.01"],
              ["B",    "4.06", "4.83", "4.03", "+0.77", "-0.03"],
              ["B-",   "4.03", "4.80", "3.98", "+0.77", "-0.05"],
              ["CCC+", "3.99", "4.77", "3.90", "+0.78", "-0.09"],
              ["CCC",  "3.94", "4.73", "3.83", "+0.79", "-0.11"],
              ["CCC-", "3.86", "4.67", "3.70", "+0.81", "-0.16"],
          ],
          col_widths=[1.0, 1.2, 1.2, 1.4, 1.2, 1.2])

add_multi_text(slide, 9.0, 1.7, 4, 5.5, [
    ("Key patterns:", True, DARK_BG),
    ("", False, DARK_BG),
    ("1. PIK extends duration by ~0.75yr", True, ACCENT_BLUE),
    ("All coupons accreted to maturity as terminal notional, concentrating cashflow at the longest point.", False, DARK_TEXT),
    ("", False, DARK_BG),
    ("2. Duration gap is stable across ratings", True, ACCENT_BLUE),
    ("BB+ to CCC-: PIK - Cash ranges from +0.75 to +0.81yr. Cashflow structure \u2014 not credit \u2014 drives the gap.", False, DARK_TEXT),
    ("", False, DARK_BG),
    ("3. Toggle converges to cash at high quality", True, ACCENT_BLUE),
    ("BB+: toggle barely fires (0.2%), so duration matches cash. CCC-: toggle fires 78%, but duration stays near cash since non-PIK periods revert to cash timing.", False, DARK_TEXT),
    ("", False, DARK_BG),
    ("Hedging implication: PIK bonds have ~18% greater rate sensitivity per dollar of notional.", True, RED),
], font_size=12)


# ════════════════════════════════════════════════════════════════════════
# APPENDIX F: Calibration — Barrier vs Vol
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, MED_GRAY)

add_textbox(slide, 0.5, 0.3, 12, 0.6, "Appendix F: Calibration \u2014 Barrier vs Asset Vol",
            font_size=28, color=DARK_BG, bold=True)

add_callout_box(slide,
    "The PIK premium is robust to calibration method \u2014 whether we calibrate via barrier or volatility, "
    "the full-PIK premium is identical. This validates that our results are not an artifact of the calibration approach.")

add_table(slide, 0.5, 1.6, 12.3, 3.5,
          ["Issuer", "Mkt", "Param", "B (cal)", "\u03c3 (cal)", "Cash Z", "PIK Z", "PIK-Cash", "Toggle-Cash"],
          [
              ["BB+",  "150bp", "barrier", "147.2", "20.0%", "149bp", "206bp", "+58bp", "+1bp"],
              ["",     "",      "vol",     "136.8", "22.3%", "149bp", "206bp", "+58bp", "+3bp"],
              ["B+",   "250bp", "barrier", "117.5", "25.0%", "253bp", "358bp", "+105bp", "+43bp"],
              ["",     "",      "vol",     "106.8", "28.1%", "253bp", "358bp", "+105bp", "+56bp"],
              ["B-",   "400bp", "barrier", "81.6",  "35.0%", "407bp", "616bp", "+209bp", "+214bp"],
              ["",     "",      "vol",     "70.0",  "40.2%", "407bp", "616bp", "+209bp", "+214bp"],
              ["CCC",  "700bp", "barrier", "84.9",  "40.0%", "698bp", "1096bp", "+398bp", "+397bp"],
              ["",     "",      "vol",     "69.6",  "47.8%", "698bp", "1096bp", "+398bp", "+397bp"],
          ],
          col_widths=[1.0, 1.0, 1.0, 1.2, 1.2, 1.2, 1.2, 1.5, 1.5])

add_multi_text(slide, 0.5, 5.4, 12, 1.5, [
    ("Result: PIK premiums are identical regardless of calibration parameter", True, GREEN),
    ("\u2022 Full-PIK spread premium is the same under barrier or vol calibration", False, DARK_TEXT),
    ("\u2022 Toggle premiums show small differences (toggle exercise responds differently to barrier vs vol)", False, DARK_TEXT),
    ("\u2022 This validates model robustness \u2014 PIK pricing is not an artifact of the calibration method", False, DARK_TEXT),
], font_size=14)

add_footnote(slide, "Calibration via bisection on either barrier B or asset volatility \u03c3 "
             "to match cash Z-spread to observed market spread. See [7] for MC methodology.")


# ════════════════════════════════════════════════════════════════════════
# APPENDIX G: Uncalibrated Breakeven Spreads
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_accent_bar(slide, 0, 0, 13.333, 0.08, MED_GRAY)

add_textbox(slide, 0.5, 0.3, 10, 0.6, "Appendix G: Uncalibrated Breakeven Spreads",
            font_size=28, color=WHITE, bold=True)

add_textbox(slide, 0.5, 0.9, 12, 0.4,
            "Terminal-barrier Merton [1] (not calibrated to market). Cash spreads are narrower than market due to first-passage mismatch [2].",
            font_size=14, color=MED_GRAY)

add_table(slide, 0.5, 1.4, 12.3, 3.8,
          ["Issuer", "Cash Z-Spread", "PIK Z-Spread", "Toggle Z-Spread", "PIK - Cash", "Toggle - Cash"],
          [
              ["BB+",  "115bp", "142bp", "116bp", "+27bp", "+1bp"],
              ["BB",   "136bp", "171bp", "138bp", "+35bp", "+2bp"],
              ["BB-",  "156bp", "198bp", "183bp", "+42bp", "+27bp"],
              ["B+",   "194bp", "255bp", "232bp", "+61bp", "+38bp"],
              ["B",    "235bp", "320bp", "327bp", "+85bp", "+93bp"],
              ["B-",   "312bp", "448bp", "462bp", "+136bp", "+150bp"],
              ["CCC+", "431bp", "638bp", "647bp", "+207bp", "+216bp"],
              ["CCC",  "545bp", "818bp", "825bp", "+273bp", "+280bp"],
              ["CCC-", "748bp", "1158bp", "1156bp", "+409bp", "+408bp"],
          ],
          col_widths=[1.3, 2.0, 2.0, 2.0, 1.8, 1.8])

add_multi_text(slide, 0.5, 5.5, 12, 1.8, [
    ("These are CASH-EQUIVALENT Z-spreads from the uncalibrated MC run.", False, LIGHT_GRAY),
    ("PIK differentials are tighter here because the uncalibrated barrier inflates cash-pay spreads disproportionately.", False, MED_GRAY),
    ("See MC-Calibrated results (main deck slide 8) for market-consistent PIK premiums.", False, ACCENT_GOLD),
], font_size=14)

add_footnote(slide, "[1] Merton (1974)  [2] Black & Cox (1976). "
             "Terminal-barrier calibration: from_target_pd matches issuer PD to Merton barrier.",
             color=MED_GRAY)


# ════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════
output_path = "finstack-py/examples/notebooks/valuations/PIK_Structural_Credit_Pricing.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
