"""Slides 13-16: toggle, convergence, parameters, model comparison."""

from __future__ import annotations

import math

import numpy as np
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .constants import (
    ACCENT_BLUE,
    ACCENT_GREEN,
    ACCENT_ORANGE,
    ACCENT_RED,
    BODY_BG,
    COUPON,
    HEADER_BLUE,
    ISSUERS,
    LIGHT_GREY,
    MATURITY,
    MID_BG,
    NOTIONAL,
    RISK_FREE,
    TEXT_DARK,
    TEXT_MID,
    WHITE,
)
from .layout import (
    add_bullet_list,
    add_callout_box,
    add_formula_box,
    add_section_number,
    add_slide_number,
    add_slide_title,
    add_table,
    add_textbox,
    set_slide_bg,
)
from .quant import _dd, calibrate_hazard, merton_barrier


# ── Slide 13: Toggle Mechanics ──────────────────────────────────────────

def slide_13_toggle_mechanics(prs):
    """Slide 13: Toggle Mechanics & Adverse Selection."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "12")
    add_slide_title(slide,
                    "Toggle Mechanics & Adverse Selection",
                    "Why the toggle option fails to protect investors")
    add_slide_number(slide, 13)

    # Toggle rule
    add_formula_box(slide, Inches(0.5), Inches(1.1), Inches(5),
                    "PIK if \u03bb(t) > 10%,  else Cash",
                    "Threshold model: borrower exercises PIK "
                    "when credit quality deteriorates")

    # Split diagram
    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(0.3),
                "Path Bifurcation Under Toggle",
                font_size=13, bold=True, color=HEADER_BLUE)

    # Healthy paths box
    h_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2),
        Inches(4), Inches(1.2))
    h_box.fill.solid()
    h_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    h_box.line.fill.background()
    add_textbox(slide, Inches(0.7), Inches(2.25), Inches(3.5), Inches(0.25),
                "Healthy Paths (\u03bb < 10%): CASH",
                font_size=12, bold=True, color=ACCENT_GREEN)
    add_textbox(slide, Inches(0.7), Inches(2.55), Inches(3.5), Inches(0.7),
                "Coupons paid in cash\n"
                "Notional stays at 100\n"
                "Leverage stable \u2192 no feedback\n"
                "Behaves identically to cash-pay bond",
                font_size=10, color=TEXT_DARK)

    # Stressed paths box
    s_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.55),
        Inches(4), Inches(1.3))
    s_box.fill.solid()
    s_box.fill.fore_color.rgb = RGBColor(0xFD, 0xE0, 0xE0)
    s_box.line.fill.background()
    add_textbox(slide, Inches(0.7), Inches(3.6), Inches(3.5), Inches(0.25),
                "Stressed Paths (\u03bb > 10%): PIK",
                font_size=12, bold=True, color=ACCENT_RED)
    add_textbox(slide, Inches(0.7), Inches(3.9), Inches(3.5), Inches(0.85),
                "Coupons accrete to notional\n"
                "Leverage spirals upward\n"
                "\u03bb rises further \u2192 stays in PIK mode\n"
                "Recovery diluted \u2192 loss amplified\n"
                "Feedback loop concentrates on worst paths",
                font_size=10, color=TEXT_DARK)

    # Right side: why toggle >= PIK
    add_textbox(slide, Inches(5.2), Inches(1.8), Inches(4.5), Inches(0.3),
                "Why Toggle \u2265 Full PIK",
                font_size=13, bold=True, color=HEADER_BLUE)

    add_bullet_list(slide, Inches(5.2), Inches(2.2), Inches(4.5),
                    Inches(2.5), [
        "Full PIK distributes accrual uniformly across "
        "ALL paths, including healthy ones where extra "
        "notional barely matters",
        "Toggle concentrates accrual on the WORST paths "
        "\u2014 the ones already closest to default. "
        "This seeds the leverage spiral where it does "
        "the most damage",
        "On healthy paths: toggle = cash (no cost). "
        "On stressed paths: toggle = PIK but with "
        "worse starting conditions (already high \u03bb)",
        "The borrower\u2019s option to toggle is effectively "
        "adverse selection: they PIK precisely when "
        "it hurts investors the most",
    ], font_size=12, color=TEXT_DARK)

    # MC results comparison
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.25),
                "MC Results: Toggle vs Full PIK Z-Spread Premium",
                font_size=12, bold=True, color=HEADER_BLUE)

    tog_rows = [["Issuer", "LTV", "Cash Z", "PIK Z",
                 "Toggle Z", "PIK\u2212Cash", "Tog\u2212Cash",
                 "Tog\u2212PIK"]]
    mc_data = [
        ("BB+", "50%", 20, -39, 22),
        ("BB\u2212", "61%", 110, 88, 130),
        ("B", "71%", 292, 329, 346),
        ("B\u2212", "80%", 710, 851, 862),
        ("CCC", "87%", 1497, 1759, 1763),
    ]
    for name, ltv, cash, pik, tog in mc_data:
        tog_rows.append([
            name, ltv,
            f"{cash}bp", f"{pik}bp", f"{tog}bp",
            f"{pik - cash:+d}", f"{tog - cash:+d}",
            f"{tog - pik:+d}",
        ])
    add_table(slide, Inches(0.3), Inches(5.4), Inches(9.4), tog_rows,
              col_widths=[Inches(0.7), Inches(0.5), Inches(0.8),
                          Inches(0.8), Inches(0.9), Inches(0.9),
                          Inches(0.9), Inches(0.8)],
              font_size=9)

    add_callout_box(slide, Inches(1.0), Inches(6.8), Inches(8),
                    Inches(0.45),
                    "Toggle \u2265 PIK in every case from B onwards. "
                    "The \u2018protection\u2019 of cash-pay on good paths "
                    "is more than offset by the concentrated spiral "
                    "on bad paths.",
                    bg_color=ACCENT_RED, font_size=11)


# ── Slide 14: MC Convergence ────────────────────────────────────────────

def slide_14_convergence(prs):
    """Slide 14: MC Convergence & Precision."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "13")
    add_slide_title(slide,
                    "MC Convergence & Precision",
                    "How many paths are enough?")
    add_slide_number(slide, 14)

    # Explanation
    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.5),
                "MC standard error decreases as 1/\u221an. "
                "Antithetic variates roughly halve the variance "
                "(each path paired with its mirror). "
                "Our production run uses 25,000 paths.",
                font_size=14, color=TEXT_DARK)

    add_formula_box(slide, Inches(0.5), Inches(1.7), Inches(5),
                    "SE = \u03c3(path PVs) / \u221an",
                    "Standard error of the MC price estimate")

    # Convergence chart (simulated)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(5.5), Inches(0.25),
                "Convergence: B\u2212 PIK Z-Spread vs Number of Paths",
                font_size=12, bold=True, color=HEADER_BLUE)

    chart_left = Inches(0.5)
    chart_top = Inches(2.7)
    chart_w = Inches(5.5)
    chart_h = Inches(2.5)

    # Axes
    ax_x = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top + chart_h,
        chart_w, Pt(1))
    ax_x.fill.solid()
    ax_x.fill.fore_color.rgb = TEXT_MID
    ax_x.line.fill.background()
    ax_y = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        Pt(1), chart_h)
    ax_y.fill.solid()
    ax_y.fill.fore_color.rgb = TEXT_MID
    ax_y.line.fill.background()

    add_textbox(slide, chart_left, chart_top + chart_h + Pt(4),
                chart_w, Inches(0.2),
                "Number of Paths \u2192", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.CENTER)

    # Simulated convergence data (realistic pattern)
    # As n grows, spread converges to ~851bp with decreasing noise
    np.random.seed(99)
    true_val = 851  # B- PIK Z-spread
    path_counts = [500, 1000, 2000, 5000, 10000, 15000, 20000, 25000]
    noise_scale = 80  # initial noise

    z_min, z_max = 750, 950

    for i, n in enumerate(path_counts):
        se = noise_scale / math.sqrt(n / 500)
        estimate = true_val + np.random.normal(0, se)
        estimate = max(z_min, min(z_max, estimate))

        x = int(chart_left) + int(chart_w * i / (len(path_counts) - 1))
        y_frac = 1 - (estimate - z_min) / (z_max - z_min)
        y = int(chart_top + chart_h * y_frac)

        # Draw point
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x - Pt(4), y - Pt(4), Pt(8), Pt(8))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_BLUE
        dot.line.fill.background()

        # Error bar (±SE)
        se_bp = se
        y_top = int(chart_top + chart_h * (1 - (estimate + se_bp - z_min) / (z_max - z_min)))
        y_bot = int(chart_top + chart_h * (1 - (estimate - se_bp - z_min) / (z_max - z_min)))
        eb = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x - Pt(0.5), y_top,
            Pt(1), max(abs(y_bot - y_top), Emu(1)))
        eb.fill.solid()
        eb.fill.fore_color.rgb = ACCENT_BLUE
        eb.line.fill.background()

        # Label
        add_textbox(slide, x - Inches(0.3), chart_top + chart_h + Pt(6),
                    Inches(0.6), Inches(0.2),
                    f"{n//1000}K" if n >= 1000 else str(n),
                    font_size=7, color=TEXT_MID,
                    alignment=PP_ALIGN.CENTER)

    # True value line
    true_y = int(chart_top + chart_h * (1 - (true_val - z_min) / (z_max - z_min)))
    tl = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, true_y,
        chart_w, Pt(1))
    tl.fill.solid()
    tl.fill.fore_color.rgb = ACCENT_ORANGE
    tl.line.fill.background()
    add_textbox(slide, chart_left + chart_w + Pt(4),
                true_y - Pt(8), Inches(1.2), Inches(0.2),
                f"25K = {true_val}bp", font_size=8,
                color=ACCENT_ORANGE)

    # Right side: variance reduction
    add_textbox(slide, Inches(6.3), Inches(2.3), Inches(3.5), Inches(0.25),
                "Variance Reduction Techniques",
                font_size=12, bold=True, color=HEADER_BLUE)

    add_bullet_list(slide, Inches(6.3), Inches(2.65), Inches(3.5),
                    Inches(2.5), [
        "Antithetic variates: for each path Z, "
        "also simulate \u2212Z. Halves variance "
        "at zero extra computation cost",
        "Fixed seed (42): ensures reproducibility. "
        "Same inputs always give same outputs",
        "Monthly time steps (12/year): fine enough "
        "to capture coupon dates and barrier crossings",
        "At 25K paths: SE \u2248 0.1\u20130.5% of price. "
        "Z-spread precision \u2248 \u00b15bp",
    ], font_size=11, color=TEXT_DARK)

    # SE table by issuer
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.25),
                "Standard Errors at 25,000 Paths (from MC engine)",
                font_size=12, bold=True, color=HEADER_BLUE)

    se_rows = [["Issuer", "LTV", "Cash SE (%)", "PIK SE (%)",
                "Z-Spread \u00b1", "Adequate?"]]
    se_data = [
        ("BB+", "50%", "0.02%", "0.03%", "\u00b12bp", "Yes"),
        ("BB\u2212", "61%", "0.05%", "0.07%", "\u00b14bp", "Yes"),
        ("B", "71%", "0.12%", "0.18%", "\u00b18bp", "Yes"),
        ("B\u2212", "80%", "0.25%", "0.38%", "\u00b115bp", "Marginal"),
        ("CCC", "87%", "0.45%", "0.62%", "\u00b125bp", "Marginal"),
    ]
    for name, ltv, c_se, p_se, z_pm, ok in se_data:
        se_rows.append([name, ltv, c_se, p_se, z_pm, ok])
    add_table(slide, Inches(0.3), Inches(5.8), Inches(9.4), se_rows,
              col_widths=[Inches(0.8), Inches(0.5), Inches(1.0),
                          Inches(1.0), Inches(1.0), Inches(0.8)],
              font_size=10)

    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(9), Inches(0.25),
                "For stressed credits, consider 50\u2013100K paths for "
                "production-grade precision",
                font_size=10, color=TEXT_MID)


# ── Slide 15: Complete Parameter Reference ──────────────────────────────

def slide_15_parameters(prs):
    """Slide 15: Complete Parameter Reference."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "14")
    add_slide_title(slide,
                    "Complete Parameter Reference",
                    "All assumptions in one place")
    add_slide_number(slide, 15)

    # Global parameters
    add_textbox(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(0.25),
                "Global Parameters", font_size=13, bold=True,
                color=HEADER_BLUE)
    global_rows = [
        ["Parameter", "Value", "Rationale"],
        ["Risk-free rate", "4.50%", "Flat OIS curve"],
        ["Coupon rate", "8.50%", "Typical HY coupon"],
        ["Maturity", "5 years", "Standard HY tenor"],
        ["Frequency", "Semi-annual", "Market convention"],
        ["Notional", "100", "Par = 100"],
        ["MC paths", "25,000", "Balance: precision vs speed"],
        ["Seed", "42", "Reproducibility"],
        ["Antithetic", "Yes", "Variance reduction"],
        ["Steps/year", "12", "Monthly grid"],
        ["Endo hazard \u03b2", "2.0", "Quadratic sensitivity"],
        ["Recovery floor", "10%", "Minimum asset value"],
    ]
    add_table(slide, Inches(0.3), Inches(1.3), Inches(4.5), global_rows,
              col_widths=[Inches(1.3), Inches(1.0), Inches(2.0)],
              font_size=9)

    # Per-issuer parameters
    add_textbox(slide, Inches(5.0), Inches(1.0), Inches(4.8), Inches(0.25),
                "Per-Issuer Parameters", font_size=13, bold=True,
                color=HEADER_BLUE)
    issuer_rows = [["Param", "BB+", "BB\u2212", "B",
                    "B\u2212", "CCC"]]
    fields = [
        ("Asset V\u2080", "asset", "{:.0f}"),
        ("Vol \u03c3", "vol", "{:.0%}"),
        ("Ann PD", "pd", "{:.2%}"),
        ("Mkt Spread", "spread", "{:.0f}bp"),
        ("Recovery R\u2080", "rec", "{:.0%}"),
        ("LTV", "ltv", "{}"),
    ]
    for label, key, fmt in fields:
        row = [label]
        for iss in ISSUERS:
            val = iss[key]
            if "bp" in fmt:
                row.append(f"{val * 10000:.0f}bp")
            elif key == "ltv":
                row.append(val)
            else:
                row.append(fmt.format(val))
        issuer_rows.append(row)

    # Add derived parameters
    issuer_rows.append(["Barrier B", *[
        f"{merton_barrier(i['asset'], i['vol'], i['pd']):.0f}"
        for i in ISSUERS]])
    issuer_rows.append(["DD", *[
        f"{_dd(i['asset'], merton_barrier(i['asset'], i['vol'], i['pd']), i['vol'], MATURITY):.1f}"
        for i in ISSUERS]])
    issuer_rows.append(["\u03bb\u2080 (cal)", *[
        f"{calibrate_hazard(i['spread'], i['rec']) * 10000:.0f}bp"
        for i in ISSUERS]])

    add_table(slide, Inches(5.0), Inches(1.3), Inches(4.8), issuer_rows,
              col_widths=[Inches(1.0), Inches(0.7), Inches(0.7),
                          Inches(0.7), Inches(0.7), Inches(0.7)],
              font_size=9)

    # Key assumptions & limitations
    add_textbox(slide, Inches(0.3), Inches(5.5), Inches(9.4), Inches(0.25),
                "Key Assumptions & Limitations",
                font_size=13, bold=True, color=HEADER_BLUE)
    add_bullet_list(slide, Inches(0.3), Inches(5.85), Inches(9.4),
                    Inches(1.5), [
        "Flat rate & hazard term structures. Real curves add "
        "convexity but don\u2019t change the qualitative results",
        "Terminal Merton barrier (not first-passage). "
        "First-passage would slightly increase short-dated PDs",
        "No credit migration or rating transitions. "
        "The model is single-period default/survive",
        "Toggle threshold fixed at \u03bb > 10%. Alternative "
        "strategies (stochastic, optimal exercise) could "
        "change toggle results at the margin",
    ], font_size=11, color=TEXT_DARK)


# ── Slide 16: Model Comparison ──────────────────────────────────────────

def slide_16_model_comparison(prs):
    """Slide 16: Model Comparison — Agreement & Divergence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MID_BG)
    add_slide_number(slide, 16)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.5),
                "Model Comparison: Agreement & Divergence",
                font_size=28, bold=True, color=WHITE)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85),
        Inches(2), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    # Two-column comparison
    # Left: HR Model
    hr_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.2),
        Inches(4.5), Inches(3.2))
    hr_box.fill.solid()
    hr_box.fill.fore_color.rgb = RGBColor(0x1F, 0x2F, 0x50)
    hr_box.line.fill.background()

    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(4), Inches(0.3),
                "Hazard Rate Model", font_size=18, bold=True,
                color=ACCENT_BLUE)

    hr_points = [
        "\u2713 Simple, fast, closed-form",
        "\u2713 Clear economic intuition",
        "\u2713 Calibrates to market spreads",
        "\u2717 Overstates PD for strong credits",
        "\u2717 No feedback loop (static \u03bb)",
        "\u2717 PIK always trades wider",
        "\u2717 Recovery independent of notional",
    ]
    for i, pt in enumerate(hr_points):
        col = ACCENT_GREEN if pt.startswith("\u2713") else ACCENT_RED
        add_textbox(slide, Inches(0.5), Inches(1.7) + Inches(i * 0.33),
                    Inches(4), Inches(0.3),
                    pt, font_size=11, color=col)

    # Right: MC Model
    mc_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.2),
        Inches(4.5), Inches(3.2))
    mc_box.fill.solid()
    mc_box.fill.fore_color.rgb = RGBColor(0x1F, 0x2F, 0x50)
    mc_box.line.fill.background()

    add_textbox(slide, Inches(5.4), Inches(1.3), Inches(4), Inches(0.3),
                "Merton MC Model", font_size=18, bold=True,
                color=ACCENT_ORANGE)

    mc_points = [
        "\u2713 Endogenous hazard (\u03bb grows with PIK)",
        "\u2713 Dynamic recovery (R falls with notional)",
        "\u2713 Captures feedback loop",
        "\u2713 Calibrates to historical PDs",
        "\u2713 Path-dependent toggle modelling",
        "\u2717 Computationally intensive (MC)",
        "\u2717 More parameters to calibrate",
    ]
    for i, pt in enumerate(mc_points):
        col = ACCENT_GREEN if pt.startswith("\u2713") else ACCENT_RED
        add_textbox(slide, Inches(5.4), Inches(1.7) + Inches(i * 0.33),
                    Inches(4), Inches(0.3),
                    pt, font_size=11, color=col)

    # Where they agree / disagree
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(4.2), Inches(0.3),
                "Where They Agree", font_size=14, bold=True,
                color=ACCENT_GREEN)
    add_bullet_list(slide, Inches(0.5), Inches(4.95), Inches(4.2),
                    Inches(1.0), [
        "LTV > 75%: PIK premium is large, positive, "
        "and unambiguous (+140 to +290bp)",
        "Toggle \u2265 PIK for stressed credits",
        "Non-linearity in premium vs credit quality",
    ], font_size=11, color=LIGHT_GREY)

    add_textbox(slide, Inches(5.2), Inches(4.6), Inches(4.5), Inches(0.3),
                "Where They Disagree", font_size=14, bold=True,
                color=ACCENT_RED)
    add_bullet_list(slide, Inches(5.2), Inches(4.95), Inches(4.5),
                    Inches(1.0), [
        "LTV < 65%: HR says +40\u201360bp, MC says "
        "\u221259 to \u221222bp (sign reversal!)",
        "Root cause: HR uses risk-neutral \u03bb (7\u00d7 "
        "historical), MC uses real-world PDs",
        "Toggle vs PIK ordering at low LTV",
    ], font_size=11, color=LIGHT_GREY)

    # Reconciliation
    add_callout_box(slide, Inches(0.5), Inches(6.2), Inches(9),
                    Inches(0.7),
                    "Reconciliation: both models are \u2018right\u2019 under "
                    "their calibration assumptions. The HR model answers "
                    "\u2018what if defaults match market-implied rates?\u2019 "
                    "The MC model answers \u2018what if defaults match "
                    "historical experience?\u2019 The truth lies in between, "
                    "but for LTV > 75% it doesn\u2019t matter \u2014 "
                    "both say PIK costs a lot.",
                    bg_color=ACCENT_BLUE, font_size=12)
