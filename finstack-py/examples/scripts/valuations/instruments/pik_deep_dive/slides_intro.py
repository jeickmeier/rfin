"""Slides 1-6: title, result, bond setup, HR model, sensitivity, calibration gap."""

from __future__ import annotations

import math

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .constants import (
    ACCENT_BLUE,
    ACCENT_GREEN,
    ACCENT_ORANGE,
    ACCENT_RED,
    BODY_BG,
    COUPON,
    DARK_BG,
    HEADER_BLUE,
    ISSUERS,
    LIGHT_GREY,
    MATURITY,
    NOTIONAL,
    TEXT_DARK,
    TEXT_MID,
    WHITE,
)
from .layout import (
    add_bullet_list,
    add_callout_box,
    add_formula_box,
    add_section_number,
    add_slide_number,
    add_slide_title,
    add_table,
    add_textbox,
    set_slide_bg,
)
from .quant import calibrate_hazard, hr_bond_price, price_to_zspread


# ── Slide 1: Title ───────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(1), Inches(2), Inches(8), Inches(1.2),
                "PIK Coupon Modelling", font_size=40, bold=True, color=WHITE)
    add_textbox(slide, Inches(1), Inches(3.1), Inches(8), Inches(0.8),
                "The Deep Dive", font_size=28, color=ACCENT_BLUE)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.8),
                "Model mechanics, parameters, Monte Carlo paths,\n"
                "and the structural feedback loop",
                font_size=16, color=LIGHT_GREY)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.95), Inches(3), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()
    add_textbox(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.4),
                "Companion to: PIK Coupon Pricing \u2014 "
                "How Much Extra Spread Is Enough?",
                font_size=12, color=TEXT_MID)


# ── Slide 2: The Result to Explain ──────────────────────────────────────

def slide_02_result_to_explain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "1")
    add_slide_title(slide, "The Result to Explain",
                    "From the executive deck: the hockey-stick PIK premium")
    add_slide_number(slide, 2)

    # Recreate the hockey-stick bar chart
    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(4), Inches(0.3),
                "PIK Z-Spread Premium by Market Spread",
                font_size=14, bold=True, color=HEADER_BLUE)

    premiums = [
        ("50bp", -65), ("100bp", -54), ("200bp", -31), ("300bp", -2),
        ("400bp", +37), ("600bp", +107), ("850bp", +216), ("1200bp", +287),
    ]
    max_abs = 287

    for i, (label, prem) in enumerate(premiums):
        y = Inches(1.5) + Inches(i * 0.42)
        add_textbox(slide, Inches(0.5), y, Inches(0.7), Inches(0.3),
                    label, font_size=9, color=TEXT_MID)
        zero_x = Inches(2.8)
        if prem >= 0:
            bar_w = Inches(1.8) * prem / max_abs
            col = (ACCENT_RED if prem > 80
                   else ACCENT_ORANGE if prem > 30 else ACCENT_GREEN)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(zero_x), y + Pt(2),
                int(bar_w), Pt(14))
            bar.fill.solid()
            bar.fill.fore_color.rgb = col
            bar.line.fill.background()
            add_textbox(slide, int(zero_x) + int(bar_w) + Pt(4),
                        y - Pt(1), Inches(0.6), Inches(0.3),
                        f"+{prem}", font_size=9, bold=True, color=col)
        else:
            bar_w = Inches(1.8) * abs(prem) / max_abs
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(zero_x) - int(bar_w),
                y + Pt(2), int(bar_w), Pt(14))
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT_GREEN
            bar.line.fill.background()
            add_textbox(slide, int(zero_x) - int(bar_w) - Inches(0.5),
                        y - Pt(1), Inches(0.5), Inches(0.3),
                        str(prem), font_size=9, bold=True,
                        color=ACCENT_GREEN, alignment=PP_ALIGN.RIGHT)

    # Zero line
    zl = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, int(Inches(2.8)), Inches(1.5),
        Pt(2), Inches(3.4))
    zl.fill.solid()
    zl.fill.fore_color.rgb = TEXT_MID
    zl.line.fill.background()

    # Right side: key questions
    add_textbox(slide, Inches(5.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "Questions This Deck Answers",
                font_size=14, bold=True, color=HEADER_BLUE)
    add_bullet_list(slide, Inches(5.3), Inches(1.5), Inches(4.5),
                    Inches(4.5), [
        "Why does the PIK premium flip sign around 300bp?",
        "What drives the non-linearity \u2014 why isn\u2019t it "
        "proportional to spread?",
        "How does the hazard-rate model price PIK, and "
        "where does it break down?",
        "What is the Merton structural model and how is "
        "it calibrated?",
        "How do endogenous hazard and dynamic recovery "
        "create the feedback loop?",
        "What do Monte Carlo paths actually look like?",
        "Why does the toggle option fail to protect "
        "investors?",
        "What are the key parameter assumptions and "
        "how sensitive are results?",
    ], font_size=12, color=TEXT_DARK)


# ── Slide 3: Bond Setup & Cash Flows ────────────────────────────────────

def slide_03_bond_setup(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "2")
    add_slide_title(slide, "Bond Setup & Cash Flow Timing",
                    "The fundamental difference between cash-pay and PIK")
    add_slide_number(slide, 3)

    # Parameters box
    add_callout_box(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.4),
                    f"5-year  |  {COUPON:.1%} semi-annual  |  "
                    f"Par = {NOTIONAL:.0f}  |  "
                    f"Risk-free = {0.045:.2%} flat  |  "
                    f"Recovery = issuer-dependent",
                    bg_color=HEADER_BLUE, font_size=12)

    # Cash-pay timeline
    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(4.2), Inches(0.3),
                "CASH-PAY: 10 coupons + par at maturity",
                font_size=13, bold=True, color=ACCENT_GREEN)

    # Draw timeline
    cpn_amt = COUPON / 2 * NOTIONAL
    cash_times = [f"t={i/2:.1f}" for i in range(1, 11)]
    for i in range(10):
        x = Inches(0.5 + i * 0.42)
        y = Inches(2.2)
        add_textbox(slide, x, y, Inches(0.42), Inches(0.22),
                    f"{cpn_amt:.2f}", font_size=7, color=ACCENT_GREEN,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.2), Inches(0.42), Inches(0.18),
                    f"{(i+1)/2:.1f}Y", font_size=6, color=TEXT_MID,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(4.5), Inches(2.2), Inches(0.8), Inches(0.22),
                f"+ {NOTIONAL:.0f}", font_size=8, bold=True,
                color=ACCENT_GREEN)

    # PIK timeline
    add_textbox(slide, Inches(0.5), Inches(2.8), Inches(4.5), Inches(0.3),
                "FULL PIK: zero coupons \u2192 inflated notional at maturity",
                font_size=13, bold=True, color=ACCENT_ORANGE)

    terminal_ntl = NOTIONAL * (1 + COUPON / 2) ** 10
    for i in range(10):
        x = Inches(0.5 + i * 0.42)
        y = Inches(3.2)
        ntl_i = NOTIONAL * (1 + COUPON / 2) ** (i + 1)
        add_textbox(slide, x, y, Inches(0.42), Inches(0.22),
                    f"N={ntl_i:.1f}", font_size=6, color=ACCENT_ORANGE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.2), Inches(0.42), Inches(0.18),
                    f"{(i+1)/2:.1f}Y", font_size=6, color=TEXT_MID,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(4.5), Inches(3.15), Inches(1.5), Inches(0.3),
                f"\u2192 {terminal_ntl:.2f}", font_size=9, bold=True,
                color=ACCENT_ORANGE)

    # Key insight box
    add_callout_box(slide, Inches(0.5), Inches(3.8), Inches(4.5),
                    Inches(0.55),
                    f"PIK terminal notional = {terminal_ntl:.2f} "
                    f"(+{terminal_ntl - NOTIONAL:.1f}% above par). "
                    f"All risk concentrated at maturity.",
                    bg_color=ACCENT_ORANGE, font_size=11)

    # Right side: why this matters
    add_textbox(slide, Inches(5.5), Inches(1.8), Inches(4.3), Inches(0.3),
                "Why Timing Matters for Pricing",
                font_size=14, bold=True, color=HEADER_BLUE)

    add_bullet_list(slide, Inches(5.5), Inches(2.2), Inches(4.3),
                    Inches(2.5), [
        "Cash-pay: spreads risk across 10 coupon dates. "
        "Early coupons are almost certain to be paid",
        "PIK: concentrates ALL cash flow at maturity. "
        "The single payment is weighted by S(T), the "
        "5-year survival probability",
        "For a credit with S(5Y) = 70%, cash-pay "
        "collects ~95% of early coupons but PIK "
        "gets nothing if the issuer defaults",
        "This timing asymmetry is the root cause "
        "of the PIK premium under hazard-rate pricing",
    ], font_size=12, color=TEXT_DARK)

    # Survival probability table
    add_textbox(slide, Inches(5.5), Inches(4.5), Inches(4.3), Inches(0.25),
                "Survival Probabilities by Issuer",
                font_size=12, bold=True, color=HEADER_BLUE)

    surv_rows = [["Issuer", "LTV", "\u03bb (bp)", "S(1Y)", "S(3Y)", "S(5Y)"]]
    for iss in ISSUERS:
        lam = calibrate_hazard(iss["spread"], iss["rec"])
        surv_rows.append([
            iss["name"].split(" (")[0], iss["ltv"],
            f"{lam * 10000:.0f}",
            f"{math.exp(-lam * 1):.1%}",
            f"{math.exp(-lam * 3):.1%}",
            f"{math.exp(-lam * 5):.1%}",
        ])
    add_table(slide, Inches(5.5), Inches(4.85), Inches(4.3), surv_rows,
              col_widths=[Inches(0.7), Inches(0.5), Inches(0.6),
                          Inches(0.7), Inches(0.7), Inches(0.7)],
              font_size=9)

    add_textbox(slide, Inches(5.5), Inches(6.7), Inches(4), Inches(0.25),
                "S(t) = exp(\u2212\u03bb \u00d7 t) under flat hazard rate",
                font_size=9, color=TEXT_MID)


# ── Slide 4: HR Model ───────────────────────────────────────────────────

def slide_04_hr_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "3")
    add_slide_title(slide, "Hazard Rate Model: How \u03bb Prices the Bond",
                    "Reduced-form pricing under flat hazard rates")
    add_slide_number(slide, 4)

    # Formula
    add_formula_box(slide, Inches(0.5), Inches(1.1), Inches(9),
                    "PV = \u03a3 cpn \u00d7 D(t) \u00d7 S(t)  +  "
                    "N \u00d7 D(T) \u00d7 S(T)  +  "
                    "R \u00d7 N \u00d7 \u03a3 D(t) \u00d7 [S(t\u22121) \u2212 S(t)]",
                    "D(t) = exp(\u2212r\u00d7t)  |  "
                    "S(t) = exp(\u2212\u03bb\u00d7t)  |  "
                    "\u03bb calibrated from market Z-spread via bisection")

    # Three components explained
    components = [
        ("Coupon PV", "\u03a3 cpn \u00d7 D(t) \u00d7 S(t)",
         "Each coupon discounted for time value AND survival. "
         "Cash-pay: 10 small payments. PIK: zero (coupons accrete).",
         ACCENT_GREEN),
        ("Redemption PV", "N \u00d7 D(T) \u00d7 S(T)",
         "Terminal notional discounted by full survival. "
         "Cash-pay: N = 100. PIK: N = 151.26 (accreted).",
         ACCENT_BLUE),
        ("Recovery PV", "R \u00d7 N \u00d7 \u03a3 D(t) \u00d7 \u0394S",
         "Expected recovery on default. Proportional to "
         "incremental default probability each period.",
         ACCENT_ORANGE),
    ]

    for i, (title, formula, desc, col) in enumerate(components):
        y = Inches(2.0) + Inches(i * 0.85)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y,
            Inches(4.5), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = col
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_top = Pt(4)
        p = tf.paragraphs[0]
        p.text = f"{title}: {formula}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        add_textbox(slide, Inches(5.2), y + Pt(4), Inches(4.5),
                    Inches(0.6), desc, font_size=11, color=TEXT_DARK)

    # HR results table
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(9), Inches(0.25),
                "HR Prices: Cash vs PIK at Each Issuer\u2019s Market Hazard Rate",
                font_size=13, bold=True, color=HEADER_BLUE)

    # Library HR results (from finstack hazard-rate engine)
    hr_tbl = [["Issuer", "LTV", "\u03bb (bp)", "Cash PV",
               "PIK PV", "\u0394Price", "Cash Z", "PIK Z", "\u0394Z (bp)"]]
    _hr_data = [
        ("BB+", "50%", "143", "113.35", "111.46", "\u22121.89",
         "84", "124", "+40"),
        ("BB\u2212", "61%", "334", "107.56", "104.88", "\u22122.68",
         "209", "269", "+60"),
        ("B", "71%", "591", "99.76", "96.05", "\u22123.71",
         "388", "479", "+91"),
        ("B\u2212", "80%", "911", "90.33", "85.40", "\u22124.93",
         "627", "764", "+137"),
        ("CCC", "87%", "1468", "76.13", "69.59", "\u22126.54",
         "1047", "1271", "+224"),
    ]
    for row in _hr_data:
        hr_tbl.append(list(row))
    add_table(slide, Inches(0.3), Inches(4.95), Inches(9.4), hr_tbl,
              col_widths=[Inches(0.9), Inches(0.5), Inches(0.7),
                          Inches(0.9), Inches(0.9), Inches(0.8),
                          Inches(0.9), Inches(0.9), Inches(0.8)],
              font_size=9)

    add_callout_box(slide, Inches(1.0), Inches(6.5), Inches(8), Inches(0.6),
                    "Under HR: PIK always trades wider because "
                    "the inflated terminal notional is discounted by S(T), "
                    "the lowest survival probability in the term structure. "
                    "The penalty grows with \u03bb.",
                    bg_color=HEADER_BLUE, font_size=11)


# ── Slide 5: HR Parameter Sensitivity ───────────────────────────────────

def slide_05_hr_sensitivity(prs):
    """Slide 5: HR Parameter Sensitivity."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "4")
    add_slide_title(slide, "HR Model: Parameter Sensitivity",
                    "How maturity, coupon rate, and recovery affect \u0394Z")
    add_slide_number(slide, 5)

    # Sensitivity: vary maturity
    add_textbox(slide, Inches(0.3), Inches(1.1), Inches(3), Inches(0.25),
                "Maturity Sensitivity (B\u2212, \u03bb=630bp)",
                font_size=12, bold=True, color=HEADER_BLUE)

    iss_b_minus = ISSUERS[3]
    lam_bm = calibrate_hazard(iss_b_minus["spread"], iss_b_minus["rec"])
    mat_rows = [["Maturity", "Cash Z", "PIK Z", "\u0394Z"]]
    for mat in [3, 5, 7]:
        c_pv = hr_bond_price(lam_bm, iss_b_minus["rec"], "cash",
                             maturity=mat)
        p_pv = hr_bond_price(lam_bm, iss_b_minus["rec"], "pik",
                             maturity=mat)
        c_z = price_to_zspread(c_pv, maturity=mat) * 10000
        p_z = price_to_zspread(p_pv, maturity=mat) * 10000
        mat_rows.append([f"{mat}Y", f"{c_z:.0f}bp", f"{p_z:.0f}bp",
                         f"{p_z - c_z:+.0f}bp"])
    add_table(slide, Inches(0.3), Inches(1.45), Inches(3), mat_rows,
              col_widths=[Inches(0.6), Inches(0.8), Inches(0.8),
                          Inches(0.7)],
              font_size=10)

    # Sensitivity: vary coupon
    add_textbox(slide, Inches(3.5), Inches(1.1), Inches(3), Inches(0.25),
                "Coupon Sensitivity (B\u2212, 5Y)",
                font_size=12, bold=True, color=HEADER_BLUE)

    cpn_rows = [["Coupon", "Cash Z", "PIK Z", "\u0394Z"]]
    for cpn in [0.06, 0.085, 0.11]:
        c_pv = hr_bond_price(lam_bm, iss_b_minus["rec"], "cash",
                             coupon_rate=cpn)
        p_pv = hr_bond_price(lam_bm, iss_b_minus["rec"], "pik",
                             coupon_rate=cpn)
        c_z = price_to_zspread(c_pv, maturity=MATURITY) * 10000
        p_z = price_to_zspread(p_pv, maturity=MATURITY) * 10000
        cpn_rows.append([f"{cpn:.1%}", f"{c_z:.0f}bp", f"{p_z:.0f}bp",
                         f"{p_z - c_z:+.0f}bp"])
    add_table(slide, Inches(3.5), Inches(1.45), Inches(3), cpn_rows,
              col_widths=[Inches(0.6), Inches(0.8), Inches(0.8),
                          Inches(0.7)],
              font_size=10)

    # Sensitivity: vary recovery
    add_textbox(slide, Inches(6.7), Inches(1.1), Inches(3), Inches(0.25),
                "Recovery Sensitivity (B\u2212, 5Y)",
                font_size=12, bold=True, color=HEADER_BLUE)

    rec_rows = [["Recovery", "Cash Z", "PIK Z", "\u0394Z"]]
    for rec in [0.25, 0.35, 0.45]:
        lam_r = calibrate_hazard(iss_b_minus["spread"], rec)
        c_pv = hr_bond_price(lam_r, rec, "cash")
        p_pv = hr_bond_price(lam_r, rec, "pik")
        c_z = price_to_zspread(c_pv) * 10000
        p_z = price_to_zspread(p_pv) * 10000
        rec_rows.append([f"{rec:.0%}", f"{c_z:.0f}bp", f"{p_z:.0f}bp",
                         f"{p_z - c_z:+.0f}bp"])
    add_table(slide, Inches(6.7), Inches(1.45), Inches(3), rec_rows,
              col_widths=[Inches(0.7), Inches(0.8), Inches(0.8),
                          Inches(0.7)],
              font_size=10)

    # Interpretation
    add_bullet_list(slide, Inches(0.5), Inches(3.6), Inches(9),
                    Inches(1.5), [
        "Longer maturity \u2192 more compounding periods \u2192 "
        "higher terminal notional \u2192 larger PIK penalty. "
        "The 7Y PIK \u0394Z is roughly double the 3Y",
        "Higher coupon \u2192 more accrual per period \u2192 "
        "faster notional growth. An 11% PIK bond has ~30% "
        "more notional at maturity than 6%",
        "Higher recovery \u2192 hazard rate must rise to match "
        "spread \u2192 amplifies survival discount \u2192 larger "
        "\u0394Z. Recovery is often underappreciated as a PIK "
        "sensitivity",
    ], font_size=13, color=TEXT_DARK)

    # Cross-issuer sensitivity
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(9), Inches(0.25),
                "\u0394Z Across All Issuers (5Y, 8.5% coupon, issuer recovery)",
                font_size=12, bold=True, color=HEADER_BLUE)

    # Library HR results across all issuers
    sweep_rows = [["Issuer", "LTV", "\u03bb (bp)", "Cash Z",
                   "PIK Z", "\u0394Z (bp)"]]
    _sweep = [
        ("BB+ (Solid HY)", "50%", "143", "84", "124", "+40"),
        ("BB\u2212 (Mid HY)", "61%", "334", "209", "269", "+60"),
        ("B (Weak HY)", "71%", "591", "388", "479", "+91"),
        ("B\u2212 (Stressed)", "80%", "911", "627", "764", "+137"),
        ("CCC (Deeply Stressed)", "87%", "1468", "1047", "1271", "+224"),
    ]
    for row in _sweep:
        sweep_rows.append(list(row))
    add_table(slide, Inches(0.5), Inches(5.5), Inches(6.5), sweep_rows,
              col_widths=[Inches(1.5), Inches(0.5), Inches(0.7),
                          Inches(0.8), Inches(0.8), Inches(0.8)],
              font_size=10)

    add_callout_box(slide, Inches(7.2), Inches(5.5), Inches(2.5),
                    Inches(1.2),
                    "Key: \u0394Z grows super-linearly with \u03bb. "
                    "A flat +50bp bump across all issuers is "
                    "too much for BB+ and too little for CCC.",
                    bg_color=ACCENT_ORANGE, font_size=10)


# ── Slide 6: The Calibration Gap ────────────────────────────────────────

def slide_06_calibration_gap(prs):
    """Slide 6: The Calibration Gap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "5")
    add_slide_title(slide, "The Calibration Gap",
                    "Market spreads vs historical default rates")
    add_slide_number(slide, 6)

    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.6),
                "Market spreads are not pure default compensation. They "
                "include a credit risk premium, liquidity premium, and "
                "systematic risk loading. Using market \u03bb in the HR model "
                "overstates default probabilities \u2014 especially for "
                "strong credits.",
                font_size=14, color=TEXT_DARK)

    # Gap table
    gap_rows = [["Issuer", "LTV", "Mkt Spread", "\u03bb (cal)",
                 "HR 5Y PD", "Hist 5Y PD", "Ratio",
                 "Risk Premium"]]
    for iss in ISSUERS:
        lam = calibrate_hazard(iss["spread"], iss["rec"])
        hr_pd = 1 - math.exp(-lam * MATURITY)
        hist_pd = 1 - math.exp(-iss["pd"] * MATURITY)
        ratio = hr_pd / hist_pd if hist_pd > 0 else float("inf")
        gap_rows.append([
            iss["name"], iss["ltv"],
            f"{iss['spread'] * 10000:.0f}bp",
            f"{lam * 10000:.0f}bp",
            f"{hr_pd:.1%}", f"{hist_pd:.1%}",
            f"{ratio:.1f}\u00d7",
            "Very high" if ratio > 5 else
            "High" if ratio > 2 else "Moderate",
        ])
    add_table(slide, Inches(0.3), Inches(1.9), Inches(9.4), gap_rows,
              col_widths=[Inches(1.5), Inches(0.5), Inches(0.8),
                          Inches(0.7), Inches(0.8), Inches(0.8),
                          Inches(0.6), Inches(0.9)],
              font_size=10)

    # Bar chart: ratio by issuer
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(4), Inches(0.25),
                "Market-Implied / Historical PD Ratio",
                font_size=12, bold=True, color=HEADER_BLUE)

    ratios = []
    for iss in ISSUERS:
        lam = calibrate_hazard(iss["spread"], iss["rec"])
        hr_pd = 1 - math.exp(-lam * MATURITY)
        hist_pd = 1 - math.exp(-iss["pd"] * MATURITY)
        ratios.append((iss["name"].split(" (")[0], hr_pd / hist_pd))

    max_ratio = max(r for _, r in ratios)
    for i, (name, ratio) in enumerate(ratios):
        y = Inches(4.55) + Inches(i * 0.4)
        add_textbox(slide, Inches(0.5), y, Inches(0.8), Inches(0.25),
                    name, font_size=9, color=TEXT_MID)
        bar_w = Inches(2.5) * ratio / max_ratio
        col = ACCENT_RED if ratio > 4 else (
            ACCENT_ORANGE if ratio > 2 else ACCENT_GREEN)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.4), y + Pt(2),
            int(bar_w), Pt(14))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        add_textbox(slide, Inches(1.4) + int(bar_w) + Pt(4), y - Pt(1),
                    Inches(0.6), Inches(0.25),
                    f"{ratio:.1f}\u00d7", font_size=9, bold=True,
                    color=col)

    # Right side: implications
    add_textbox(slide, Inches(5.0), Inches(4.2), Inches(4.8), Inches(0.25),
                "Implications for PIK Pricing",
                font_size=12, bold=True, color=HEADER_BLUE)
    add_bullet_list(slide, Inches(5.0), Inches(4.55), Inches(4.8),
                    Inches(2.5), [
        "BB+ spread implies 7\u00d7 the historical default "
        "rate. The HR model\u2019s PIK penalty (+40bp) is "
        "driven by an overstated \u03bb",
        "CCC spread implies only 1.3\u00d7 historical PD. "
        "At this level, the PIK penalty (+224bp) is "
        "closer to a \u2018real\u2019 default cost",
        "The structural model uses historical PDs to "
        "calibrate barriers. It answers: what is the "
        "PIK penalty under realistic default assumptions?",
    ], font_size=12, color=TEXT_DARK)

    add_callout_box(slide, Inches(1.0), Inches(6.6), Inches(8),
                    Inches(0.5),
                    "The HR model is an upper bound on PIK cost. "
                    "For strong credits, the true premium may be "
                    "negative under historical defaults.",
                    bg_color=ACCENT_ORANGE, font_size=12)
