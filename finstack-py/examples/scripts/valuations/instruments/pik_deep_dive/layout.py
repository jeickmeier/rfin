"""PPTX layout helper functions."""

from __future__ import annotations

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .constants import (
    ACCENT_BLUE,
    HEADER_BLUE,
    TABLE_ALT_BG,
    TABLE_HEADER_BG,
    TABLE_WHITE_BG,
    TEXT_DARK,
    TEXT_MID,
    TOTAL_SLIDES,
    WHITE,
)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text: str, *,
                font_size=18, bold=False, color=TEXT_DARK,
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items: list[str], *,
                    font_size=16, color=TEXT_DARK, spacing_after=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.level = 0
        p.space_after = spacing_after
    return txBox


def add_callout_box(slide, left, top, width, height, text: str, *,
                    bg_color=ACCENT_BLUE, text_color=WHITE, font_size=14):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = "Calibri"
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_section_number(slide, number: str, left=Inches(0.5), top=Inches(0.3)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_table(slide, left, top, width, rows_data: list[list[str]], *,
              col_widths=None, font_size=11, header_row=True):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, Inches(0.35 * n_rows))
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0 and header_row:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = TEXT_DARK
                    paragraph.alignment = (
                        PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT)
            if r == 0 and header_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_WHITE_BG
    return tbl_shape


def add_slide_number(slide, num: int):
    add_textbox(slide, Inches(8.8), Inches(7.1), Inches(1.2), Inches(0.3),
                f"{num}/{TOTAL_SLIDES}", font_size=10, color=TEXT_MID,
                alignment=PP_ALIGN.RIGHT)


def add_slide_title(slide, title: str, subtitle: str = ""):
    add_textbox(slide, Inches(1.1), Inches(0.2), Inches(8), Inches(0.5),
                title, font_size=26, bold=True, color=HEADER_BLUE)
    if subtitle:
        add_textbox(slide, Inches(1.1), Inches(0.65), Inches(8), Inches(0.4),
                    subtitle, font_size=14, color=TEXT_MID)


def draw_line(slide, x1, y1, x2, y2, color=TEXT_MID, width=Pt(1)):
    """Draw a thin line between two points."""
    left = min(x1, x2)
    top = min(y1, y2)
    w = abs(x2 - x1) or Emu(1)
    h = abs(y2 - y1) or Emu(1)
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, w, h)
    connector.fill.solid()
    connector.fill.fore_color.rgb = color
    connector.line.fill.background()
    return connector


def add_formula_box(slide, left, top, width, formula: str,
                    caption: str = ""):
    """Formula in bold + optional caption below."""
    add_textbox(slide, left, top, width, Inches(0.35),
                formula, font_size=16, bold=True, color=HEADER_BLUE)
    if caption:
        add_textbox(slide, left, top + Inches(0.35), width, Inches(0.25),
                    caption, font_size=11, color=TEXT_MID)
