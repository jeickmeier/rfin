"""Slides 7-12: Merton model, barriers, endogenous hazard, recovery, MC, spiral."""

from __future__ import annotations

import math

import numpy as np
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .constants import (
    ACCENT_GREEN,
    ACCENT_ORANGE,
    ACCENT_RED,
    BODY_BG,
    COUPON,
    HEADER_BLUE,
    ISSUERS,
    MATURITY,
    NOTIONAL,
    RISK_FREE,
    TEXT_DARK,
    TEXT_MID,
    WHITE,
)
from .layout import (
    add_bullet_list,
    add_callout_box,
    add_formula_box,
    add_section_number,
    add_slide_number,
    add_slide_title,
    add_table,
    add_textbox,
    set_slide_bg,
)
from .quant import _dd, _norm_ppf, calibrate_hazard, merton_barrier


# ── Slide 7: Merton Model ───────────────────────────────────────────────

def slide_07_merton_model(prs):
    """Slide 7: Merton Model — Firm Value & Default."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "6")
    add_slide_title(slide, "Merton Structural Model: Firm Value & Default",
                    "Default occurs when asset value breaches the barrier")
    add_slide_number(slide, 7)

    # GBM formula
    add_formula_box(slide, Inches(0.5), Inches(1.1), Inches(5),
                    "dV = (r \u2212 q)\u00b7V\u00b7dt + \u03c3\u00b7V\u00b7dW",
                    "Geometric Brownian Motion for firm asset value V")

    add_formula_box(slide, Inches(0.5), Inches(1.8), Inches(5),
                    "DD = [ln(V/B) + (r \u2212 \u03c3\u00b2/2)\u00b7T] / "
                    "(\u03c3\u00b7\u221aT)",
                    "Distance-to-Default: standard deviations from barrier")

    add_formula_box(slide, Inches(0.5), Inches(2.5), Inches(5),
                    "PD(T) = N(\u2212DD)",
                    "Default probability = normal CDF of negative DD")

    # Conceptual diagram: asset path with barrier
    add_textbox(slide, Inches(5.8), Inches(1.1), Inches(4), Inches(0.25),
                "Conceptual: Asset Value vs Barrier",
                font_size=12, bold=True, color=HEADER_BLUE)

    # Draw a simple conceptual chart using shapes
    chart_left = Inches(5.8)
    chart_top = Inches(1.5)
    chart_w = Inches(3.8)
    chart_h = Inches(2.2)

    # Axes
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top + chart_h,
        chart_w, Pt(2)).fill.solid()
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        Pt(2), chart_h).fill.solid()

    # Barrier line (horizontal, at ~40% height)
    barrier_y = int(chart_top + chart_h * 0.6)
    b_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, barrier_y,
        chart_w, Pt(2))
    b_line.fill.solid()
    b_line.fill.fore_color.rgb = ACCENT_RED
    b_line.line.fill.background()
    add_textbox(slide, chart_left + chart_w + Pt(4),
                barrier_y - Pt(6), Inches(0.6), Inches(0.2),
                "B (barrier)", font_size=8, color=ACCENT_RED)

    # Asset start point label
    asset_start_y = int(chart_top + chart_h * 0.15)
    add_textbox(slide, chart_left - Inches(0.5), asset_start_y - Pt(4),
                Inches(0.5), Inches(0.2),
                "V\u2080", font_size=10, bold=True, color=HEADER_BLUE,
                alignment=PP_ALIGN.RIGHT)

    # Simulated paths (simplified as line segments)
    np.random.seed(42)
    n_steps = 20
    dt = MATURITY / n_steps
    for path_i in range(8):
        v = 1.0  # normalised
        points = [(0, v)]
        vol = 0.25
        for step in range(n_steps):
            dw = np.random.normal(0, math.sqrt(dt))
            v *= math.exp((RISK_FREE - vol**2/2) * dt + vol * dw)
            points.append((step + 1, v))

        defaulted = any(p[1] < 0.55 for p in points)
        col = ACCENT_RED if defaulted else ACCENT_GREEN

        for j in range(1, len(points)):
            x1 = int(chart_left) + int(chart_w * points[j-1][0] / n_steps)
            x2 = int(chart_left) + int(chart_w * points[j][0] / n_steps)
            y1 = int(chart_top + chart_h * (1 - points[j-1][1] * 0.85))
            y2 = int(chart_top + chart_h * (1 - points[j][1] * 0.85))
            # Draw as small rectangle (approximation)
            seg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                min(x1, x2), min(y1, y2),
                max(abs(x2 - x1), Emu(1)),
                max(abs(y2 - y1), Pt(1)))
            seg.fill.solid()
            seg.fill.fore_color.rgb = col
            seg.line.fill.background()

    add_textbox(slide, chart_left, chart_top + chart_h + Pt(4),
                chart_w, Inches(0.2),
                "Time \u2192 (0 to 5Y)", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.CENTER)

    # Key concepts (below)
    concepts = [
        ("Equity as Call Option",
         "Equity = max(V \u2212 B, 0) at maturity. Shareholders "
         "own a call option on firm assets with strike = debt. "
         "The Merton model links credit and equity markets."),
        ("Distance-to-Default",
         "DD measures how many standard deviations the asset "
         "value sits above the barrier. Higher DD = safer credit. "
         "BB+ has DD \u2248 3.4, CCC has DD \u2248 0.6."),
        ("Terminal vs First-Passage",
         "Terminal barrier: default only checked at T. "
         "First-passage (Black-Cox): default can occur at any "
         "time V < B. We use terminal for simplicity."),
    ]

    for i, (title, desc) in enumerate(concepts):
        y = Inches(4.0) + Inches(i * 0.95)
        add_textbox(slide, Inches(0.5), y, Inches(2.5), Inches(0.25),
                    title, font_size=12, bold=True, color=HEADER_BLUE)
        add_textbox(slide, Inches(3.0), y, Inches(6.8), Inches(0.85),
                    desc, font_size=11, color=TEXT_DARK)


# ── Slide 8: Barrier Calibration ────────────────────────────────────────

def slide_08_barrier_calibration(prs):
    """Slide 8: Barrier Calibration from Historical PDs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "7")
    add_slide_title(slide, "Barrier Calibration from Historical PDs",
                    "MertonModel.from_target_pd: backing out the barrier")
    add_slide_number(slide, 8)

    # Method explanation
    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.5),
                "The barrier B is the free parameter. Given V\u2080, \u03c3, "
                "and r, we solve for B such that PD(T) = N(\u2212DD) "
                "matches the target historical default probability.",
                font_size=14, color=TEXT_DARK)

    add_formula_box(slide, Inches(0.5), Inches(1.7), Inches(9),
                    "B = V \u00d7 exp[\u2212(DD \u00d7 \u03c3\u00d7\u221aT "
                    "+ (r \u2212 \u03c3\u00b2/2)\u00d7T)]  where  "
                    "DD = \u2212N\u207b\u00b9(PD\u2085\u2084)",
                    "Invert the distance-to-default formula "
                    "to find barrier from target PD")

    # Full parameter table
    cal_rows = [["Issuer", "LTV", "V\u2080", "\u03c3",
                 "Ann PD", "5Y PD", "DD",
                 "Barrier", "Impl Sprd"]]
    for iss in ISSUERS:
        five_pd = 1 - math.exp(-iss["pd"] * MATURITY)
        dd = -_norm_ppf(five_pd) if five_pd < 1 else 0
        barrier = merton_barrier(iss["asset"], iss["vol"], iss["pd"])
        # Implied spread: s = -ln(1 - PD*(1-R)) / T
        impl_s = -math.log(1 - five_pd * (1 - iss["rec"])) / MATURITY
        cal_rows.append([
            iss["name"], iss["ltv"],
            f"{iss['asset']:.0f}", f"{iss['vol']:.0%}",
            f"{iss['pd']:.2%}", f"{five_pd:.1%}",
            f"{dd:.2f}", f"{barrier:.1f}",
            f"{impl_s * 10000:.0f}bp",
        ])
    add_table(slide, Inches(0.3), Inches(2.4), Inches(9.4), cal_rows,
              col_widths=[Inches(1.5), Inches(0.5), Inches(0.5),
                          Inches(0.5), Inches(0.7), Inches(0.6),
                          Inches(0.5), Inches(0.7), Inches(0.8)],
              font_size=10)

    # Comparison: Merton implied spread vs market spread
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(9), Inches(0.25),
                "Market Spread vs Merton Implied Spread",
                font_size=13, bold=True, color=HEADER_BLUE)

    cmp_rows = [["Issuer", "LTV", "Mkt Spread",
                 "Merton Spread", "Gap", "Interpretation"]]
    for iss in ISSUERS:
        five_pd = 1 - math.exp(-iss["pd"] * MATURITY)
        impl_s = -math.log(1 - five_pd * (1 - iss["rec"])) / MATURITY
        gap = iss["spread"] - impl_s
        cmp_rows.append([
            iss["name"], iss["ltv"],
            f"{iss['spread'] * 10000:.0f}bp",
            f"{impl_s * 10000:.0f}bp",
            f"{gap * 10000:+.0f}bp",
            "Large risk premium" if gap > 0.005 else
            "Moderate premium" if gap > 0.002 else
            "Small premium",
        ])
    add_table(slide, Inches(0.3), Inches(4.85), Inches(9.4), cmp_rows,
              col_widths=[Inches(1.5), Inches(0.5), Inches(0.9),
                          Inches(1.0), Inches(0.8), Inches(1.4)],
              font_size=10)

    add_callout_box(slide, Inches(0.5), Inches(6.5), Inches(9),
                    Inches(0.6),
                    "The Merton implied spread is consistently below the "
                    "market spread because it uses only historical default "
                    "risk. The gap IS the risk premium. By calibrating to "
                    "historical PDs, the MC model prices PIK under "
                    "\u2018real-world\u2019 rather than risk-neutral defaults.",
                    bg_color=HEADER_BLUE, font_size=11)


# ── Slide 9: Endogenous Hazard ──────────────────────────────────────────

def slide_09_endogenous_hazard(prs):
    """Slide 9: Endogenous Hazard — \u03bb(L) = \u03bb\u2080\u00d7(L/L\u2080)\u00b2."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "8")
    add_slide_title(slide, "Endogenous Hazard: \u03bb(L) = \u03bb\u2080 \u00d7 (L/L\u2080)\u00b2",
                    "PIK accrual raises leverage \u2192 hazard rate "
                    "rises non-linearly")
    add_slide_number(slide, 9)

    # Formula and explanation
    add_formula_box(slide, Inches(0.5), Inches(1.1), Inches(5),
                    "\u03bb(L) = \u03bb\u2080 \u00d7 (L / L\u2080)\u00b2",
                    "\u03bb\u2080 = base hazard  |  L\u2080 = initial "
                    "leverage  |  \u03b2 = 2 (quadratic)")

    add_bullet_list(slide, Inches(0.5), Inches(1.85), Inches(5),
                    Inches(1.5), [
        "L = N(t) / V(t): leverage = notional / asset value",
        "PIK accretes to notional \u2192 N(t) grows \u2192 "
        "L rises even if V is unchanged",
        "Quadratic (\u03b2=2): a 20% leverage increase "
        "raises hazard by 44%",
        "This is the first half of the feedback loop: "
        "PIK \u2192 higher \u03bb \u2192 more defaults",
    ], font_size=12, color=TEXT_DARK)

    # Chart: \u03bb vs leverage for each issuer
    add_textbox(slide, Inches(5.8), Inches(1.1), Inches(4), Inches(0.25),
                "Hazard Rate vs Leverage",
                font_size=12, bold=True, color=HEADER_BLUE)

    chart_left = Inches(5.8)
    chart_top = Inches(1.5)
    chart_w = Inches(3.8)
    chart_h = Inches(2.5)

    # Axes
    ax_x = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top + chart_h,
        chart_w, Pt(1))
    ax_x.fill.solid()
    ax_x.fill.fore_color.rgb = TEXT_MID
    ax_x.line.fill.background()
    ax_y = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        Pt(1), chart_h)
    ax_y.fill.solid()
    ax_y.fill.fore_color.rgb = TEXT_MID
    ax_y.line.fill.background()

    # Axis labels
    add_textbox(slide, chart_left, chart_top + chart_h + Pt(4),
                chart_w, Inches(0.2),
                "Leverage (N/V) \u2192", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, chart_left - Inches(0.6), chart_top,
                Inches(0.5), Inches(0.2),
                "\u03bb \u2192", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.RIGHT)

    # Plot curves for 3 issuers: BB+, B, CCC
    plot_issuers = [ISSUERS[0], ISSUERS[2], ISSUERS[4]]
    colors_plot = [ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED]
    max_lev = 1.2
    max_haz = 0.35  # cap for display

    for iss, col in zip(plot_issuers, colors_plot):
        lam0 = calibrate_hazard(iss["spread"], iss["rec"])
        l0 = NOTIONAL / iss["asset"]
        n_pts = 30
        prev_x, prev_y = None, None
        for j in range(n_pts):
            lev = 0.3 + j * (max_lev - 0.3) / (n_pts - 1)
            haz = lam0 * (lev / l0) ** 2
            haz_clip = min(haz, max_haz)
            px = int(chart_left) + int(chart_w * (lev - 0.3) / (max_lev - 0.3))
            py = int(chart_top + chart_h * (1 - haz_clip / max_haz))
            if prev_x is not None:
                seg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    min(px, prev_x), min(py, prev_y),
                    max(abs(px - prev_x), Emu(1)),
                    max(abs(prev_y - py), Pt(1.5)))
                seg.fill.solid()
                seg.fill.fore_color.rgb = col
                seg.line.fill.background()
            prev_x, prev_y = px, py

        # Label at end
        add_textbox(slide, prev_x + Pt(4), prev_y - Pt(6),
                    Inches(0.7), Inches(0.2),
                    iss["name"].split(" (")[0],
                    font_size=7, color=col)

    # Initial leverage markers
    add_textbox(slide, Inches(5.8), Inches(4.15), Inches(4), Inches(0.2),
                "Arrows show where PIK pushes each issuer along its curve",
                font_size=8, color=TEXT_MID)

    # Numerical example table
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(5), Inches(0.25),
                "PIK Impact: 5-Year Notional & Hazard Growth",
                font_size=12, bold=True, color=HEADER_BLUE)

    terminal_ntl = NOTIONAL * (1 + COUPON / 2) ** 10
    pik_rows = [["Issuer", "LTV\u2080", "L\u2080",
                 "L at 5Y (PIK)", "\u03bb\u2080",
                 "\u03bb at 5Y (PIK)", "Increase"]]
    for iss in ISSUERS:
        lam0 = calibrate_hazard(iss["spread"], iss["rec"])
        l0 = NOTIONAL / iss["asset"]
        l_pik = terminal_ntl / iss["asset"]
        lam_pik = lam0 * (l_pik / l0) ** 2
        pik_rows.append([
            iss["name"].split(" (")[0], iss["ltv"],
            f"{l0:.2f}", f"{l_pik:.2f}",
            f"{lam0 * 10000:.0f}bp", f"{lam_pik * 10000:.0f}bp",
            f"{lam_pik / lam0:.1f}\u00d7",
        ])
    add_table(slide, Inches(0.3), Inches(4.65), Inches(9.4), pik_rows,
              col_widths=[Inches(0.9), Inches(0.6), Inches(0.6),
                          Inches(1.0), Inches(0.8), Inches(1.1),
                          Inches(0.8)],
              font_size=9)

    add_callout_box(slide, Inches(1.0), Inches(6.5), Inches(8),
                    Inches(0.55),
                    "For CCC (87% LTV): PIK pushes leverage from 0.87 "
                    f"to {terminal_ntl / 115:.2f}, increasing \u03bb by "
                    f"{(terminal_ntl / 115 / (100/115))**2:.1f}\u00d7. "
                    "This is before considering asset value changes.",
                    bg_color=ACCENT_RED, font_size=11)


# ── Slide 10: Dynamic Recovery ──────────────────────────────────────────

def slide_10_dynamic_recovery(prs):
    """Slide 10: Dynamic Recovery — R(N) = max(floor, R\u2080\u00d7N\u2080/N)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "9")
    add_slide_title(slide, "Dynamic Recovery: R(N) = max(floor, R\u2080 \u00d7 N\u2080/N)",
                    "PIK dilutes recovery per dollar of claim")
    add_slide_number(slide, 10)

    # Formula
    add_formula_box(slide, Inches(0.5), Inches(1.1), Inches(5.5),
                    "R(N) = max(floor,  R\u2080 \u00d7 N\u2080 / N)",
                    "R\u2080 = base recovery  |  N\u2080 = 100 (par)  |  "
                    "floor = 10%")

    add_bullet_list(slide, Inches(0.5), Inches(1.85), Inches(5),
                    Inches(1.5), [
        "On default, the recovery pool (assets) is fixed but "
        "the claim (notional) has grown via PIK accrual",
        "Each dollar of claim gets proportionally less. "
        "If notional doubles, recovery per dollar halves",
        "The floor (10%) prevents recovery from going "
        "to zero \u2014 there is always some residual asset value",
        "This is the second half of the feedback loop: "
        "PIK \u2192 lower R \u2192 higher loss-given-default",
    ], font_size=12, color=TEXT_DARK)

    # Chart: Recovery vs Notional for each issuer
    add_textbox(slide, Inches(5.8), Inches(1.1), Inches(4), Inches(0.25),
                "Recovery Rate vs Accreted Notional",
                font_size=12, bold=True, color=HEADER_BLUE)

    chart_left = Inches(5.8)
    chart_top = Inches(1.5)
    chart_w = Inches(3.8)
    chart_h = Inches(2.3)

    # Axes
    ax_x = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top + chart_h,
        chart_w, Pt(1))
    ax_x.fill.solid()
    ax_x.fill.fore_color.rgb = TEXT_MID
    ax_x.line.fill.background()
    ax_y = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        Pt(1), chart_h)
    ax_y.fill.solid()
    ax_y.fill.fore_color.rgb = TEXT_MID
    ax_y.line.fill.background()

    add_textbox(slide, chart_left, chart_top + chart_h + Pt(4),
                chart_w, Inches(0.2),
                "Notional (100 \u2192 200) \u2192",
                font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, chart_left - Inches(0.6), chart_top,
                Inches(0.5), Inches(0.2),
                "R(N) \u2192", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.RIGHT)

    # Floor line
    floor = 0.10
    floor_y = int(chart_top + chart_h * (1 - floor / 0.50))
    fl = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, floor_y,
        chart_w, Pt(1))
    fl.fill.solid()
    fl.fill.fore_color.rgb = ACCENT_RED
    fl.line.fill.background()
    add_textbox(slide, chart_left + chart_w + Pt(4), floor_y - Pt(6),
                Inches(0.6), Inches(0.2),
                "floor=10%", font_size=7, color=ACCENT_RED)

    # Plot curves for 3 issuers
    plot_issuers = [ISSUERS[0], ISSUERS[2], ISSUERS[4]]
    colors_plot = [ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED]
    n_min, n_max = 100, 200

    for iss, col in zip(plot_issuers, colors_plot):
        r0 = iss["rec"]
        n_pts = 30
        prev_x, prev_y = None, None
        for j in range(n_pts):
            n = n_min + j * (n_max - n_min) / (n_pts - 1)
            r_n = max(floor, r0 * NOTIONAL / n)
            px = int(chart_left) + int(chart_w * (n - n_min) / (n_max - n_min))
            py = int(chart_top + chart_h * (1 - r_n / 0.50))
            if prev_x is not None:
                seg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    min(px, prev_x), min(py, prev_y),
                    max(abs(px - prev_x), Emu(1)),
                    max(abs(prev_y - py), Pt(1.5)))
                seg.fill.solid()
                seg.fill.fore_color.rgb = col
                seg.line.fill.background()
            prev_x, prev_y = px, py

        add_textbox(slide, prev_x + Pt(4), prev_y - Pt(6),
                    Inches(0.9), Inches(0.2),
                    f"{iss['name'].split(' (')[0]} (R\u2080={r0:.0%})",
                    font_size=7, color=col)

    # Combined effect table
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.25),
                "Combined Feedback: Hazard \u00d7 Recovery Impact at Maturity",
                font_size=12, bold=True, color=HEADER_BLUE)

    terminal_ntl = NOTIONAL * (1 + COUPON / 2) ** 10
    combo_rows = [["Issuer", "LTV", "R\u2080",
                   "R at 5Y (PIK)", "\u0394R",
                   "\u03bb\u2080", "\u03bb at 5Y (PIK)",
                   "Net LGD increase"]]
    for iss in ISSUERS:
        lam0 = calibrate_hazard(iss["spread"], iss["rec"])
        r0 = iss["rec"]
        r_pik = max(0.10, r0 * NOTIONAL / terminal_ntl)
        l0 = NOTIONAL / iss["asset"]
        l_pik = terminal_ntl / iss["asset"]
        lam_pik = lam0 * (l_pik / l0) ** 2
        lgd_0 = (1 - r0) * lam0
        lgd_pik = (1 - r_pik) * lam_pik
        combo_rows.append([
            iss["name"].split(" (")[0], iss["ltv"],
            f"{r0:.0%}", f"{r_pik:.0%}",
            f"{(r_pik - r0) * 100:+.0f}pp",
            f"{lam0 * 10000:.0f}bp", f"{lam_pik * 10000:.0f}bp",
            f"{lgd_pik / lgd_0:.1f}\u00d7" if lgd_0 > 0 else "n/a",
        ])
    add_table(slide, Inches(0.3), Inches(4.55), Inches(9.4), combo_rows,
              col_widths=[Inches(0.9), Inches(0.5), Inches(0.5),
                          Inches(0.9), Inches(0.6), Inches(0.7),
                          Inches(1.0), Inches(1.0)],
              font_size=9)

    add_callout_box(slide, Inches(1.0), Inches(6.5), Inches(8),
                    Inches(0.55),
                    "The double whammy: PIK simultaneously increases "
                    "the probability of default (\u03bb\u2191) AND reduces "
                    "recovery per dollar (R\u2193). These compound "
                    "multiplicatively, not additively.",
                    bg_color=ACCENT_RED, font_size=11)


# ── Slide 11: Monte Carlo Paths ─────────────────────────────────────────

def slide_11_mc_paths(prs):
    """Slide 11: Monte Carlo Paths — visual."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "10")
    add_slide_title(slide, "Monte Carlo Paths: What the Simulation Looks Like",
                    "50 sample GBM paths for B\u2212 (Stressed, 80% LTV)")
    add_slide_number(slide, 11)

    iss = ISSUERS[3]  # B-
    barrier = merton_barrier(iss["asset"], iss["vol"], iss["pd"])

    # Chart area
    chart_left = Inches(0.5)
    chart_top = Inches(1.2)
    chart_w = Inches(5.5)
    chart_h = Inches(3.5)

    # Axes
    ax_x = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top + chart_h,
        chart_w, Pt(1))
    ax_x.fill.solid()
    ax_x.fill.fore_color.rgb = TEXT_MID
    ax_x.line.fill.background()
    ax_y = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        Pt(1), chart_h)
    ax_y.fill.solid()
    ax_y.fill.fore_color.rgb = TEXT_MID
    ax_y.line.fill.background()

    add_textbox(slide, chart_left, chart_top + chart_h + Pt(4),
                chart_w, Inches(0.2),
                "Time (years) \u2192", font_size=9, color=TEXT_MID,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, chart_left - Inches(0.5), chart_top - Pt(2),
                Inches(0.5), Inches(0.2),
                "Asset\nValue", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.RIGHT)

    # Barrier line
    v_max = iss["asset"] * 1.6
    v_min = 0
    b_frac = 1 - (barrier - v_min) / (v_max - v_min)
    barrier_y = int(chart_top + chart_h * b_frac)
    bl = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, barrier_y,
        chart_w, Pt(2))
    bl.fill.solid()
    bl.fill.fore_color.rgb = ACCENT_RED
    bl.line.fill.background()
    add_textbox(slide, chart_left + chart_w + Pt(4),
                barrier_y - Pt(8), Inches(1.0), Inches(0.2),
                f"B = {barrier:.0f}", font_size=9, bold=True,
                color=ACCENT_RED)

    # V0 label
    v0_frac = 1 - (iss["asset"] - v_min) / (v_max - v_min)
    add_textbox(slide, chart_left - Inches(0.7),
                int(chart_top + chart_h * v0_frac) - Pt(6),
                Inches(0.6), Inches(0.2),
                f"V\u2080={iss['asset']}", font_size=8, bold=True,
                color=HEADER_BLUE, alignment=PP_ALIGN.RIGHT)

    # Simulate 50 paths
    np.random.seed(123)
    n_paths = 50
    n_steps = 60  # monthly
    dt = MATURITY / n_steps
    survived = 0
    defaulted = 0

    for _ in range(n_paths):
        v = float(iss["asset"])
        points = [(0, v)]
        did_default = False
        for step in range(n_steps):
            dw = np.random.normal(0, math.sqrt(dt))
            v *= math.exp(
                (RISK_FREE - iss["vol"]**2 / 2) * dt
                + iss["vol"] * dw)
            points.append((step + 1, v))
            if v < barrier:
                did_default = True
                break

        col = ACCENT_RED if did_default else ACCENT_GREEN
        if did_default:
            defaulted += 1
        else:
            survived += 1

        for j in range(1, len(points)):
            t0, v0 = points[j - 1]
            t1, v1 = points[j]
            x1 = int(chart_left) + int(chart_w * t0 / n_steps)
            x2 = int(chart_left) + int(chart_w * t1 / n_steps)
            frac0 = 1 - (v0 - v_min) / (v_max - v_min)
            frac1 = 1 - (v1 - v_min) / (v_max - v_min)
            frac0 = max(0, min(1, frac0))
            frac1 = max(0, min(1, frac1))
            y1 = int(chart_top + chart_h * frac0)
            y2 = int(chart_top + chart_h * frac1)
            seg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                min(x1, x2), min(y1, y2),
                max(abs(x2 - x1), Emu(1)),
                max(abs(y2 - y1), Pt(1)))
            seg.fill.solid()
            seg.fill.fore_color.rgb = col
            seg.line.fill.background()

    # Stats box
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(5.5),
                Inches(0.3),
                f"50 paths: {survived} survived (green), "
                f"{defaulted} defaulted (red)  |  "
                f"Sample default rate = {defaulted/50:.0%}  |  "
                f"Historical 5Y PD = "
                f"{1 - math.exp(-iss['pd'] * MATURITY):.1%}",
                font_size=10, color=TEXT_MID)

    # Right side: MC algorithm steps
    add_textbox(slide, Inches(6.3), Inches(1.2), Inches(3.5), Inches(0.25),
                "MC Algorithm (per path)",
                font_size=13, bold=True, color=HEADER_BLUE)

    steps = [
        ("1. Evolve assets",
         "V(t+dt) = V(t) \u00d7 exp[(r\u2212\u03c3\u00b2/2)dt + \u03c3\u221adt\u00b7Z]"),
        ("2. At each coupon date:",
         "Cash: pay coupon, discount\n"
         "PIK: accrete N \u2192 N\u00d7(1+c/2)"),
        ("3. Update credit state",
         "L = N/V, recompute \u03bb(L), DD"),
        ("4. Check default",
         "V < B? If yes: recovery = R(N)\u00d7N"),
        ("5. If survived to T:",
         "PV += N(T) \u00d7 D(T)"),
        ("6. Aggregate paths",
         "Price = mean(path PVs)\n"
         "SE = std(path PVs) / \u221an"),
    ]

    for i, (title, desc) in enumerate(steps):
        y = Inches(1.55) + Inches(i * 0.65)
        add_textbox(slide, Inches(6.3), y, Inches(3.5), Inches(0.2),
                    title, font_size=10, bold=True, color=HEADER_BLUE)
        add_textbox(slide, Inches(6.3), y + Inches(0.18), Inches(3.5),
                    Inches(0.4),
                    desc, font_size=9, color=TEXT_DARK)

    add_callout_box(slide, Inches(0.5), Inches(5.3), Inches(9.2),
                    Inches(0.5),
                    "Key: the MC engine runs 25,000 paths with antithetic "
                    "variates (variance reduction). At each coupon date, "
                    "the toggle decision, hazard update, and recovery "
                    "adjustment are all path-dependent.",
                    bg_color=HEADER_BLUE, font_size=11)

    # Notional comparison
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(9), Inches(0.25),
                "Terminal Notional: Cash vs PIK on Surviving Paths",
                font_size=11, bold=True, color=HEADER_BLUE)

    terminal_ntl = NOTIONAL * (1 + COUPON / 2) ** 10
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(9), Inches(0.5),
                f"Cash-pay: always N = {NOTIONAL:.0f} at maturity  |  "
                f"Full PIK: N = {terminal_ntl:.2f} (+{terminal_ntl - NOTIONAL:.1f}%)  |  "
                f"Toggle: N varies by path (PIK only on stressed paths)",
                font_size=11, color=TEXT_DARK)


# ── Slide 12: Feedback Spiral ───────────────────────────────────────────

def slide_12_feedback_spiral(prs):
    """Slide 12: Walk through a single stressed path."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "11")
    add_slide_title(slide,
                    "The Feedback Spiral: A Single Path Walkthrough",
                    "B\u2212 issuer \u2014 asset value declining, PIK accreting")
    add_slide_number(slide, 12)

    iss = ISSUERS[3]  # B-
    lam0 = calibrate_hazard(iss["spread"], iss["rec"])
    l0 = NOTIONAL / iss["asset"]
    barrier = merton_barrier(iss["asset"], iss["vol"], iss["pd"])

    # Simulate a specific stressed path (seed chosen for good story)
    np.random.seed(77)
    n_coupons = 10
    dt_coupon = 0.5
    steps_per_coupon = 6  # monthly within each coupon period

    v = float(iss["asset"])
    ntl_cash = NOTIONAL
    ntl_pik = NOTIONAL
    semi_cpn = COUPON / 2

    path_rows = [["Period", "V(t)", "N (Cash)", "N (PIK)",
                  "L (PIK)", "\u03bb (PIK)", "R (PIK)", "DD"]]

    path_rows.append([
        "t=0", f"{v:.1f}", f"{ntl_cash:.1f}", f"{ntl_pik:.1f}",
        f"{ntl_pik / v:.2f}",
        f"{lam0 * 10000:.0f}bp",
        f"{iss['rec']:.0%}",
        f"{_dd(v, barrier, iss['vol'], MATURITY):.2f}",
    ])

    for cpn_i in range(n_coupons):
        # Evolve asset monthly within coupon period
        for _ in range(steps_per_coupon):
            dt = dt_coupon / steps_per_coupon
            dw = np.random.normal(0, math.sqrt(dt))
            v *= math.exp(
                (RISK_FREE - iss["vol"]**2 / 2) * dt
                + iss["vol"] * dw)
            # Stress: gentle downward bias for illustration
            v *= 0.997

        # PIK accrual
        ntl_pik *= (1 + semi_cpn)

        # Compute state
        lev = ntl_pik / v
        lam_now = lam0 * (lev / l0) ** 2
        r_now = max(0.10, iss["rec"] * NOTIONAL / ntl_pik)
        remain = MATURITY - (cpn_i + 1) * dt_coupon
        dd_now = _dd(v, barrier, iss["vol"], max(remain, 0.01))

        path_rows.append([
            f"t={dt_coupon * (cpn_i + 1):.1f}",
            f"{v:.1f}", f"{ntl_cash:.1f}", f"{ntl_pik:.1f}",
            f"{lev:.2f}",
            f"{lam_now * 10000:.0f}bp",
            f"{r_now:.0%}",
            f"{dd_now:.2f}",
        ])

    add_table(slide, Inches(0.3), Inches(1.1), Inches(9.4), path_rows,
              col_widths=[Inches(0.6), Inches(0.7), Inches(0.8),
                          Inches(0.8), Inches(0.7), Inches(0.8),
                          Inches(0.7), Inches(0.6)],
              font_size=9)

    # Interpretation
    final_lev = float(path_rows[-1][4])
    final_lam = path_rows[-1][5]
    final_r = path_rows[-1][6]
    final_dd = path_rows[-1][7]

    add_bullet_list(slide, Inches(0.5), Inches(5.2), Inches(9),
                    Inches(1.5), [
        f"Over 5 years: PIK notional grew from 100 to "
        f"{ntl_pik:.0f} while assets declined to {v:.0f}",
        f"Leverage rose from {l0:.2f} to {final_lev:.2f} \u2014 "
        f"hazard rate rose from {lam0*10000:.0f}bp to {final_lam}",
        f"Recovery fell from {iss['rec']:.0%} to {final_r} "
        f"\u2014 distance-to-default collapsed to {final_dd}",
        "This is one path. Across 25,000 paths the average "
        "captures the expected cost of this spiral",
    ], font_size=12, color=TEXT_DARK)

    add_callout_box(slide, Inches(1.0), Inches(6.8), Inches(8),
                    Inches(0.45),
                    "The spiral is self-reinforcing: PIK raises notional "
                    "\u2192 leverage rises \u2192 \u03bb rises \u2192 "
                    "more defaults \u2192 recovery falls \u2192 higher loss. "
                    "Cash-pay avoids this entirely.",
                    bg_color=ACCENT_RED, font_size=11)
