#!/usr/bin/env python3
"""Generate a PowerPoint deck on PIK coupon breakeven modelling.

Audience: credit team familiar with coupon types but low on modelling knowledge.
Approach: "Building Blocks" — start simple, add complexity, show results.

Usage:
    python generate_pik_presentation.py
    # => writes pik_coupon_pricing.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ────────────────────────────────────────────────────────

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
MID_BG = RGBColor(0x16, 0x21, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_BLUE = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT_ORANGE = RGBColor(0xE8, 0x8D, 0x3F)
ACCENT_RED = RGBColor(0xE0, 0x4F, 0x4F)
ACCENT_GREEN = RGBColor(0x5C, 0xB8, 0x5C)
BODY_BG = RGBColor(0xF5, 0xF5, 0xFA)
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x3D)
TEXT_MID = RGBColor(0x55, 0x55, 0x70)
HEADER_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
TABLE_HEADER_BG = RGBColor(0x2C, 0x3E, 0x6B)
TABLE_ALT_BG = RGBColor(0xE8, 0xEB, 0xF5)
TABLE_WHITE_BG = RGBColor(0xFF, 0xFF, 0xFF)

TOTAL_SLIDES = 18


# ── Helpers ───────────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text: str, *,
                font_size=18, bold=False, color=TEXT_DARK,
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items: list[str], *,
                    font_size=16, color=TEXT_DARK, spacing_after=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.level = 0
        p.space_after = spacing_after
    return txBox


def add_callout_box(slide, left, top, width, height, text: str, *,
                    bg_color=ACCENT_BLUE, text_color=WHITE, font_size=14):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = "Calibri"
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_section_number(slide, number: str, left=Inches(0.5), top=Inches(0.3)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_table(slide, left, top, width, rows_data: list[list[str]], *,
              col_widths=None, font_size=11, header_row=True):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, Inches(0.35 * n_rows))
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0 and header_row:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = TEXT_DARK
                    paragraph.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            if r == 0 and header_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_WHITE_BG
    return tbl_shape


def add_slide_number(slide, num: int):
    add_textbox(slide, Inches(8.8), Inches(7.1), Inches(1.2), Inches(0.3),
                f"{num}/{TOTAL_SLIDES}", font_size=10, color=TEXT_MID,
                alignment=PP_ALIGN.RIGHT)


def add_slide_title(slide, title: str, subtitle: str = ""):
    add_textbox(slide, Inches(1.1), Inches(0.2), Inches(8), Inches(0.5),
                title, font_size=26, bold=True, color=HEADER_BLUE)
    if subtitle:
        add_textbox(slide, Inches(1.1), Inches(0.65), Inches(8), Inches(0.4),
                    subtitle, font_size=14, color=TEXT_MID)


# ── Slide builders ────────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(1), Inches(2), Inches(8), Inches(1.2),
                "PIK Coupon Pricing", font_size=40, bold=True, color=WHITE)
    add_textbox(slide, Inches(1), Inches(3.1), Inches(8), Inches(0.8),
                "How Much Extra Spread Is Enough?", font_size=24, color=ACCENT_BLUE)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.6),
                "Modelling breakeven Z-spreads for Cash, PIK, and Toggle coupons\n"
                "across issuer credit quality",
                font_size=16, color=LIGHT_GREY)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.95), Inches(3), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()


def slide_02_the_question(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "1")
    add_slide_title(slide, "The Question")
    add_slide_number(slide, 2)

    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(8.5), Inches(0.7),
                "At what spread should a PIK bond trade versus a\n"
                "cash-pay bond for the same issuer?",
                font_size=20, bold=True, color=HEADER_BLUE,
                alignment=PP_ALIGN.CENTER)

    add_callout_box(slide, Inches(0.5), Inches(2.2), Inches(4.2), Inches(0.45),
                    "Naive View", bg_color=TEXT_MID)
    add_bullet_list(slide, Inches(0.5), Inches(2.8), Inches(4.2), Inches(2.5), [
        '"PIK just defers coupons \u2014 add 50bp"',
        '"The coupon rate is the same either way"',
        '"Compounding offsets the deferral"',
    ], font_size=14, color=TEXT_DARK)

    add_callout_box(slide, Inches(5.3), Inches(2.2), Inches(4.2), Inches(0.45),
                    "What the Models Show", bg_color=ACCENT_BLUE)
    add_bullet_list(slide, Inches(5.3), Inches(2.8), Inches(4.2), Inches(2.5), [
        "For strong credits (BB+): PIK premium is small (\u00b160bp) "
        "and sign is model-dependent",
        "For stressed credits (CCC): PIK premium > +260bp, "
        "robust across models",
        "The non-linearity means flat spread bumps are dangerous",
    ], font_size=14, color=TEXT_DARK)

    add_callout_box(slide, Inches(1.5), Inches(5.3), Inches(7), Inches(0.7),
                    "The answer depends on credit quality in a way that is impossible "
                    "to guess without a model. The premium is non-linear and, for "
                    "strong credits, depends on calibration assumptions.",
                    bg_color=ACCENT_ORANGE, font_size=14)


def slide_03_coupon_recap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "1")
    add_slide_title(slide, "Coupon Structures: Cash, PIK, and Toggle",
                    "What the investor receives and when")
    add_slide_number(slide, 3)

    # ── Cash-Pay ──────────────────────────────────────────────────────
    y_cash = Inches(1.3)
    add_textbox(slide, Inches(0.5), y_cash, Inches(2), Inches(0.3),
                "Cash-Pay", font_size=16, bold=True, color=ACCENT_GREEN)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(2.5), y_cash + Pt(8), Inches(6.5), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GREEN
    bar.line.fill.background()
    add_textbox(slide, Inches(2.5), y_cash + Inches(0.3), Inches(6.5), Inches(0.4),
                "4.25    4.25    4.25    4.25    4.25    4.25    4.25    4.25"
                "    4.25   104.25",
                font_size=10, color=TEXT_MID, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.5), y_cash + Inches(0.55), Inches(6.5), Inches(0.3),
                "Regular coupon payments + par at maturity. "
                "Cash flows spread across time.",
                font_size=12, color=TEXT_DARK)

    # ── Full PIK ──────────────────────────────────────────────────────
    y_pik = Inches(2.5)
    add_textbox(slide, Inches(0.5), y_pik, Inches(2), Inches(0.3),
                "Full PIK", font_size=16, bold=True, color=ACCENT_RED)
    bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(2.5), y_pik + Pt(8), Inches(6.5), Pt(4))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = ACCENT_RED
    bar2.line.fill.background()
    add_textbox(slide, Inches(2.5), y_pik + Inches(0.3), Inches(6.5), Inches(0.4),
                "  0         0         0         0         0         0"
                "         0         0         0       151.76",
                font_size=10, color=TEXT_MID, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.5), y_pik + Inches(0.55), Inches(6.5), Inches(0.3),
                "No coupons. Entire value at maturity: "
                "N \u00d7 (1 + c/2)\u00b9\u2070 = 151.76. Concentrated exposure.",
                font_size=12, color=TEXT_DARK)

    # ── PIK Toggle (detailed) ────────────────────────────────────────
    # Example: cash for periods 1-4, PIK for 5-6 (credit deteriorates),
    # back to cash for 7-10.  PIK accrual inflates notional:
    #   100 -> 104.25 -> 108.68
    # So post-PIK cash coupons are 108.68 x 4.25% = 4.62 (not 4.25).
    y_tog = Inches(3.7)
    add_textbox(slide, Inches(0.5), y_tog, Inches(2), Inches(0.3),
                "PIK Toggle", font_size=16, bold=True, color=ACCENT_ORANGE)
    bar3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(2.5), y_tog + Pt(8), Inches(6.5), Pt(4))
    bar3.fill.solid()
    bar3.fill.fore_color.rgb = ACCENT_ORANGE
    bar3.line.fill.background()

    # Toggle cashflow + notional table
    toggle_rows = [
        ["Period", "1", "2", "3", "4", "5", "6", "7", "8", "9", "10"],
        ["Mode", "Cash", "Cash", "Cash", "Cash",
         "PIK", "PIK", "Cash", "Cash", "Cash", "Cash"],
        ["Cashflow", "4.25", "4.25", "4.25", "4.25",
         "0", "0", "4.62", "4.62", "4.62", "113.30"],
        ["Notional", "100", "100", "100", "100",
         "104.25", "108.68", "108.68", "108.68", "108.68", "108.68"],
    ]
    tbl_shape = slide.shapes.add_table(
        4, 11, Inches(0.5), y_tog + Inches(0.35), Inches(9), Inches(1.2))
    table = tbl_shape.table
    # Set column widths: narrow first col, equal others
    table.columns[0].width = Inches(0.8)
    for ci in range(1, 11):
        table.columns[ci].width = Inches(0.82)

    for r, row in enumerate(toggle_rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(9)
                para.font.name = "Calibri"
                para.alignment = PP_ALIGN.CENTER
                if r == 0:
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                elif r == 1:
                    # Color the mode row: green for cash, red for PIK
                    if val == "PIK":
                        para.font.color.rgb = ACCENT_RED
                        para.font.bold = True
                    else:
                        para.font.color.rgb = ACCENT_GREEN
                        para.font.bold = True
                elif r == 2:
                    # Highlight the larger post-PIK coupons
                    if c >= 7 and c <= 10:
                        para.font.bold = True
                        para.font.color.rgb = ACCENT_ORANGE
                    elif c >= 5 and c <= 6:
                        para.font.color.rgb = ACCENT_RED
                    else:
                        para.font.color.rgb = TEXT_DARK
                elif r == 3:
                    # Highlight growing notional during PIK periods
                    if c >= 5:
                        para.font.bold = True
                        para.font.color.rgb = (ACCENT_RED if c <= 6
                                               else ACCENT_ORANGE)
                    else:
                        para.font.color.rgb = TEXT_MID
                if c == 0:
                    para.font.bold = True
                    para.alignment = PP_ALIGN.LEFT

            # Cell backgrounds
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif c >= 5 and c <= 6 and r > 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_WHITE_BG

    # Annotation below table
    add_textbox(slide, Inches(0.5), y_tog + Inches(1.6), Inches(9), Inches(0.6),
                "PIK periods (5\u20136): notional compounds 100 \u2192 104.25 "
                "\u2192 108.68.  When cash resumes (period 7+), coupons are "
                "4.62 not 4.25 \u2014 computed on the inflated notional.  "
                "At maturity the investor receives 108.68 + 4.62 = 113.30.",
                font_size=12, color=TEXT_DARK)

    add_callout_box(slide, Inches(1.0), Inches(6.0), Inches(8), Inches(0.7),
                    "PIK periods permanently ratchet the notional upward. "
                    "Later cash coupons are larger, but so is the principal "
                    "at risk in default \u2014 feeding the leverage spiral.",
                    bg_color=HEADER_BLUE, font_size=13)


def slide_04_hazard_rate(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "2")
    add_slide_title(slide, "The Hazard Rate: Measuring Default Risk",
                    "The simplest credit model")
    add_slide_number(slide, 4)

    add_bullet_list(slide, Inches(0.5), Inches(1.2), Inches(5), Inches(1.5), [
        "The hazard rate \u03bb is the instantaneous rate of default",
        "Think of it as: 'each year, there is a \u03bb chance of default'",
        "Survival probability decays over time: S(t) = e^(\u2212\u03bb \u00d7 t)",
    ], font_size=15, color=TEXT_DARK)

    add_textbox(slide, Inches(0.5), Inches(3.0), Inches(3), Inches(0.3),
                "Survival Probability S(t)", font_size=14, bold=True, color=HEADER_BLUE)
    rows = [
        ["Hazard Rate \u03bb", "1Y", "3Y", "5Y", "7Y"],
        ["143bp (BB+)", "98.6%", "95.8%", "93.1%", "90.5%"],
        ["591bp (B)", "94.3%", "83.8%", "74.4%", "66.2%"],
        ["1468bp (CCC)", "86.4%", "64.5%", "48.2%", "36.0%"],
    ]
    add_table(slide, Inches(0.5), Inches(3.4), Inches(5.5), rows, font_size=12)

    add_callout_box(slide, Inches(6.2), Inches(1.5), Inches(3.3), Inches(2.5),
                    "A CCC issuer at \u03bb=1468bp has only a 48% chance of surviving "
                    "to year 5.\n\nEvery cashflow scheduled beyond that point has "
                    "\u226448% probability of being received.",
                    bg_color=HEADER_BLUE, font_size=12)

    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.8),
                "The hazard rate is calibrated from the issuer's observed market "
                "spread. Higher \u03bb \u2192 faster survival decay \u2192 lower present value "
                "of future cashflows.",
                font_size=14, color=TEXT_DARK)


def slide_05_pricing_with_hazard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "2")
    add_slide_title(slide, "Pricing Bonds with Hazard Rates",
                    "Survival-weighted present value")
    add_slide_number(slide, 5)

    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.5),
                "PV  =  \u03a3  Cashflow(t)  \u00d7  S(t)  \u00d7  D(t)    +    Recovery Leg",
                font_size=20, bold=True, color=HEADER_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(9), Inches(0.3),
                "where S(t) = survival probability, D(t) = risk-free discount factor",
                font_size=13, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, Inches(0.5), Inches(2.3), Inches(9), Inches(2), [
        "Each future cashflow is weighted by the probability it gets paid "
        "(survival) and time value (discount)",
        "Later cashflows are penalised more: both S(t) and D(t) are smaller at "
        "longer horizons",
        "The recovery leg adds value for the default scenario: "
        "R \u00d7 par \u00d7 default_prob(t\u2081, t\u2082) \u00d7 D(t)",
    ], font_size=14, color=TEXT_DARK)

    add_callout_box(slide, Inches(0.5), Inches(4.2), Inches(4.2), Inches(0.4),
                    "Cash-Pay Bond", bg_color=ACCENT_GREEN, font_size=13)
    add_bullet_list(slide, Inches(0.5), Inches(4.7), Inches(4.2), Inches(1.8), [
        "10 semi-annual cashflows of 4.25",
        "Plus 100 at maturity",
        "Survival exposure: spread across 10 dates",
        "Early coupons have high S(t) \u2248 1.0",
    ], font_size=12, color=TEXT_DARK, spacing_after=Pt(4))

    add_callout_box(slide, Inches(5.3), Inches(4.2), Inches(4.2), Inches(0.4),
                    "Full PIK Bond", bg_color=ACCENT_RED, font_size=13)
    add_bullet_list(slide, Inches(5.3), Inches(4.7), Inches(4.2), Inches(1.8), [
        "Zero interim cashflows",
        "Single payment of 151.76 at maturity",
        "Survival exposure: 100% at year 5",
        "Entirely dependent on S(5Y)",
    ], font_size=12, color=TEXT_DARK, spacing_after=Pt(4))

    add_callout_box(slide, Inches(1.5), Inches(6.5), Inches(7), Inches(0.5),
                    "PIK concentrates all value at the worst survival point with the "
                    "largest notional.",
                    bg_color=ACCENT_ORANGE, font_size=13)


def slide_06_pik_in_hazard_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "2")
    add_slide_title(slide, "PIK in the Hazard Rate Model",
                    "Risk-neutral pricing: hazard rates from market spreads")
    add_slide_number(slide, 6)

    rows = [
        ["Issuer", "\u03bb (bp)", "Cash Price", "PIK Price", "\u0394 Price"],
        ["BB+ (Solid HY)", "143", "113.35", "110.89", "\u22122.46"],
        ["BB\u2212 (Mid HY)", "334", "107.56", "103.76", "\u22123.80"],
        ["B (Weak HY)", "591", "99.76", "94.44", "\u22125.32"],
        ["B\u2212 (Stressed)", "911", "90.33", "83.45", "\u22126.88"],
        ["CCC (Deeply Stressed)", "1468", "76.13", "67.34", "\u22128.79"],
    ]
    add_table(slide, Inches(0.8), Inches(1.4), Inches(8.4), rows,
              col_widths=[Inches(2.5), Inches(1.2), Inches(1.5),
                          Inches(1.5), Inches(1.5)],
              font_size=13)

    add_bullet_list(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(1.5), [
        "Under risk-neutral hazard rates (calibrated from market spreads), "
        "PIK always prices below cash \u2014 concentrated maturity exposure "
        "penalises PIK at every credit level",
        "The gap widens with credit risk: from \u22122.5 pts (BB+) to \u22128.8 "
        "pts (CCC)",
        "Note: market spreads embed a credit risk premium above expected "
        "losses. For BB+ the spread (85bp) implies ~7% 5Y default prob, "
        "versus the historical rate of ~1%",
    ], font_size=14, color=TEXT_DARK)

    add_callout_box(slide, Inches(1.5), Inches(6.0), Inches(7), Inches(0.7),
                    "Under risk-neutral pricing, PIK always costs the investor. "
                    "The question is whether the structural model \u2014 which uses "
                    "historical PDs \u2014 tells a different story.",
                    bg_color=HEADER_BLUE, font_size=12)


def slide_07_what_simple_misses(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "2")
    add_slide_title(slide, "The Calibration Gap",
                    "Why two models disagree on PIK for strong credits")
    add_slide_number(slide, 7)

    rows = [
        ["Issuer", "Mkt Sprd", "HR 5Y PD", "Hist 5Y PD", "Ratio"],
        ["BB+ (Solid HY)", "85bp", "6.9%", "1.0%", "6.9\u00d7"],
        ["BB\u2212 (Mid HY)", "210bp", "15.5%", "4.9%", "3.2\u00d7"],
        ["B (Weak HY)", "390bp", "25.6%", "11.8%", "2.2\u00d7"],
        ["B\u2212 (Stressed)", "630bp", "36.5%", "24.0%", "1.5\u00d7"],
        ["CCC (Deeply Stressed)", "1050bp", "52.0%", "39.3%", "1.3\u00d7"],
    ]
    add_table(slide, Inches(0.8), Inches(1.3), Inches(8.4), rows,
              col_widths=[Inches(2.5), Inches(1.2), Inches(1.2),
                          Inches(1.4), Inches(1.0)],
              font_size=12)

    add_textbox(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(0.5),
                "HR 5Y PD = implied from market spread.  "
                "Hist 5Y PD = from historical annual default rate.\n"
                "Ratio = how much the market \u2018overstates\u2019 default risk "
                "vs historical experience.",
                font_size=12, color=TEXT_MID)

    add_bullet_list(slide, Inches(0.5), Inches(4.4), Inches(9), Inches(1.5), [
        "Market spreads embed a risk premium \u2014 compensation for bearing "
        "credit risk beyond expected losses. For BB+, the spread implies "
        "7\u00d7the historical default rate",
        "The HR model uses market-implied hazard \u2192 PIK always costs. "
        "The MC model uses historical PDs for barriers \u2192 few defaults "
        "\u2192 PIK compounding survives \u2192 PIK appears to benefit investors",
        "The gap shrinks for weaker credits (1.3\u00d7 for CCC) \u2014 "
        "the risk premium is a smaller fraction of the total spread",
    ], font_size=14, color=TEXT_DARK)

    add_callout_box(slide, Inches(1.0), Inches(6.2), Inches(8), Inches(0.6),
                    "Whether PIK benefits investors for strong credits depends on "
                    "which default measure you use. Both models agree the effect "
                    "is small (\u00b160bp) for strong credits and large (100\u2013260bp) "
                    "for weak ones.",
                    bg_color=ACCENT_ORANGE, font_size=12)


def slide_08_merton_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "3")
    add_slide_title(slide, "The Merton Model: Default as a Balance Sheet Event",
                    "From fixed hazard rates to structural credit")
    add_slide_number(slide, 8)

    add_bullet_list(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.5), [
        "Firm has assets that fluctuate randomly (geometric Brownian motion)",
        "Firm has a debt barrier \u2014 if assets fall below it, the firm defaults",
        "The barrier is calibrated from historical default probabilities",
        "Higher asset vol \u2192 more chance of hitting the barrier \u2192 higher "
        "default probability",
    ], font_size=14, color=TEXT_DARK)

    add_callout_box(slide, Inches(6.5), Inches(1.3), Inches(3), Inches(0.35),
                    "Key Inputs", bg_color=TABLE_HEADER_BG, font_size=12)
    inputs = [
        ["Input", "Meaning"],
        ["Asset value (V)", "Market value of firm"],
        ["Asset vol (\u03c3)", "How volatile is V?"],
        ["Debt barrier (B)", "Default trigger level"],
        ["Annual PD", "Calibrates barrier"],
    ]
    add_table(slide, Inches(6.5), Inches(1.75), Inches(3), inputs, font_size=10)

    add_textbox(slide, Inches(0.5), Inches(3.9), Inches(5), Inches(0.3),
                "Our Five Test Issuers:", font_size=14, bold=True,
                color=HEADER_BLUE)

    profile_rows = [
        ["Issuer", "Assets", "Vol", "Barrier", "Ann PD", "DD", "PD(5Y)"],
        ["BB+ (Solid HY)", "200", "20%", "80.0", "0.20%", "2.33", "1.0%"],
        ["BB\u2212 (Mid HY)", "165", "25%", "70.0", "1.00%", "1.66", "4.9%"],
        ["B (Weak HY)", "140", "30%", "63.1", "2.50%", "1.19", "11.8%"],
        ["B\u2212 (Stressed)", "125", "35%", "66.4", "5.50%", "0.70", "24.0%"],
        ["CCC (Deeply Stressed)", "115", "40%", "75.8", "10.00%", "0.27", "39.4%"],
    ]
    add_table(slide, Inches(0.5), Inches(4.3), Inches(8.5), profile_rows,
              font_size=11)

    add_textbox(slide, Inches(0.5), Inches(6.4), Inches(9), Inches(0.5),
                "Barrier is calibrated via from_target_pd() to match each "
                "issuer\u2019s historical annual default rate.",
                font_size=12, color=TEXT_MID)


def slide_09_endogenous_hazard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "3")
    add_slide_title(slide, "Endogenous Hazard: PIK Raises Default Risk",
                    "Leverage-dependent hazard rate")
    add_slide_number(slide, 9)

    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5),
                "\u03bb(L) = \u03bb\u2080 \u00d7 (L / L\u2080)\u00b2",
                font_size=22, bold=True, color=HEADER_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(1.85), Inches(9), Inches(0.3),
                "Hazard rate at current leverage = base hazard \u00d7 "
                "(current leverage / base leverage)\u00b2",
                font_size=13, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, Inches(0.5), Inches(2.4), Inches(5), Inches(1.5), [
        "PIK accrual grows the debt (notional increases)",
        "Higher debt \u2192 higher leverage ratio L = B/V",
        "Higher leverage \u2192 higher hazard rate \u03bb(L)",
        "The exponent (\u00b2) makes it non-linear: doubling "
        "leverage quadruples hazard",
    ], font_size=14, color=TEXT_DARK)

    # B- issuer: base_leverage = 66.4/125 = 0.531, h0 = 9.11%
    haz_rows = [
        ["Leverage", "Hazard Rate", "vs Base"],
        ["0.35", "3.9%", "\u221257%"],
        ["0.40", "5.2%", "\u221243%"],
        ["0.53 (base)", "9.1%", "\u2014"],
        ["0.60", "11.6%", "+28%"],
        ["0.70", "15.8%", "+74%"],
        ["0.80", "20.7%", "+127%"],
        ["0.90", "26.2%", "+187%"],
    ]
    add_table(slide, Inches(5.5), Inches(2.4), Inches(4), haz_rows,
              font_size=11)

    add_textbox(slide, Inches(5.5), Inches(5.3), Inches(4), Inches(0.3),
                "Example: B\u2212 issuer (barrier 66.4, assets 125)",
                font_size=11, color=TEXT_MID)

    add_callout_box(slide, Inches(1), Inches(5.8), Inches(8), Inches(0.6),
                    "When PIK accrual pushes leverage from 0.53 to 0.90, the hazard "
                    "rate nearly triples. This is the core mechanism that makes PIK "
                    "dangerous for stressed credits.",
                    bg_color=ACCENT_ORANGE, font_size=12)


def slide_10_dynamic_recovery(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "3")
    add_slide_title(slide, "Dynamic Recovery: PIK Dilutes Recovery Value",
                    "More debt competing for the same assets")
    add_slide_number(slide, 10)

    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(0.5),
                "R(N) = max( floor,  R\u2080 \u00d7 N\u2080 / N )",
                font_size=22, bold=True, color=HEADER_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.5), Inches(1.85), Inches(9), Inches(0.3),
                "Recovery at current notional = base recovery \u00d7 "
                "(base notional / current notional), floored at 10%",
                font_size=13, color=TEXT_MID, alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, Inches(0.5), Inches(2.4), Inches(5), Inches(1.5), [
        "In default, creditors split the remaining assets",
        "PIK grows the total debt claim (notional rises)",
        "Same assets \u00f7 more debt = lower recovery per dollar",
        "Floor at 10% prevents recovery from reaching zero",
    ], font_size=14, color=TEXT_DARK)

    rec_rows = [
        ["Notional", "Recovery", "vs Base"],
        ["100 (base)", "30.0%", "\u2014"],
        ["112.5", "26.7%", "\u221211%"],
        ["125", "24.0%", "\u221220%"],
        ["137.5", "21.8%", "\u221227%"],
        ["151.8 (full PIK)", "19.8%", "\u221234%"],
    ]
    add_table(slide, Inches(5.5), Inches(2.4), Inches(4), rec_rows,
              font_size=11)

    add_textbox(slide, Inches(5.5), Inches(4.5), Inches(4), Inches(0.3),
                "Example: B\u2212 issuer (base recovery 30%)",
                font_size=11, color=TEXT_MID)

    add_callout_box(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.7),
                    "Double penalty: PIK simultaneously raises the probability of "
                    "default (via endogenous hazard) AND lowers the amount recovered "
                    "in default (via diluted recovery).",
                    bg_color=ACCENT_RED, font_size=12)


def slide_11_feedback_loop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "3")
    add_slide_title(slide, "The Feedback Loop",
                    "Why PIK risk is self-reinforcing")
    add_slide_number(slide, 11)

    positions = [
        (Inches(3.6), Inches(1.5), "PIK Accrual"),
        (Inches(6.5), Inches(2.2), "Higher\nNotional"),
        (Inches(7.0), Inches(3.8), "Higher\nLeverage"),
        (Inches(5.0), Inches(4.8), "Higher\nHazard Rate"),
        (Inches(2.5), Inches(3.8), "More Defaults\nBefore Maturity"),
        (Inches(1.5), Inches(2.2), "Lower\nRecovery"),
    ]
    colors = [ACCENT_ORANGE, ACCENT_ORANGE, ACCENT_RED, ACCENT_RED,
              ACCENT_RED, ACCENT_RED]

    for (x, y, label), color in zip(positions, colors):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.6), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(4)
        tf.margin_top = tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    arrow_data = [
        (Inches(5.15), Inches(1.6), "\u2192"),
        (Inches(7.5), Inches(2.9), "\u2193"),
        (Inches(6.4), Inches(4.7), "\u2190"),
        (Inches(3.8), Inches(5.0), "\u2190"),
        (Inches(2.0), Inches(3.0), "\u2191"),
        (Inches(3.0), Inches(1.5), "\u2192"),
    ]
    for x, y, arrow in arrow_data:
        add_textbox(slide, x, y, Inches(0.5), Inches(0.5),
                    arrow, font_size=20, bold=True, color=HEADER_BLUE,
                    alignment=PP_ALIGN.CENTER)

    add_callout_box(slide, Inches(0.5), Inches(6.0), Inches(9), Inches(0.7),
                    "This is why the PIK premium is non-linear: it is a "
                    "self-reinforcing spiral. Each cycle amplifies the next. "
                    "For stressed credits, the spiral dominates pricing.",
                    bg_color=HEADER_BLUE, font_size=13)


def slide_12_monte_carlo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "3")
    add_slide_title(slide, "Monte Carlo Simulation",
                    "Bringing the feedback loop to life")
    add_slide_number(slide, 12)

    add_textbox(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(0.6),
                "We simulate 25,000 paths of firm asset values over 5 years.\n"
                "At each time step, the model updates hazard and recovery based "
                "on current leverage.",
                font_size=15, color=TEXT_DARK)

    steps = [
        ("1", "Simulate asset paths",
         "Assets follow geometric Brownian motion with calibrated volatility"),
        ("2", "Check for default",
         "First-passage: has the asset path hit or crossed the debt barrier?"),
        ("3", "Update PIK notional",
         "If PIK coupon date: accrete the notional by the coupon amount"),
        ("4", "Recalculate hazard",
         "\u03bb(L) = \u03bb\u2080 \u00d7 (L/L\u2080)\u00b2  \u2014 leverage has "
         "changed due to PIK accrual"),
        ("5", "Recalculate recovery",
         "R(N) = R\u2080 \u00d7 N\u2080/N \u2014 more debt dilutes recovery"),
        ("6", "Discount and average",
         "PV across all paths \u2192 price, Z-spread, expected loss"),
    ]

    for i, (num, title, desc) in enumerate(steps):
        y = Inches(2.2) + Inches(i * 0.7)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.35), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Pt(0)
        tf.margin_top = tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_textbox(slide, Inches(1.0), y - Pt(2), Inches(3), Inches(0.35),
                    title, font_size=14, bold=True, color=TEXT_DARK)
        add_textbox(slide, Inches(4.0), y - Pt(2), Inches(5.5), Inches(0.35),
                    desc, font_size=12, color=TEXT_MID)

    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.4),
                "Steps 3\u20135 create the feedback loop: PIK \u2192 higher notional "
                "\u2192 recalculated hazard and recovery \u2192 different default/"
                "survival outcomes.",
                font_size=13, bold=True, color=HEADER_BLUE)


def slide_13_breakeven_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "4")
    add_slide_title(slide, "Structural Model: Breakeven Z-Spreads",
                    "MC results using historical PDs for barrier calibration")
    add_slide_number(slide, 13)

    rows = [
        ["Issuer", "Cash", "PIK", "Toggle", "PIK\u2212Cash", "Tog\u2212Cash"],
        ["BB+ (Solid HY)", "20bp", "\u221239bp", "22bp", "\u221259", "+1"],
        ["BB\u2212 (Mid HY)", "110bp", "88bp", "130bp", "\u221222", "+20"],
        ["B (Weak HY)", "292bp", "329bp", "346bp", "+37", "+54"],
        ["B\u2212 (Stressed)", "710bp", "851bp", "862bp", "+141", "+151"],
        ["CCC (Deeply Stressed)", "1,497bp", "1,759bp", "1,763bp", "+262", "+266"],
    ]
    add_table(slide, Inches(0.5), Inches(1.4), Inches(9), rows,
              col_widths=[Inches(2.5), Inches(1.1), Inches(1.1),
                          Inches(1.1), Inches(1.2), Inches(1.2)],
              font_size=12)

    add_textbox(slide, Inches(0.5), Inches(4.0), Inches(9), Inches(0.4),
                "Z-spreads are cash-equivalent.  MC barriers calibrated from "
                "historical PDs (BB+ = 0.2%/yr \u2192 ~2% 5Y default rate).",
                font_size=11, color=TEXT_MID)

    add_bullet_list(slide, Inches(0.5), Inches(4.5), Inches(9), Inches(2), [
        "BB+ to BB\u2212: the MC model shows PIK Z-spread below cash. "
        "Under historical default rates (1\u20135% over 5Y), PIK compounding "
        "survives on ~98% of paths \u2014 but recall this uses lower PDs "
        "than market spreads imply (slide 7)",
        "Crossover near B: the feedback loop starts to dominate "
        "(\u0394Z = +37bp). Both models agree PIK costs from here",
        "B\u2212 to CCC: PIK premium of +141bp to +262bp. The leverage "
        "spiral is unambiguous \u2014 this result is robust across "
        "calibration assumptions",
    ], font_size=14, color=TEXT_DARK)

    add_callout_box(slide, Inches(1.0), Inches(6.5), Inches(8), Inches(0.5),
                    "For weak credits (B and below): PIK premium is large "
                    "and robust. For strong credits: the sign depends on "
                    "whether you calibrate to market or historical defaults.",
                    bg_color=ACCENT_ORANGE, font_size=13)


def slide_14_credit_quality_sweep(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "4")
    add_slide_title(slide, "PIK Premium vs Credit Quality",
                    "The hockey-stick relationship")
    add_slide_number(slide, 14)

    rows = [
        ["Mkt Spread", "Cash Z", "PIK Z", "Premium"],
        ["50bp", "10bp", "\u221255bp", "\u221265bp"],
        ["100bp", "31bp", "\u221224bp", "\u221254bp"],
        ["200bp", "87bp", "56bp", "\u221231bp"],
        ["300bp", "168bp", "166bp", "\u22122bp"],
        ["400bp", "292bp", "329bp", "+37bp"],
        ["600bp", "552bp", "659bp", "+107bp"],
        ["850bp", "1,119bp", "1,335bp", "+216bp"],
        ["1,200bp", "1,702bp", "1,990bp", "+287bp"],
    ]
    add_table(slide, Inches(0.5), Inches(1.3), Inches(5.5), rows,
              col_widths=[Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2)],
              font_size=11)

    # Visual bar representation
    add_textbox(slide, Inches(6.3), Inches(1.3), Inches(3.2), Inches(0.3),
                "PIK Z-Spread Premium", font_size=13, bold=True,
                color=HEADER_BLUE)

    premiums = [
        ("50bp", -65), ("100bp", -54), ("200bp", -31), ("300bp", -2),
        ("400bp", 37), ("600bp", 107), ("850bp", 216), ("1200bp", 287),
    ]
    max_abs = 287

    for i, (label, prem) in enumerate(premiums):
        y = Inches(1.7) + Inches(i * 0.5)
        add_textbox(slide, Inches(6.1), y, Inches(0.6), Inches(0.3),
                    label, font_size=9, color=TEXT_MID)
        # zero line at a fixed x position
        zero_x = Inches(7.8)
        if prem >= 0:
            bar_w = Inches(1.5) * prem / max_abs
            color = ACCENT_RED if prem > 80 else (
                ACCENT_ORANGE if prem > 30 else ACCENT_GREEN)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(zero_x), y + Pt(2),
                int(bar_w), Pt(14))
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()
            add_textbox(slide, int(zero_x) + int(bar_w) + Pt(4),
                        y - Pt(1), Inches(0.8), Inches(0.3),
                        f"+{prem}", font_size=9, bold=True, color=color)
        else:
            bar_w = Inches(1.5) * abs(prem) / max_abs
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(zero_x) - int(bar_w),
                y + Pt(2), int(bar_w), Pt(14))
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT_GREEN
            bar.line.fill.background()
            add_textbox(slide, int(zero_x) - int(bar_w) - Inches(0.6),
                        y - Pt(1), Inches(0.6), Inches(0.3),
                        str(prem), font_size=9, bold=True,
                        color=ACCENT_GREEN, alignment=PP_ALIGN.RIGHT)

    # Zero line
    zero_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, int(Inches(7.8)), Inches(1.7),
        Pt(2), Inches(4.0))
    zero_line.fill.solid()
    zero_line.fill.fore_color.rgb = TEXT_MID
    zero_line.line.fill.background()

    # Zones
    add_callout_box(slide, Inches(0.5), Inches(5.5), Inches(2.8), Inches(0.5),
                    "LOW IMPACT ZONE", bg_color=ACCENT_GREEN,
                    font_size=11)
    add_textbox(slide, Inches(0.5), Inches(6.1), Inches(2.8), Inches(0.4),
                "Market spread < 300bp:\nPremium small, sign model-dep.",
                font_size=11, color=TEXT_DARK)

    add_callout_box(slide, Inches(3.5), Inches(5.5), Inches(2.8), Inches(0.5),
                    "CROSSOVER ZONE", bg_color=ACCENT_ORANGE, font_size=11)
    add_textbox(slide, Inches(3.5), Inches(6.1), Inches(2.8), Inches(0.4),
                "300\u2013500bp market spread:\nPIK premium 0\u201375bp",
                font_size=11, color=TEXT_DARK)

    add_callout_box(slide, Inches(6.5), Inches(5.5), Inches(2.8), Inches(0.5),
                    "STRUCTURAL RISK", bg_color=ACCENT_RED, font_size=11)
    add_textbox(slide, Inches(6.5), Inches(6.1), Inches(2.8), Inches(0.4),
                "Market spread > 500bp:\nPIK premium 75\u2013290+ bp",
                font_size=11, color=TEXT_DARK)


def slide_15_toggle_not_free(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "4")
    add_slide_title(slide, "The Toggle Is Not Free",
                    "B\u2212 issuer: comparing toggle strategies")
    add_slide_number(slide, 15)

    rows = [
        ["Strategy", "Price", "Z-Spread", "E[Loss]", "PIK %", "Term. Ntl"],
        ["Cash-Pay", "87.28", "710bp", "25.7%", "0%", "100.0"],
        ["Full PIK", "82.39", "851bp", "29.9%", "100%", "151.6"],
        ["Threshold (h > 10%)", "82.04", "862bp", "30.2%", "80%", "137.3"],
        ["Stochastic (sigmoid)", "82.23", "856bp", "30.0%", "73%", "132.9"],
        ["Optimal (nested MC)", "83.40", "821bp", "29.0%", "40%", "114.8"],
    ]
    add_table(slide, Inches(0.5), Inches(1.4), Inches(9), rows,
              col_widths=[Inches(2.2), Inches(1.0), Inches(1.2),
                          Inches(1.2), Inches(0.9), Inches(1.5)],
              font_size=11)

    add_bullet_list(slide, Inches(0.5), Inches(3.9), Inches(9), Inches(2), [
        "Threshold toggle Z-spread (862bp) exceeds full PIK (851bp) \u2014 "
        "adverse selection against the investor",
        "The borrower PIKs on 80% of coupons, concentrated on the worst "
        "paths \u2014 seeding the leverage spiral where it does the most damage",
        "Even the optimal exercise model (821bp) carries a +111bp premium "
        "over cash",
        "The toggle premium (+151bp over cash) is comparable to the full PIK "
        "premium (+141bp)",
    ], font_size=14, color=TEXT_DARK)

    add_callout_box(slide, Inches(1.5), Inches(6.0), Inches(7), Inches(0.7),
                    "A toggle option does not protect the investor. It gives the "
                    "borrower the ability to accelerate leverage precisely when "
                    "it does the most damage.",
                    bg_color=ACCENT_RED, font_size=12)


def slide_16_model_gap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "4")
    add_slide_title(slide, "Why the Models Disagree",
                    "Calibration assumptions drive the answer for strong credits")
    add_slide_number(slide, 16)

    rows = [
        ["Issuer", "HR \u0394Price", "MC \u0394Price", "HR 5Y PD",
         "MC DefRate", "Ratio"],
        ["BB+", "\u22122.5", "+3.0", "6.9%", "1.9%", "3.6\u00d7"],
        ["BB\u2212", "\u22123.8", "+1.1", "15.5%", "5.8%", "2.7\u00d7"],
        ["B", "\u22125.3", "\u22121.6", "25.6%", "13.2%", "1.9\u00d7"],
        ["B\u2212", "\u22126.9", "\u22124.9", "36.5%", "25.7%", "1.4\u00d7"],
        ["CCC", "\u22128.8", "\u22126.2", "52.0%", "42.3%", "1.2\u00d7"],
    ]
    add_table(slide, Inches(0.5), Inches(1.3), Inches(9), rows,
              col_widths=[Inches(2.0), Inches(1.1), Inches(1.1),
                          Inches(1.2), Inches(1.2), Inches(1.0)],
              font_size=11)

    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(9), Inches(0.3),
                "\u0394Price = PIK minus Cash (points). "
                "HR 5Y PD = market-implied. MC DefRate = simulation default rate.",
                font_size=11, color=TEXT_MID)

    add_bullet_list(slide, Inches(0.5), Inches(4.0), Inches(9), Inches(2), [
        "The HR model uses market-implied defaults (including risk "
        "premium) \u2192 PIK always costs. The MC model uses historical "
        "PDs for barriers \u2192 fewer defaults \u2192 PIK compounding survives",
        "For BB+: the HR model sees 3.6\u00d7 more defaults than the MC "
        "model. This gap drives the sign reversal in PIK pricing",
        "For CCC: both models see similar default rates (ratio 1.2\u00d7). "
        "The PIK penalty is robust regardless of calibration",
        "The practical implication: for strong credits, the PIK premium/"
        "discount is small either way (\u00b160bp). For weak credits, it's "
        "large (+100\u2013260bp) and model choice barely matters",
    ], font_size=13, color=TEXT_DARK)

    add_callout_box(slide, Inches(1.0), Inches(6.2), Inches(8), Inches(0.6),
                    "Model choice matters most where the answer matters least "
                    "(strong credits). Where the premium is large and actionable "
                    "(weak credits), both models agree.",
                    bg_color=HEADER_BLUE, font_size=12)


def slide_17_when_does_it_matter(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "5")
    add_slide_title(slide, "When Does the PIK Premium Matter?",
                    "A practical framework")
    add_slide_number(slide, 17)

    zones = [
        ("LOW IMPACT ZONE", ACCENT_GREEN, Inches(0.5),
         "BB+ to BB  |  Mkt spread < 300bp",
         ["PIK premium is small (\u00b160bp)",
          "Sign is calibration-dependent",
          "Feedback loop is negligible",
          "PIK vs cash difference is minor"]),
        ("CROSSOVER ZONE", ACCENT_ORANGE, Inches(3.6),
         "B  |  Mkt spread 300\u2013500bp",
         ["PIK premium 0\u201375bp",
          "Simple models understate risk",
          "Structural model adds value",
          "Scrutinise leverage trajectory"]),
        ("STRUCTURAL RISK", ACCENT_RED, Inches(6.7),
         "B\u2212 to CCC  |  Mkt spread > 500bp",
         ["PIK premium 100\u2013290+ bp",
          "Feedback loop dominates",
          "Structural model is essential",
          "Flat spread bump is dangerous"]),
    ]

    for title, color, x, desc, bullets in zones:
        add_callout_box(slide, x, Inches(1.3), Inches(2.8), Inches(0.5),
                        title, bg_color=color, font_size=13)
        add_textbox(slide, x, Inches(1.9), Inches(2.8), Inches(0.3),
                    desc, font_size=12, bold=True, color=TEXT_DARK)
        add_bullet_list(slide, x, Inches(2.3), Inches(2.8), Inches(1.5),
                        bullets, font_size=12, color=TEXT_DARK,
                        spacing_after=Pt(4))

    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.3),
                "Practical Decision Guide:", font_size=14, bold=True,
                color=HEADER_BLUE)

    guide_rows = [
        ["Question", "Low Impact", "Crossover", "Structural Risk"],
        ["PIK Z-spread premium", "\u00b160bp (model-dep.)", "0 to +75bp",
         "+100 to +290bp"],
        ["Model needed?", "Either model adequate",
         "Structural recommended", "Structural required"],
        ["Toggle vs full PIK", "Minimal difference",
         "Toggle may exceed PIK", "Toggle \u2248 PIK or worse"],
        ["Key risk factor", "Calibration assumption",
         "Leverage trajectory", "Feedback loop intensity"],
    ]
    add_table(slide, Inches(0.5), Inches(4.6), Inches(9), guide_rows,
              col_widths=[Inches(2.2), Inches(2.3), Inches(2.3), Inches(2.3)],
              font_size=11)


def slide_18_takeaways(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MID_BG)
    add_slide_number(slide, 18)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
                "Key Takeaways", font_size=32, bold=True, color=WHITE)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(2), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    takeaways = [
        ("1", "PIK premium is non-linear and credit-quality-dependent",
         "For strong credits (BB+/BB\u2212), the premium is small (\u00b160bp) "
         "and its sign depends on calibration. For stressed credits "
         "(CCC), the premium exceeds +260bp and is robust across models."),
        ("2", "Model disagreement is largest where the stakes are smallest",
         "For strong credits, models disagree on the sign but the magnitude "
         "is small. For weak credits, models agree: PIK premium is large "
         "and the feedback spiral dominates."),
        ("3", "The toggle option does not protect the investor",
         "Borrowers PIK when credit deteriorates \u2014 adverse selection. "
         "Toggle Z-spreads can exceed full PIK because the feedback loop "
         "is concentrated on the worst paths."),
        ("4", "Above ~400bp market spread, structural modelling is essential",
         "In the structural risk zone, the feedback loop dominates pricing "
         "regardless of calibration. A flat spread bump approach will "
         "materially misprice PIK risk."),
    ]

    for i, (num, title, detail) in enumerate(takeaways):
        y = Inches(1.5) + Inches(i * 1.3)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.45), Inches(0.45))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.margin_left = tf.margin_right = Pt(0)
        tf.margin_top = tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_textbox(slide, Inches(1.5), y - Pt(2), Inches(7.5), Inches(0.35),
                    title, font_size=17, bold=True, color=WHITE)
        add_textbox(slide, Inches(1.5), y + Inches(0.35), Inches(7.5),
                    Inches(0.6),
                    detail, font_size=13, color=LIGHT_GREY)


# ── Main ──────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_the_question(prs)
    slide_03_coupon_recap(prs)
    slide_04_hazard_rate(prs)
    slide_05_pricing_with_hazard(prs)
    slide_06_pik_in_hazard_model(prs)
    slide_07_what_simple_misses(prs)
    slide_08_merton_model(prs)
    slide_09_endogenous_hazard(prs)
    slide_10_dynamic_recovery(prs)
    slide_11_feedback_loop(prs)
    slide_12_monte_carlo(prs)
    slide_13_breakeven_table(prs)
    slide_14_credit_quality_sweep(prs)
    slide_15_toggle_not_free(prs)
    slide_16_model_gap(prs)
    slide_17_when_does_it_matter(prs)
    slide_18_takeaways(prs)

    out = Path(__file__).parent / "pik_coupon_pricing.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
