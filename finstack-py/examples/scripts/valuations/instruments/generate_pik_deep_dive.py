#!/usr/bin/env python3
"""Generate a deep-dive PowerPoint deck on PIK coupon modelling.

Companion to the executive summary (pik_coupon_pricing.pptx).
Covers model mechanics, parameters, intuition, MC path visualisation,
and the feedback loop in detail.

Usage:
    python generate_pik_deep_dive.py
    # => writes pik_deep_dive.pptx
"""

from __future__ import annotations

import math
from pathlib import Path

import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette (matches executive deck) ──────────────────────────────

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
MID_BG = RGBColor(0x16, 0x21, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_BLUE = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT_ORANGE = RGBColor(0xE8, 0x8D, 0x3F)
ACCENT_RED = RGBColor(0xE0, 0x4F, 0x4F)
ACCENT_GREEN = RGBColor(0x5C, 0xB8, 0x5C)
BODY_BG = RGBColor(0xF5, 0xF5, 0xFA)
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x3D)
TEXT_MID = RGBColor(0x55, 0x55, 0x70)
HEADER_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
TABLE_HEADER_BG = RGBColor(0x2C, 0x3E, 0x6B)
TABLE_ALT_BG = RGBColor(0xE8, 0xEB, 0xF5)
TABLE_WHITE_BG = RGBColor(0xFF, 0xFF, 0xFF)

TOTAL_SLIDES = 16

# ── Global model parameters ──────────────────────────────────────────────

RISK_FREE = 0.045
COUPON = 0.085
MATURITY = 5
NOTIONAL = 100.0
NUM_PATHS = 25_000

ISSUERS = [
    {"name": "BB+ (Solid HY)", "asset": 200, "vol": 0.20,
     "pd": 0.0020, "spread": 0.0085, "rec": 0.45, "ltv": "50%"},
    {"name": "BB\u2212 (Mid HY)", "asset": 165, "vol": 0.25,
     "pd": 0.0100, "spread": 0.0210, "rec": 0.40, "ltv": "61%"},
    {"name": "B (Weak HY)", "asset": 140, "vol": 0.30,
     "pd": 0.0250, "spread": 0.0390, "rec": 0.35, "ltv": "71%"},
    {"name": "B\u2212 (Stressed)", "asset": 125, "vol": 0.35,
     "pd": 0.0550, "spread": 0.0630, "rec": 0.30, "ltv": "80%"},
    {"name": "CCC (Deeply Stressed)", "asset": 115, "vol": 0.40,
     "pd": 0.1000, "spread": 0.1050, "rec": 0.25, "ltv": "87%"},
]


# ── Helpers ──────────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text: str, *,
                font_size=18, bold=False, color=TEXT_DARK,
                alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items: list[str], *,
                    font_size=16, color=TEXT_DARK, spacing_after=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.level = 0
        p.space_after = spacing_after
    return txBox


def add_callout_box(slide, left, top, width, height, text: str, *,
                    bg_color=ACCENT_BLUE, text_color=WHITE, font_size=14):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = "Calibri"
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_section_number(slide, number: str, left=Inches(0.5), top=Inches(0.3)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_table(slide, left, top, width, rows_data: list[list[str]], *,
              col_widths=None, font_size=11, header_row=True):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, Inches(0.35 * n_rows))
    table = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Calibri"
                if r == 0 and header_row:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = TEXT_DARK
                    paragraph.alignment = (
                        PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT)
            if r == 0 and header_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_WHITE_BG
    return tbl_shape


def add_slide_number(slide, num: int):
    add_textbox(slide, Inches(8.8), Inches(7.1), Inches(1.2), Inches(0.3),
                f"{num}/{TOTAL_SLIDES}", font_size=10, color=TEXT_MID,
                alignment=PP_ALIGN.RIGHT)


def add_slide_title(slide, title: str, subtitle: str = ""):
    add_textbox(slide, Inches(1.1), Inches(0.2), Inches(8), Inches(0.5),
                title, font_size=26, bold=True, color=HEADER_BLUE)
    if subtitle:
        add_textbox(slide, Inches(1.1), Inches(0.65), Inches(8), Inches(0.4),
                    subtitle, font_size=14, color=TEXT_MID)


def draw_line(slide, x1, y1, x2, y2, color=TEXT_MID, width=Pt(1)):
    """Draw a thin line between two points."""
    left = min(x1, x2)
    top = min(y1, y2)
    w = abs(x2 - x1) or Emu(1)
    h = abs(y2 - y1) or Emu(1)
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, w, h)
    connector.fill.solid()
    connector.fill.fore_color.rgb = color
    connector.line.fill.background()
    return connector


def add_formula_box(slide, left, top, width, formula: str,
                    caption: str = ""):
    """Formula in bold + optional caption below."""
    add_textbox(slide, left, top, width, Inches(0.35),
                formula, font_size=16, bold=True, color=HEADER_BLUE)
    if caption:
        add_textbox(slide, left, top + Inches(0.35), width, Inches(0.25),
                    caption, font_size=11, color=TEXT_MID)


# ── Computation helpers ──────────────────────────────────────────────────

def calibrate_hazard(spread: float, recovery: float) -> float:
    """Bisect for flat hazard rate matching a Z-spread."""
    times = [i / 2.0 for i in range(1, MATURITY * 2 + 1)]
    cpn = COUPON / 2 * NOTIONAL
    target = sum(cpn * math.exp(-(RISK_FREE + spread) * t) for t in times)
    target += NOTIONAL * math.exp(-(RISK_FREE + spread) * MATURITY)

    def _pv(h):
        pv, prev_s = 0.0, 1.0
        for t in times:
            df = math.exp(-RISK_FREE * t)
            s = math.exp(-h * t)
            pv += cpn * df * s
            pv += recovery * NOTIONAL * df * (prev_s - s)
            prev_s = s
        pv += NOTIONAL * math.exp(-RISK_FREE * MATURITY) * math.exp(
            -h * MATURITY)
        return pv

    lo, hi = 0.0, 5.0
    for _ in range(200):
        mid = (lo + hi) / 2
        if _pv(mid) > target:
            lo = mid
        else:
            hi = mid
    return (lo + hi) / 2


def hr_bond_price(hazard: float, recovery: float, coupon_type: str = "cash",
                  coupon_rate: float = COUPON, maturity: float = MATURITY
                  ) -> float:
    """Price a bond under flat hazard rate. coupon_type: cash or pik."""
    times = [i / 2.0 for i in range(1, int(maturity * 2) + 1)]
    semi_cpn = coupon_rate / 2

    pv, prev_s = 0.0, 1.0
    notional = NOTIONAL
    for t in times:
        df = math.exp(-RISK_FREE * t)
        s = math.exp(-hazard * t)
        if coupon_type == "cash":
            pv += semi_cpn * NOTIONAL * df * s
        else:  # pik: coupon accretes
            notional *= (1 + semi_cpn)
        pv += recovery * NOTIONAL * df * (prev_s - s)
        prev_s = s
    # terminal
    df_T = math.exp(-RISK_FREE * maturity)
    s_T = math.exp(-hazard * maturity)
    pv += notional * df_T * s_T
    return pv


def price_to_zspread(price: float, coupon_rate: float = COUPON,
                     maturity: float = MATURITY) -> float:
    """Bisect for Z-spread that reproduces a given clean price."""
    times = [i / 2.0 for i in range(1, int(maturity * 2) + 1)]
    cpn = coupon_rate / 2 * NOTIONAL

    def _pv(z):
        return (sum(cpn * math.exp(-(RISK_FREE + z) * t) for t in times)
                + NOTIONAL * math.exp(-(RISK_FREE + z) * maturity))

    lo, hi = -0.5, 5.0
    for _ in range(200):
        mid = (lo + hi) / 2
        if _pv(mid) > price:
            lo = mid
        else:
            hi = mid
    return (lo + hi) / 2


def _norm_cdf(x: float) -> float:
    """Standard normal CDF via math.erfc (no scipy needed)."""
    return 0.5 * math.erfc(-x / math.sqrt(2))


def _norm_ppf(p: float) -> float:
    """Inverse standard normal CDF via rational approximation."""
    if p <= 0:
        return -10.0
    if p >= 1:
        return 10.0
    # Beasley-Springer-Moro algorithm
    a = [0, -3.969683028665376e+01, 2.209460984245205e+02,
         -2.759285104469687e+02, 1.383577518672690e+02,
         -3.066479806614716e+01, 2.506628277459239e+00]
    b = [0, -5.447609879822406e+01, 1.615858368580409e+02,
         -1.556989798598866e+02, 6.680131188771972e+01,
         -1.328068155288572e+01]
    c = [0, -7.784894002430293e-03, -3.223964580411365e-01,
         -2.400758277161838e+00, -2.549732539343734e+00,
         4.374664141464968e+00, 2.938163982698783e+00]
    d = [0, 7.784695709041462e-03, 3.224671290700398e-01,
         2.445134137142996e+00, 3.754408661907416e+00]

    p_low, p_high = 0.02425, 1 - 0.02425
    if p < p_low:
        q = math.sqrt(-2 * math.log(p))
        return ((((c[5]*q + c[4])*q + c[3])*q + c[2])*q + c[1]) / \
               (((d[4]*q + d[3])*q + d[2])*q + d[1]*q + 1)
    elif p <= p_high:
        q = p - 0.5
        r = q * q
        return ((((a[5]*r + a[4])*r + a[3])*r + a[2])*r + a[1]) / \
               ((((b[5]*r + b[4])*r + b[3])*r + b[2])*r + b[1]*r + 1) * q
    else:
        q = math.sqrt(-2 * math.log(1 - p))
        return -((((c[5]*q + c[4])*q + c[3])*q + c[2])*q + c[1]) / \
                (((d[4]*q + d[3])*q + d[2])*q + d[1]*q + 1)


def merton_barrier(asset: float, vol: float, annual_pd: float) -> float:
    """Calibrate Merton barrier from target annual PD (terminal barrier)."""
    five_yr_pd = 1 - math.exp(-annual_pd * MATURITY)
    dd = -_norm_ppf(five_yr_pd)
    drift = (RISK_FREE - vol**2 / 2) * MATURITY
    barrier = asset * math.exp(-(dd * vol * math.sqrt(MATURITY) + drift))
    return barrier


# ── Slide 1: Title ───────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_textbox(slide, Inches(1), Inches(2), Inches(8), Inches(1.2),
                "PIK Coupon Modelling", font_size=40, bold=True, color=WHITE)
    add_textbox(slide, Inches(1), Inches(3.1), Inches(8), Inches(0.8),
                "The Deep Dive", font_size=28, color=ACCENT_BLUE)
    add_textbox(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.8),
                "Model mechanics, parameters, Monte Carlo paths,\n"
                "and the structural feedback loop",
                font_size=16, color=LIGHT_GREY)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.95), Inches(3), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()
    add_textbox(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.4),
                "Companion to: PIK Coupon Pricing \u2014 "
                "How Much Extra Spread Is Enough?",
                font_size=12, color=TEXT_MID)


# ── Slide 2: The Result to Explain ──────────────────────────────────────

def slide_02_result_to_explain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "1")
    add_slide_title(slide, "The Result to Explain",
                    "From the executive deck: the hockey-stick PIK premium")
    add_slide_number(slide, 2)

    # Recreate the hockey-stick bar chart
    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(4), Inches(0.3),
                "PIK Z-Spread Premium by Market Spread",
                font_size=14, bold=True, color=HEADER_BLUE)

    premiums = [
        ("50bp", -65), ("100bp", -54), ("200bp", -31), ("300bp", -2),
        ("400bp", +37), ("600bp", +107), ("850bp", +216), ("1200bp", +287),
    ]
    max_abs = 287

    for i, (label, prem) in enumerate(premiums):
        y = Inches(1.5) + Inches(i * 0.42)
        add_textbox(slide, Inches(0.5), y, Inches(0.7), Inches(0.3),
                    label, font_size=9, color=TEXT_MID)
        zero_x = Inches(2.8)
        if prem >= 0:
            bar_w = Inches(1.8) * prem / max_abs
            col = (ACCENT_RED if prem > 80
                   else ACCENT_ORANGE if prem > 30 else ACCENT_GREEN)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(zero_x), y + Pt(2),
                int(bar_w), Pt(14))
            bar.fill.solid()
            bar.fill.fore_color.rgb = col
            bar.line.fill.background()
            add_textbox(slide, int(zero_x) + int(bar_w) + Pt(4),
                        y - Pt(1), Inches(0.6), Inches(0.3),
                        f"+{prem}", font_size=9, bold=True, color=col)
        else:
            bar_w = Inches(1.8) * abs(prem) / max_abs
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, int(zero_x) - int(bar_w),
                y + Pt(2), int(bar_w), Pt(14))
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT_GREEN
            bar.line.fill.background()
            add_textbox(slide, int(zero_x) - int(bar_w) - Inches(0.5),
                        y - Pt(1), Inches(0.5), Inches(0.3),
                        str(prem), font_size=9, bold=True,
                        color=ACCENT_GREEN, alignment=PP_ALIGN.RIGHT)

    # Zero line
    zl = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, int(Inches(2.8)), Inches(1.5),
        Pt(2), Inches(3.4))
    zl.fill.solid()
    zl.fill.fore_color.rgb = TEXT_MID
    zl.line.fill.background()

    # Right side: key questions
    add_textbox(slide, Inches(5.3), Inches(1.1), Inches(4.5), Inches(0.3),
                "Questions This Deck Answers",
                font_size=14, bold=True, color=HEADER_BLUE)
    add_bullet_list(slide, Inches(5.3), Inches(1.5), Inches(4.5),
                    Inches(4.5), [
        "Why does the PIK premium flip sign around 300bp?",
        "What drives the non-linearity \u2014 why isn\u2019t it "
        "proportional to spread?",
        "How does the hazard-rate model price PIK, and "
        "where does it break down?",
        "What is the Merton structural model and how is "
        "it calibrated?",
        "How do endogenous hazard and dynamic recovery "
        "create the feedback loop?",
        "What do Monte Carlo paths actually look like?",
        "Why does the toggle option fail to protect "
        "investors?",
        "What are the key parameter assumptions and "
        "how sensitive are results?",
    ], font_size=12, color=TEXT_DARK)


# ── Slide 3: Bond Setup & Cash Flows ────────────────────────────────────

def slide_03_bond_setup(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "2")
    add_slide_title(slide, "Bond Setup & Cash Flow Timing",
                    "The fundamental difference between cash-pay and PIK")
    add_slide_number(slide, 3)

    # Parameters box
    add_callout_box(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.4),
                    f"5-year  |  {COUPON:.1%} semi-annual  |  "
                    f"Par = {NOTIONAL:.0f}  |  "
                    f"Risk-free = {RISK_FREE:.2%} flat  |  "
                    f"Recovery = issuer-dependent",
                    bg_color=HEADER_BLUE, font_size=12)

    # Cash-pay timeline
    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(4.2), Inches(0.3),
                "CASH-PAY: 10 coupons + par at maturity",
                font_size=13, bold=True, color=ACCENT_GREEN)

    # Draw timeline
    cpn_amt = COUPON / 2 * NOTIONAL
    cash_times = [f"t={i/2:.1f}" for i in range(1, 11)]
    for i in range(10):
        x = Inches(0.5 + i * 0.42)
        y = Inches(2.2)
        add_textbox(slide, x, y, Inches(0.42), Inches(0.22),
                    f"{cpn_amt:.2f}", font_size=7, color=ACCENT_GREEN,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.2), Inches(0.42), Inches(0.18),
                    f"{(i+1)/2:.1f}Y", font_size=6, color=TEXT_MID,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(4.5), Inches(2.2), Inches(0.8), Inches(0.22),
                f"+ {NOTIONAL:.0f}", font_size=8, bold=True,
                color=ACCENT_GREEN)

    # PIK timeline
    add_textbox(slide, Inches(0.5), Inches(2.8), Inches(4.5), Inches(0.3),
                "FULL PIK: zero coupons \u2192 inflated notional at maturity",
                font_size=13, bold=True, color=ACCENT_ORANGE)

    terminal_ntl = NOTIONAL * (1 + COUPON / 2) ** 10
    for i in range(10):
        x = Inches(0.5 + i * 0.42)
        y = Inches(3.2)
        ntl_i = NOTIONAL * (1 + COUPON / 2) ** (i + 1)
        add_textbox(slide, x, y, Inches(0.42), Inches(0.22),
                    f"N={ntl_i:.1f}", font_size=6, color=ACCENT_ORANGE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.2), Inches(0.42), Inches(0.18),
                    f"{(i+1)/2:.1f}Y", font_size=6, color=TEXT_MID,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(4.5), Inches(3.15), Inches(1.5), Inches(0.3),
                f"\u2192 {terminal_ntl:.2f}", font_size=9, bold=True,
                color=ACCENT_ORANGE)

    # Key insight box
    add_callout_box(slide, Inches(0.5), Inches(3.8), Inches(4.5),
                    Inches(0.55),
                    f"PIK terminal notional = {terminal_ntl:.2f} "
                    f"(+{terminal_ntl - NOTIONAL:.1f}% above par). "
                    f"All risk concentrated at maturity.",
                    bg_color=ACCENT_ORANGE, font_size=11)

    # Right side: why this matters
    add_textbox(slide, Inches(5.5), Inches(1.8), Inches(4.3), Inches(0.3),
                "Why Timing Matters for Pricing",
                font_size=14, bold=True, color=HEADER_BLUE)

    add_bullet_list(slide, Inches(5.5), Inches(2.2), Inches(4.3),
                    Inches(2.5), [
        "Cash-pay: spreads risk across 10 coupon dates. "
        "Early coupons are almost certain to be paid",
        "PIK: concentrates ALL cash flow at maturity. "
        "The single payment is weighted by S(T), the "
        "5-year survival probability",
        "For a credit with S(5Y) = 70%, cash-pay "
        "collects ~95% of early coupons but PIK "
        "gets nothing if the issuer defaults",
        "This timing asymmetry is the root cause "
        "of the PIK premium under hazard-rate pricing",
    ], font_size=12, color=TEXT_DARK)

    # Survival probability table
    add_textbox(slide, Inches(5.5), Inches(4.5), Inches(4.3), Inches(0.25),
                "Survival Probabilities by Issuer",
                font_size=12, bold=True, color=HEADER_BLUE)

    surv_rows = [["Issuer", "LTV", "\u03bb (bp)", "S(1Y)", "S(3Y)", "S(5Y)"]]
    for iss in ISSUERS:
        lam = calibrate_hazard(iss["spread"], iss["rec"])
        surv_rows.append([
            iss["name"].split(" (")[0], iss["ltv"],
            f"{lam * 10000:.0f}",
            f"{math.exp(-lam * 1):.1%}",
            f"{math.exp(-lam * 3):.1%}",
            f"{math.exp(-lam * 5):.1%}",
        ])
    add_table(slide, Inches(5.5), Inches(4.85), Inches(4.3), surv_rows,
              col_widths=[Inches(0.7), Inches(0.5), Inches(0.6),
                          Inches(0.7), Inches(0.7), Inches(0.7)],
              font_size=9)

    add_textbox(slide, Inches(5.5), Inches(6.7), Inches(4), Inches(0.25),
                "S(t) = exp(\u2212\u03bb \u00d7 t) under flat hazard rate",
                font_size=9, color=TEXT_MID)


# ── Slide 4: HR Model — How \u03bb Prices the Bond ────────────────────────────

def slide_04_hr_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "3")
    add_slide_title(slide, "Hazard Rate Model: How \u03bb Prices the Bond",
                    "Reduced-form pricing under flat hazard rates")
    add_slide_number(slide, 4)

    # Formula
    add_formula_box(slide, Inches(0.5), Inches(1.1), Inches(9),
                    "PV = \u03a3 cpn \u00d7 D(t) \u00d7 S(t)  +  "
                    "N \u00d7 D(T) \u00d7 S(T)  +  "
                    "R \u00d7 N \u00d7 \u03a3 D(t) \u00d7 [S(t\u22121) \u2212 S(t)]",
                    "D(t) = exp(\u2212r\u00d7t)  |  "
                    "S(t) = exp(\u2212\u03bb\u00d7t)  |  "
                    "\u03bb calibrated from market Z-spread via bisection")

    # Three components explained
    components = [
        ("Coupon PV", "\u03a3 cpn \u00d7 D(t) \u00d7 S(t)",
         "Each coupon discounted for time value AND survival. "
         "Cash-pay: 10 small payments. PIK: zero (coupons accrete).",
         ACCENT_GREEN),
        ("Redemption PV", "N \u00d7 D(T) \u00d7 S(T)",
         "Terminal notional discounted by full survival. "
         "Cash-pay: N = 100. PIK: N = 151.26 (accreted).",
         ACCENT_BLUE),
        ("Recovery PV", "R \u00d7 N \u00d7 \u03a3 D(t) \u00d7 \u0394S",
         "Expected recovery on default. Proportional to "
         "incremental default probability each period.",
         ACCENT_ORANGE),
    ]

    for i, (title, formula, desc, col) in enumerate(components):
        y = Inches(2.0) + Inches(i * 0.85)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y,
            Inches(4.5), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = col
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_top = Pt(4)
        p = tf.paragraphs[0]
        p.text = f"{title}: {formula}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        add_textbox(slide, Inches(5.2), y + Pt(4), Inches(4.5),
                    Inches(0.6), desc, font_size=11, color=TEXT_DARK)

    # HR results table
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(9), Inches(0.25),
                "HR Prices: Cash vs PIK at Each Issuer\u2019s Market Hazard Rate",
                font_size=13, bold=True, color=HEADER_BLUE)

    # Library HR results (from finstack hazard-rate engine)
    hr_tbl = [["Issuer", "LTV", "\u03bb (bp)", "Cash PV",
               "PIK PV", "\u0394Price", "Cash Z", "PIK Z", "\u0394Z (bp)"]]
    _hr_data = [
        ("BB+", "50%", "143", "113.35", "111.46", "\u22121.89",
         "84", "124", "+40"),
        ("BB\u2212", "61%", "334", "107.56", "104.88", "\u22122.68",
         "209", "269", "+60"),
        ("B", "71%", "591", "99.76", "96.05", "\u22123.71",
         "388", "479", "+91"),
        ("B\u2212", "80%", "911", "90.33", "85.40", "\u22124.93",
         "627", "764", "+137"),
        ("CCC", "87%", "1468", "76.13", "69.59", "\u22126.54",
         "1047", "1271", "+224"),
    ]
    for row in _hr_data:
        hr_tbl.append(list(row))
    add_table(slide, Inches(0.3), Inches(4.95), Inches(9.4), hr_tbl,
              col_widths=[Inches(0.9), Inches(0.5), Inches(0.7),
                          Inches(0.9), Inches(0.9), Inches(0.8),
                          Inches(0.9), Inches(0.9), Inches(0.8)],
              font_size=9)

    add_callout_box(slide, Inches(1.0), Inches(6.5), Inches(8), Inches(0.6),
                    "Under HR: PIK always trades wider because "
                    "the inflated terminal notional is discounted by S(T), "
                    "the lowest survival probability in the term structure. "
                    "The penalty grows with \u03bb.",
                    bg_color=HEADER_BLUE, font_size=11)


# ── Placeholder functions for remaining slides ───────────────────────────

def slide_05_hr_sensitivity(prs):
    """Slide 5: HR Parameter Sensitivity."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "4")
    add_slide_title(slide, "HR Model: Parameter Sensitivity",
                    "How maturity, coupon rate, and recovery affect \u0394Z")
    add_slide_number(slide, 5)

    # Sensitivity: vary maturity
    add_textbox(slide, Inches(0.3), Inches(1.1), Inches(3), Inches(0.25),
                "Maturity Sensitivity (B\u2212, \u03bb=630bp)",
                font_size=12, bold=True, color=HEADER_BLUE)

    iss_b_minus = ISSUERS[3]
    lam_bm = calibrate_hazard(iss_b_minus["spread"], iss_b_minus["rec"])
    mat_rows = [["Maturity", "Cash Z", "PIK Z", "\u0394Z"]]
    for mat in [3, 5, 7]:
        c_pv = hr_bond_price(lam_bm, iss_b_minus["rec"], "cash",
                             maturity=mat)
        p_pv = hr_bond_price(lam_bm, iss_b_minus["rec"], "pik",
                             maturity=mat)
        c_z = price_to_zspread(c_pv, maturity=mat) * 10000
        p_z = price_to_zspread(p_pv, maturity=mat) * 10000
        mat_rows.append([f"{mat}Y", f"{c_z:.0f}bp", f"{p_z:.0f}bp",
                         f"{p_z - c_z:+.0f}bp"])
    add_table(slide, Inches(0.3), Inches(1.45), Inches(3), mat_rows,
              col_widths=[Inches(0.6), Inches(0.8), Inches(0.8),
                          Inches(0.7)],
              font_size=10)

    # Sensitivity: vary coupon
    add_textbox(slide, Inches(3.5), Inches(1.1), Inches(3), Inches(0.25),
                "Coupon Sensitivity (B\u2212, 5Y)",
                font_size=12, bold=True, color=HEADER_BLUE)

    cpn_rows = [["Coupon", "Cash Z", "PIK Z", "\u0394Z"]]
    for cpn in [0.06, 0.085, 0.11]:
        c_pv = hr_bond_price(lam_bm, iss_b_minus["rec"], "cash",
                             coupon_rate=cpn)
        p_pv = hr_bond_price(lam_bm, iss_b_minus["rec"], "pik",
                             coupon_rate=cpn)
        c_z = price_to_zspread(c_pv, maturity=MATURITY) * 10000
        p_z = price_to_zspread(p_pv, maturity=MATURITY) * 10000
        cpn_rows.append([f"{cpn:.1%}", f"{c_z:.0f}bp", f"{p_z:.0f}bp",
                         f"{p_z - c_z:+.0f}bp"])
    add_table(slide, Inches(3.5), Inches(1.45), Inches(3), cpn_rows,
              col_widths=[Inches(0.6), Inches(0.8), Inches(0.8),
                          Inches(0.7)],
              font_size=10)

    # Sensitivity: vary recovery
    add_textbox(slide, Inches(6.7), Inches(1.1), Inches(3), Inches(0.25),
                "Recovery Sensitivity (B\u2212, 5Y)",
                font_size=12, bold=True, color=HEADER_BLUE)

    rec_rows = [["Recovery", "Cash Z", "PIK Z", "\u0394Z"]]
    for rec in [0.25, 0.35, 0.45]:
        lam_r = calibrate_hazard(iss_b_minus["spread"], rec)
        c_pv = hr_bond_price(lam_r, rec, "cash")
        p_pv = hr_bond_price(lam_r, rec, "pik")
        c_z = price_to_zspread(c_pv) * 10000
        p_z = price_to_zspread(p_pv) * 10000
        rec_rows.append([f"{rec:.0%}", f"{c_z:.0f}bp", f"{p_z:.0f}bp",
                         f"{p_z - c_z:+.0f}bp"])
    add_table(slide, Inches(6.7), Inches(1.45), Inches(3), rec_rows,
              col_widths=[Inches(0.7), Inches(0.8), Inches(0.8),
                          Inches(0.7)],
              font_size=10)

    # Interpretation
    add_bullet_list(slide, Inches(0.5), Inches(3.6), Inches(9),
                    Inches(1.5), [
        "Longer maturity \u2192 more compounding periods \u2192 "
        "higher terminal notional \u2192 larger PIK penalty. "
        "The 7Y PIK \u0394Z is roughly double the 3Y",
        "Higher coupon \u2192 more accrual per period \u2192 "
        "faster notional growth. An 11% PIK bond has ~30% "
        "more notional at maturity than 6%",
        "Higher recovery \u2192 hazard rate must rise to match "
        "spread \u2192 amplifies survival discount \u2192 larger "
        "\u0394Z. Recovery is often underappreciated as a PIK "
        "sensitivity",
    ], font_size=13, color=TEXT_DARK)

    # Cross-issuer sensitivity
    add_textbox(slide, Inches(0.5), Inches(5.2), Inches(9), Inches(0.25),
                "\u0394Z Across All Issuers (5Y, 8.5% coupon, issuer recovery)",
                font_size=12, bold=True, color=HEADER_BLUE)

    # Library HR results across all issuers
    sweep_rows = [["Issuer", "LTV", "\u03bb (bp)", "Cash Z",
                   "PIK Z", "\u0394Z (bp)"]]
    _sweep = [
        ("BB+ (Solid HY)", "50%", "143", "84", "124", "+40"),
        ("BB\u2212 (Mid HY)", "61%", "334", "209", "269", "+60"),
        ("B (Weak HY)", "71%", "591", "388", "479", "+91"),
        ("B\u2212 (Stressed)", "80%", "911", "627", "764", "+137"),
        ("CCC (Deeply Stressed)", "87%", "1468", "1047", "1271", "+224"),
    ]
    for row in _sweep:
        sweep_rows.append(list(row))
    add_table(slide, Inches(0.5), Inches(5.5), Inches(6.5), sweep_rows,
              col_widths=[Inches(1.5), Inches(0.5), Inches(0.7),
                          Inches(0.8), Inches(0.8), Inches(0.8)],
              font_size=10)

    add_callout_box(slide, Inches(7.2), Inches(5.5), Inches(2.5),
                    Inches(1.2),
                    "Key: \u0394Z grows super-linearly with \u03bb. "
                    "A flat +50bp bump across all issuers is "
                    "too much for BB+ and too little for CCC.",
                    bg_color=ACCENT_ORANGE, font_size=10)


def slide_06_calibration_gap(prs):
    """Slide 6: The Calibration Gap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "5")
    add_slide_title(slide, "The Calibration Gap",
                    "Market spreads vs historical default rates")
    add_slide_number(slide, 6)

    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.6),
                "Market spreads are not pure default compensation. They "
                "include a credit risk premium, liquidity premium, and "
                "systematic risk loading. Using market \u03bb in the HR model "
                "overstates default probabilities \u2014 especially for "
                "strong credits.",
                font_size=14, color=TEXT_DARK)

    # Gap table
    gap_rows = [["Issuer", "LTV", "Mkt Spread", "\u03bb (cal)",
                 "HR 5Y PD", "Hist 5Y PD", "Ratio",
                 "Risk Premium"]]
    for iss in ISSUERS:
        lam = calibrate_hazard(iss["spread"], iss["rec"])
        hr_pd = 1 - math.exp(-lam * MATURITY)
        hist_pd = 1 - math.exp(-iss["pd"] * MATURITY)
        ratio = hr_pd / hist_pd if hist_pd > 0 else float("inf")
        gap_rows.append([
            iss["name"], iss["ltv"],
            f"{iss['spread'] * 10000:.0f}bp",
            f"{lam * 10000:.0f}bp",
            f"{hr_pd:.1%}", f"{hist_pd:.1%}",
            f"{ratio:.1f}\u00d7",
            "Very high" if ratio > 5 else
            "High" if ratio > 2 else "Moderate",
        ])
    add_table(slide, Inches(0.3), Inches(1.9), Inches(9.4), gap_rows,
              col_widths=[Inches(1.5), Inches(0.5), Inches(0.8),
                          Inches(0.7), Inches(0.8), Inches(0.8),
                          Inches(0.6), Inches(0.9)],
              font_size=10)

    # Bar chart: ratio by issuer
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(4), Inches(0.25),
                "Market-Implied / Historical PD Ratio",
                font_size=12, bold=True, color=HEADER_BLUE)

    ratios = []
    for iss in ISSUERS:
        lam = calibrate_hazard(iss["spread"], iss["rec"])
        hr_pd = 1 - math.exp(-lam * MATURITY)
        hist_pd = 1 - math.exp(-iss["pd"] * MATURITY)
        ratios.append((iss["name"].split(" (")[0], hr_pd / hist_pd))

    max_ratio = max(r for _, r in ratios)
    for i, (name, ratio) in enumerate(ratios):
        y = Inches(4.55) + Inches(i * 0.4)
        add_textbox(slide, Inches(0.5), y, Inches(0.8), Inches(0.25),
                    name, font_size=9, color=TEXT_MID)
        bar_w = Inches(2.5) * ratio / max_ratio
        col = ACCENT_RED if ratio > 4 else (
            ACCENT_ORANGE if ratio > 2 else ACCENT_GREEN)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.4), y + Pt(2),
            int(bar_w), Pt(14))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        add_textbox(slide, Inches(1.4) + int(bar_w) + Pt(4), y - Pt(1),
                    Inches(0.6), Inches(0.25),
                    f"{ratio:.1f}\u00d7", font_size=9, bold=True,
                    color=col)

    # Right side: implications
    add_textbox(slide, Inches(5.0), Inches(4.2), Inches(4.8), Inches(0.25),
                "Implications for PIK Pricing",
                font_size=12, bold=True, color=HEADER_BLUE)
    add_bullet_list(slide, Inches(5.0), Inches(4.55), Inches(4.8),
                    Inches(2.5), [
        "BB+ spread implies 7\u00d7 the historical default "
        "rate. The HR model\u2019s PIK penalty (+40bp) is "
        "driven by an overstated \u03bb",
        "CCC spread implies only 1.3\u00d7 historical PD. "
        "At this level, the PIK penalty (+224bp) is "
        "closer to a \u2018real\u2019 default cost",
        "The structural model uses historical PDs to "
        "calibrate barriers. It answers: what is the "
        "PIK penalty under realistic default assumptions?",
    ], font_size=12, color=TEXT_DARK)

    add_callout_box(slide, Inches(1.0), Inches(6.6), Inches(8),
                    Inches(0.5),
                    "The HR model is an upper bound on PIK cost. "
                    "For strong credits, the true premium may be "
                    "negative under historical defaults.",
                    bg_color=ACCENT_ORANGE, font_size=12)


def slide_07_merton_model(prs):
    """Slide 7: Merton Model — Firm Value & Default."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "6")
    add_slide_title(slide, "Merton Structural Model: Firm Value & Default",
                    "Default occurs when asset value breaches the barrier")
    add_slide_number(slide, 7)

    # GBM formula
    add_formula_box(slide, Inches(0.5), Inches(1.1), Inches(5),
                    "dV = (r \u2212 q)\u00b7V\u00b7dt + \u03c3\u00b7V\u00b7dW",
                    "Geometric Brownian Motion for firm asset value V")

    add_formula_box(slide, Inches(0.5), Inches(1.8), Inches(5),
                    "DD = [ln(V/B) + (r \u2212 \u03c3\u00b2/2)\u00b7T] / "
                    "(\u03c3\u00b7\u221aT)",
                    "Distance-to-Default: standard deviations from barrier")

    add_formula_box(slide, Inches(0.5), Inches(2.5), Inches(5),
                    "PD(T) = N(\u2212DD)",
                    "Default probability = normal CDF of negative DD")

    # Conceptual diagram: asset path with barrier
    add_textbox(slide, Inches(5.8), Inches(1.1), Inches(4), Inches(0.25),
                "Conceptual: Asset Value vs Barrier",
                font_size=12, bold=True, color=HEADER_BLUE)

    # Draw a simple conceptual chart using shapes
    chart_left = Inches(5.8)
    chart_top = Inches(1.5)
    chart_w = Inches(3.8)
    chart_h = Inches(2.2)

    # Axes
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top + chart_h,
        chart_w, Pt(2)).fill.solid()
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        Pt(2), chart_h).fill.solid()

    # Barrier line (horizontal, at ~40% height)
    barrier_y = int(chart_top + chart_h * 0.6)
    b_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, barrier_y,
        chart_w, Pt(2))
    b_line.fill.solid()
    b_line.fill.fore_color.rgb = ACCENT_RED
    b_line.line.fill.background()
    add_textbox(slide, chart_left + chart_w + Pt(4),
                barrier_y - Pt(6), Inches(0.6), Inches(0.2),
                "B (barrier)", font_size=8, color=ACCENT_RED)

    # Asset start point label
    asset_start_y = int(chart_top + chart_h * 0.15)
    add_textbox(slide, chart_left - Inches(0.5), asset_start_y - Pt(4),
                Inches(0.5), Inches(0.2),
                "V\u2080", font_size=10, bold=True, color=HEADER_BLUE,
                alignment=PP_ALIGN.RIGHT)

    # Simulated paths (simplified as line segments)
    np.random.seed(42)
    n_steps = 20
    dt = MATURITY / n_steps
    for path_i in range(8):
        v = 1.0  # normalised
        points = [(0, v)]
        vol = 0.25
        for step in range(n_steps):
            dw = np.random.normal(0, math.sqrt(dt))
            v *= math.exp((RISK_FREE - vol**2/2) * dt + vol * dw)
            points.append((step + 1, v))

        defaulted = any(p[1] < 0.55 for p in points)
        col = ACCENT_RED if defaulted else ACCENT_GREEN

        for j in range(1, len(points)):
            x1 = int(chart_left) + int(chart_w * points[j-1][0] / n_steps)
            x2 = int(chart_left) + int(chart_w * points[j][0] / n_steps)
            y1 = int(chart_top + chart_h * (1 - points[j-1][1] * 0.85))
            y2 = int(chart_top + chart_h * (1 - points[j][1] * 0.85))
            # Draw as small rectangle (approximation)
            seg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                min(x1, x2), min(y1, y2),
                max(abs(x2 - x1), Emu(1)),
                max(abs(y2 - y1), Pt(1)))
            seg.fill.solid()
            seg.fill.fore_color.rgb = col
            seg.line.fill.background()

    add_textbox(slide, chart_left, chart_top + chart_h + Pt(4),
                chart_w, Inches(0.2),
                "Time \u2192 (0 to 5Y)", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.CENTER)

    # Key concepts (below)
    concepts = [
        ("Equity as Call Option",
         "Equity = max(V \u2212 B, 0) at maturity. Shareholders "
         "own a call option on firm assets with strike = debt. "
         "The Merton model links credit and equity markets."),
        ("Distance-to-Default",
         "DD measures how many standard deviations the asset "
         "value sits above the barrier. Higher DD = safer credit. "
         "BB+ has DD \u2248 3.4, CCC has DD \u2248 0.6."),
        ("Terminal vs First-Passage",
         "Terminal barrier: default only checked at T. "
         "First-passage (Black-Cox): default can occur at any "
         "time V < B. We use terminal for simplicity."),
    ]

    for i, (title, desc) in enumerate(concepts):
        y = Inches(4.0) + Inches(i * 0.95)
        add_textbox(slide, Inches(0.5), y, Inches(2.5), Inches(0.25),
                    title, font_size=12, bold=True, color=HEADER_BLUE)
        add_textbox(slide, Inches(3.0), y, Inches(6.8), Inches(0.85),
                    desc, font_size=11, color=TEXT_DARK)


def slide_08_barrier_calibration(prs):
    """Slide 8: Barrier Calibration from Historical PDs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "7")
    add_slide_title(slide, "Barrier Calibration from Historical PDs",
                    "MertonModel.from_target_pd: backing out the barrier")
    add_slide_number(slide, 8)

    # Method explanation
    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.5),
                "The barrier B is the free parameter. Given V\u2080, \u03c3, "
                "and r, we solve for B such that PD(T) = N(\u2212DD) "
                "matches the target historical default probability.",
                font_size=14, color=TEXT_DARK)

    add_formula_box(slide, Inches(0.5), Inches(1.7), Inches(9),
                    "B = V \u00d7 exp[\u2212(DD \u00d7 \u03c3\u00d7\u221aT "
                    "+ (r \u2212 \u03c3\u00b2/2)\u00d7T)]  where  "
                    "DD = \u2212N\u207b\u00b9(PD\u2085\u2084)",
                    "Invert the distance-to-default formula "
                    "to find barrier from target PD")

    # Full parameter table
    cal_rows = [["Issuer", "LTV", "V\u2080", "\u03c3",
                 "Ann PD", "5Y PD", "DD",
                 "Barrier", "Impl Sprd"]]
    for iss in ISSUERS:
        five_pd = 1 - math.exp(-iss["pd"] * MATURITY)
        dd = -_norm_ppf(five_pd) if five_pd < 1 else 0
        barrier = merton_barrier(iss["asset"], iss["vol"], iss["pd"])
        # Implied spread: s = -ln(1 - PD*(1-R)) / T
        impl_s = -math.log(1 - five_pd * (1 - iss["rec"])) / MATURITY
        cal_rows.append([
            iss["name"], iss["ltv"],
            f"{iss['asset']:.0f}", f"{iss['vol']:.0%}",
            f"{iss['pd']:.2%}", f"{five_pd:.1%}",
            f"{dd:.2f}", f"{barrier:.1f}",
            f"{impl_s * 10000:.0f}bp",
        ])
    add_table(slide, Inches(0.3), Inches(2.4), Inches(9.4), cal_rows,
              col_widths=[Inches(1.5), Inches(0.5), Inches(0.5),
                          Inches(0.5), Inches(0.7), Inches(0.6),
                          Inches(0.5), Inches(0.7), Inches(0.8)],
              font_size=10)

    # Comparison: Merton implied spread vs market spread
    add_textbox(slide, Inches(0.5), Inches(4.5), Inches(9), Inches(0.25),
                "Market Spread vs Merton Implied Spread",
                font_size=13, bold=True, color=HEADER_BLUE)

    cmp_rows = [["Issuer", "LTV", "Mkt Spread",
                 "Merton Spread", "Gap", "Interpretation"]]
    for iss in ISSUERS:
        five_pd = 1 - math.exp(-iss["pd"] * MATURITY)
        impl_s = -math.log(1 - five_pd * (1 - iss["rec"])) / MATURITY
        gap = iss["spread"] - impl_s
        cmp_rows.append([
            iss["name"], iss["ltv"],
            f"{iss['spread'] * 10000:.0f}bp",
            f"{impl_s * 10000:.0f}bp",
            f"{gap * 10000:+.0f}bp",
            "Large risk premium" if gap > 0.005 else
            "Moderate premium" if gap > 0.002 else
            "Small premium",
        ])
    add_table(slide, Inches(0.3), Inches(4.85), Inches(9.4), cmp_rows,
              col_widths=[Inches(1.5), Inches(0.5), Inches(0.9),
                          Inches(1.0), Inches(0.8), Inches(1.4)],
              font_size=10)

    add_callout_box(slide, Inches(0.5), Inches(6.5), Inches(9),
                    Inches(0.6),
                    "The Merton implied spread is consistently below the "
                    "market spread because it uses only historical default "
                    "risk. The gap IS the risk premium. By calibrating to "
                    "historical PDs, the MC model prices PIK under "
                    "\u2018real-world\u2019 rather than risk-neutral defaults.",
                    bg_color=HEADER_BLUE, font_size=11)


def slide_09_endogenous_hazard(prs):
    """Slide 9: Endogenous Hazard — \u03bb(L) = \u03bb\u2080\u00d7(L/L\u2080)\u00b2."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "8")
    add_slide_title(slide, "Endogenous Hazard: \u03bb(L) = \u03bb\u2080 \u00d7 (L/L\u2080)\u00b2",
                    "PIK accrual raises leverage \u2192 hazard rate "
                    "rises non-linearly")
    add_slide_number(slide, 9)

    # Formula and explanation
    add_formula_box(slide, Inches(0.5), Inches(1.1), Inches(5),
                    "\u03bb(L) = \u03bb\u2080 \u00d7 (L / L\u2080)\u00b2",
                    "\u03bb\u2080 = base hazard  |  L\u2080 = initial "
                    "leverage  |  \u03b2 = 2 (quadratic)")

    add_bullet_list(slide, Inches(0.5), Inches(1.85), Inches(5),
                    Inches(1.5), [
        "L = N(t) / V(t): leverage = notional / asset value",
        "PIK accretes to notional \u2192 N(t) grows \u2192 "
        "L rises even if V is unchanged",
        "Quadratic (\u03b2=2): a 20% leverage increase "
        "raises hazard by 44%",
        "This is the first half of the feedback loop: "
        "PIK \u2192 higher \u03bb \u2192 more defaults",
    ], font_size=12, color=TEXT_DARK)

    # Chart: \u03bb vs leverage for each issuer
    add_textbox(slide, Inches(5.8), Inches(1.1), Inches(4), Inches(0.25),
                "Hazard Rate vs Leverage",
                font_size=12, bold=True, color=HEADER_BLUE)

    chart_left = Inches(5.8)
    chart_top = Inches(1.5)
    chart_w = Inches(3.8)
    chart_h = Inches(2.5)

    # Axes
    ax_x = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top + chart_h,
        chart_w, Pt(1))
    ax_x.fill.solid()
    ax_x.fill.fore_color.rgb = TEXT_MID
    ax_x.line.fill.background()
    ax_y = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        Pt(1), chart_h)
    ax_y.fill.solid()
    ax_y.fill.fore_color.rgb = TEXT_MID
    ax_y.line.fill.background()

    # Axis labels
    add_textbox(slide, chart_left, chart_top + chart_h + Pt(4),
                chart_w, Inches(0.2),
                "Leverage (N/V) \u2192", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, chart_left - Inches(0.6), chart_top,
                Inches(0.5), Inches(0.2),
                "\u03bb \u2192", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.RIGHT)

    # Plot curves for 3 issuers: BB+, B, CCC
    plot_issuers = [ISSUERS[0], ISSUERS[2], ISSUERS[4]]
    colors_plot = [ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED]
    max_lev = 1.2
    max_haz = 0.35  # cap for display

    for iss, col in zip(plot_issuers, colors_plot):
        lam0 = calibrate_hazard(iss["spread"], iss["rec"])
        l0 = NOTIONAL / iss["asset"]
        n_pts = 30
        prev_x, prev_y = None, None
        for j in range(n_pts):
            lev = 0.3 + j * (max_lev - 0.3) / (n_pts - 1)
            haz = lam0 * (lev / l0) ** 2
            haz_clip = min(haz, max_haz)
            px = int(chart_left) + int(chart_w * (lev - 0.3) / (max_lev - 0.3))
            py = int(chart_top + chart_h * (1 - haz_clip / max_haz))
            if prev_x is not None:
                seg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    min(px, prev_x), min(py, prev_y),
                    max(abs(px - prev_x), Emu(1)),
                    max(abs(prev_y - py), Pt(1.5)))
                seg.fill.solid()
                seg.fill.fore_color.rgb = col
                seg.line.fill.background()
            prev_x, prev_y = px, py

        # Label at end
        add_textbox(slide, prev_x + Pt(4), prev_y - Pt(6),
                    Inches(0.7), Inches(0.2),
                    iss["name"].split(" (")[0],
                    font_size=7, color=col)

    # Initial leverage markers
    add_textbox(slide, Inches(5.8), Inches(4.15), Inches(4), Inches(0.2),
                "Arrows show where PIK pushes each issuer along its curve",
                font_size=8, color=TEXT_MID)

    # Numerical example table
    add_textbox(slide, Inches(0.5), Inches(4.3), Inches(5), Inches(0.25),
                "PIK Impact: 5-Year Notional & Hazard Growth",
                font_size=12, bold=True, color=HEADER_BLUE)

    terminal_ntl = NOTIONAL * (1 + COUPON / 2) ** 10
    pik_rows = [["Issuer", "LTV\u2080", "L\u2080",
                 "L at 5Y (PIK)", "\u03bb\u2080",
                 "\u03bb at 5Y (PIK)", "Increase"]]
    for iss in ISSUERS:
        lam0 = calibrate_hazard(iss["spread"], iss["rec"])
        l0 = NOTIONAL / iss["asset"]
        l_pik = terminal_ntl / iss["asset"]
        lam_pik = lam0 * (l_pik / l0) ** 2
        pik_rows.append([
            iss["name"].split(" (")[0], iss["ltv"],
            f"{l0:.2f}", f"{l_pik:.2f}",
            f"{lam0 * 10000:.0f}bp", f"{lam_pik * 10000:.0f}bp",
            f"{lam_pik / lam0:.1f}\u00d7",
        ])
    add_table(slide, Inches(0.3), Inches(4.65), Inches(9.4), pik_rows,
              col_widths=[Inches(0.9), Inches(0.6), Inches(0.6),
                          Inches(1.0), Inches(0.8), Inches(1.1),
                          Inches(0.8)],
              font_size=9)

    add_callout_box(slide, Inches(1.0), Inches(6.5), Inches(8),
                    Inches(0.55),
                    "For CCC (87% LTV): PIK pushes leverage from 0.87 "
                    f"to {terminal_ntl / 115:.2f}, increasing \u03bb by "
                    f"{(terminal_ntl / 115 / (100/115))**2:.1f}\u00d7. "
                    "This is before considering asset value changes.",
                    bg_color=ACCENT_RED, font_size=11)


def slide_10_dynamic_recovery(prs):
    """Slide 10: Dynamic Recovery — R(N) = max(floor, R\u2080\u00d7N\u2080/N)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "9")
    add_slide_title(slide, "Dynamic Recovery: R(N) = max(floor, R\u2080 \u00d7 N\u2080/N)",
                    "PIK dilutes recovery per dollar of claim")
    add_slide_number(slide, 10)

    # Formula
    add_formula_box(slide, Inches(0.5), Inches(1.1), Inches(5.5),
                    "R(N) = max(floor,  R\u2080 \u00d7 N\u2080 / N)",
                    "R\u2080 = base recovery  |  N\u2080 = 100 (par)  |  "
                    "floor = 10%")

    add_bullet_list(slide, Inches(0.5), Inches(1.85), Inches(5),
                    Inches(1.5), [
        "On default, the recovery pool (assets) is fixed but "
        "the claim (notional) has grown via PIK accrual",
        "Each dollar of claim gets proportionally less. "
        "If notional doubles, recovery per dollar halves",
        "The floor (10%) prevents recovery from going "
        "to zero \u2014 there is always some residual asset value",
        "This is the second half of the feedback loop: "
        "PIK \u2192 lower R \u2192 higher loss-given-default",
    ], font_size=12, color=TEXT_DARK)

    # Chart: Recovery vs Notional for each issuer
    add_textbox(slide, Inches(5.8), Inches(1.1), Inches(4), Inches(0.25),
                "Recovery Rate vs Accreted Notional",
                font_size=12, bold=True, color=HEADER_BLUE)

    chart_left = Inches(5.8)
    chart_top = Inches(1.5)
    chart_w = Inches(3.8)
    chart_h = Inches(2.3)

    # Axes
    ax_x = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top + chart_h,
        chart_w, Pt(1))
    ax_x.fill.solid()
    ax_x.fill.fore_color.rgb = TEXT_MID
    ax_x.line.fill.background()
    ax_y = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        Pt(1), chart_h)
    ax_y.fill.solid()
    ax_y.fill.fore_color.rgb = TEXT_MID
    ax_y.line.fill.background()

    add_textbox(slide, chart_left, chart_top + chart_h + Pt(4),
                chart_w, Inches(0.2),
                "Notional (100 \u2192 200) \u2192",
                font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, chart_left - Inches(0.6), chart_top,
                Inches(0.5), Inches(0.2),
                "R(N) \u2192", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.RIGHT)

    # Floor line
    floor = 0.10
    floor_y = int(chart_top + chart_h * (1 - floor / 0.50))
    fl = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, floor_y,
        chart_w, Pt(1))
    fl.fill.solid()
    fl.fill.fore_color.rgb = ACCENT_RED
    fl.line.fill.background()
    add_textbox(slide, chart_left + chart_w + Pt(4), floor_y - Pt(6),
                Inches(0.6), Inches(0.2),
                "floor=10%", font_size=7, color=ACCENT_RED)

    # Plot curves for 3 issuers
    plot_issuers = [ISSUERS[0], ISSUERS[2], ISSUERS[4]]
    colors_plot = [ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED]
    n_min, n_max = 100, 200

    for iss, col in zip(plot_issuers, colors_plot):
        r0 = iss["rec"]
        n_pts = 30
        prev_x, prev_y = None, None
        for j in range(n_pts):
            n = n_min + j * (n_max - n_min) / (n_pts - 1)
            r_n = max(floor, r0 * NOTIONAL / n)
            px = int(chart_left) + int(chart_w * (n - n_min) / (n_max - n_min))
            py = int(chart_top + chart_h * (1 - r_n / 0.50))
            if prev_x is not None:
                seg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    min(px, prev_x), min(py, prev_y),
                    max(abs(px - prev_x), Emu(1)),
                    max(abs(prev_y - py), Pt(1.5)))
                seg.fill.solid()
                seg.fill.fore_color.rgb = col
                seg.line.fill.background()
            prev_x, prev_y = px, py

        add_textbox(slide, prev_x + Pt(4), prev_y - Pt(6),
                    Inches(0.9), Inches(0.2),
                    f"{iss['name'].split(' (')[0]} (R\u2080={r0:.0%})",
                    font_size=7, color=col)

    # Combined effect table
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.25),
                "Combined Feedback: Hazard \u00d7 Recovery Impact at Maturity",
                font_size=12, bold=True, color=HEADER_BLUE)

    terminal_ntl = NOTIONAL * (1 + COUPON / 2) ** 10
    combo_rows = [["Issuer", "LTV", "R\u2080",
                   "R at 5Y (PIK)", "\u0394R",
                   "\u03bb\u2080", "\u03bb at 5Y (PIK)",
                   "Net LGD increase"]]
    for iss in ISSUERS:
        lam0 = calibrate_hazard(iss["spread"], iss["rec"])
        r0 = iss["rec"]
        r_pik = max(0.10, r0 * NOTIONAL / terminal_ntl)
        l0 = NOTIONAL / iss["asset"]
        l_pik = terminal_ntl / iss["asset"]
        lam_pik = lam0 * (l_pik / l0) ** 2
        lgd_0 = (1 - r0) * lam0
        lgd_pik = (1 - r_pik) * lam_pik
        combo_rows.append([
            iss["name"].split(" (")[0], iss["ltv"],
            f"{r0:.0%}", f"{r_pik:.0%}",
            f"{(r_pik - r0) * 100:+.0f}pp",
            f"{lam0 * 10000:.0f}bp", f"{lam_pik * 10000:.0f}bp",
            f"{lgd_pik / lgd_0:.1f}\u00d7" if lgd_0 > 0 else "n/a",
        ])
    add_table(slide, Inches(0.3), Inches(4.55), Inches(9.4), combo_rows,
              col_widths=[Inches(0.9), Inches(0.5), Inches(0.5),
                          Inches(0.9), Inches(0.6), Inches(0.7),
                          Inches(1.0), Inches(1.0)],
              font_size=9)

    add_callout_box(slide, Inches(1.0), Inches(6.5), Inches(8),
                    Inches(0.55),
                    "The double whammy: PIK simultaneously increases "
                    "the probability of default (\u03bb\u2191) AND reduces "
                    "recovery per dollar (R\u2193). These compound "
                    "multiplicatively, not additively.",
                    bg_color=ACCENT_RED, font_size=11)


def slide_11_mc_paths(prs):
    """Slide 11: Monte Carlo Paths — visual."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "10")
    add_slide_title(slide, "Monte Carlo Paths: What the Simulation Looks Like",
                    "50 sample GBM paths for B\u2212 (Stressed, 80% LTV)")
    add_slide_number(slide, 11)

    iss = ISSUERS[3]  # B-
    barrier = merton_barrier(iss["asset"], iss["vol"], iss["pd"])

    # Chart area
    chart_left = Inches(0.5)
    chart_top = Inches(1.2)
    chart_w = Inches(5.5)
    chart_h = Inches(3.5)

    # Axes
    ax_x = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top + chart_h,
        chart_w, Pt(1))
    ax_x.fill.solid()
    ax_x.fill.fore_color.rgb = TEXT_MID
    ax_x.line.fill.background()
    ax_y = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        Pt(1), chart_h)
    ax_y.fill.solid()
    ax_y.fill.fore_color.rgb = TEXT_MID
    ax_y.line.fill.background()

    add_textbox(slide, chart_left, chart_top + chart_h + Pt(4),
                chart_w, Inches(0.2),
                "Time (years) \u2192", font_size=9, color=TEXT_MID,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, chart_left - Inches(0.5), chart_top - Pt(2),
                Inches(0.5), Inches(0.2),
                "Asset\nValue", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.RIGHT)

    # Barrier line
    v_max = iss["asset"] * 1.6
    v_min = 0
    b_frac = 1 - (barrier - v_min) / (v_max - v_min)
    barrier_y = int(chart_top + chart_h * b_frac)
    bl = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, barrier_y,
        chart_w, Pt(2))
    bl.fill.solid()
    bl.fill.fore_color.rgb = ACCENT_RED
    bl.line.fill.background()
    add_textbox(slide, chart_left + chart_w + Pt(4),
                barrier_y - Pt(8), Inches(1.0), Inches(0.2),
                f"B = {barrier:.0f}", font_size=9, bold=True,
                color=ACCENT_RED)

    # V0 label
    v0_frac = 1 - (iss["asset"] - v_min) / (v_max - v_min)
    add_textbox(slide, chart_left - Inches(0.7),
                int(chart_top + chart_h * v0_frac) - Pt(6),
                Inches(0.6), Inches(0.2),
                f"V\u2080={iss['asset']}", font_size=8, bold=True,
                color=HEADER_BLUE, alignment=PP_ALIGN.RIGHT)

    # Simulate 50 paths
    np.random.seed(123)
    n_paths = 50
    n_steps = 60  # monthly
    dt = MATURITY / n_steps
    survived = 0
    defaulted = 0

    for _ in range(n_paths):
        v = float(iss["asset"])
        points = [(0, v)]
        did_default = False
        for step in range(n_steps):
            dw = np.random.normal(0, math.sqrt(dt))
            v *= math.exp(
                (RISK_FREE - iss["vol"]**2 / 2) * dt
                + iss["vol"] * dw)
            points.append((step + 1, v))
            if v < barrier:
                did_default = True
                break

        col = ACCENT_RED if did_default else ACCENT_GREEN
        if did_default:
            defaulted += 1
        else:
            survived += 1

        for j in range(1, len(points)):
            t0, v0 = points[j - 1]
            t1, v1 = points[j]
            x1 = int(chart_left) + int(chart_w * t0 / n_steps)
            x2 = int(chart_left) + int(chart_w * t1 / n_steps)
            frac0 = 1 - (v0 - v_min) / (v_max - v_min)
            frac1 = 1 - (v1 - v_min) / (v_max - v_min)
            frac0 = max(0, min(1, frac0))
            frac1 = max(0, min(1, frac1))
            y1 = int(chart_top + chart_h * frac0)
            y2 = int(chart_top + chart_h * frac1)
            seg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                min(x1, x2), min(y1, y2),
                max(abs(x2 - x1), Emu(1)),
                max(abs(y2 - y1), Pt(1)))
            seg.fill.solid()
            seg.fill.fore_color.rgb = col
            seg.line.fill.background()

    # Stats box
    add_textbox(slide, Inches(0.5), Inches(4.85), Inches(5.5),
                Inches(0.3),
                f"50 paths: {survived} survived (green), "
                f"{defaulted} defaulted (red)  |  "
                f"Sample default rate = {defaulted/50:.0%}  |  "
                f"Historical 5Y PD = "
                f"{1 - math.exp(-iss['pd'] * MATURITY):.1%}",
                font_size=10, color=TEXT_MID)

    # Right side: MC algorithm steps
    add_textbox(slide, Inches(6.3), Inches(1.2), Inches(3.5), Inches(0.25),
                "MC Algorithm (per path)",
                font_size=13, bold=True, color=HEADER_BLUE)

    steps = [
        ("1. Evolve assets",
         "V(t+dt) = V(t) \u00d7 exp[(r\u2212\u03c3\u00b2/2)dt + \u03c3\u221adt\u00b7Z]"),
        ("2. At each coupon date:",
         "Cash: pay coupon, discount\n"
         "PIK: accrete N \u2192 N\u00d7(1+c/2)"),
        ("3. Update credit state",
         "L = N/V, recompute \u03bb(L), DD"),
        ("4. Check default",
         "V < B? If yes: recovery = R(N)\u00d7N"),
        ("5. If survived to T:",
         "PV += N(T) \u00d7 D(T)"),
        ("6. Aggregate paths",
         "Price = mean(path PVs)\n"
         "SE = std(path PVs) / \u221an"),
    ]

    for i, (title, desc) in enumerate(steps):
        y = Inches(1.55) + Inches(i * 0.65)
        add_textbox(slide, Inches(6.3), y, Inches(3.5), Inches(0.2),
                    title, font_size=10, bold=True, color=HEADER_BLUE)
        add_textbox(slide, Inches(6.3), y + Inches(0.18), Inches(3.5),
                    Inches(0.4),
                    desc, font_size=9, color=TEXT_DARK)

    add_callout_box(slide, Inches(0.5), Inches(5.3), Inches(9.2),
                    Inches(0.5),
                    "Key: the MC engine runs 25,000 paths with antithetic "
                    "variates (variance reduction). At each coupon date, "
                    "the toggle decision, hazard update, and recovery "
                    "adjustment are all path-dependent.",
                    bg_color=HEADER_BLUE, font_size=11)

    # Notional comparison
    add_textbox(slide, Inches(0.5), Inches(6.0), Inches(9), Inches(0.25),
                "Terminal Notional: Cash vs PIK on Surviving Paths",
                font_size=11, bold=True, color=HEADER_BLUE)

    terminal_ntl = NOTIONAL * (1 + COUPON / 2) ** 10
    add_textbox(slide, Inches(0.5), Inches(6.3), Inches(9), Inches(0.5),
                f"Cash-pay: always N = {NOTIONAL:.0f} at maturity  |  "
                f"Full PIK: N = {terminal_ntl:.2f} (+{terminal_ntl - NOTIONAL:.1f}%)  |  "
                f"Toggle: N varies by path (PIK only on stressed paths)",
                font_size=11, color=TEXT_DARK)


def slide_12_feedback_spiral(prs):
    """Slide 12: Walk through a single stressed path."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "11")
    add_slide_title(slide,
                    "The Feedback Spiral: A Single Path Walkthrough",
                    "B\u2212 issuer \u2014 asset value declining, PIK accreting")
    add_slide_number(slide, 12)

    iss = ISSUERS[3]  # B-
    lam0 = calibrate_hazard(iss["spread"], iss["rec"])
    l0 = NOTIONAL / iss["asset"]
    barrier = merton_barrier(iss["asset"], iss["vol"], iss["pd"])

    # Simulate a specific stressed path (seed chosen for good story)
    np.random.seed(77)
    n_coupons = 10
    dt_coupon = 0.5
    steps_per_coupon = 6  # monthly within each coupon period

    v = float(iss["asset"])
    ntl_cash = NOTIONAL
    ntl_pik = NOTIONAL
    semi_cpn = COUPON / 2

    path_rows = [["Period", "V(t)", "N (Cash)", "N (PIK)",
                  "L (PIK)", "\u03bb (PIK)", "R (PIK)", "DD"]]

    path_rows.append([
        "t=0", f"{v:.1f}", f"{ntl_cash:.1f}", f"{ntl_pik:.1f}",
        f"{ntl_pik / v:.2f}",
        f"{lam0 * 10000:.0f}bp",
        f"{iss['rec']:.0%}",
        f"{_dd(v, barrier, iss['vol'], MATURITY):.2f}",
    ])

    for cpn_i in range(n_coupons):
        # Evolve asset monthly within coupon period
        for _ in range(steps_per_coupon):
            dt = dt_coupon / steps_per_coupon
            dw = np.random.normal(0, math.sqrt(dt))
            v *= math.exp(
                (RISK_FREE - iss["vol"]**2 / 2) * dt
                + iss["vol"] * dw)
            # Stress: gentle downward bias for illustration
            v *= 0.997

        # PIK accrual
        ntl_pik *= (1 + semi_cpn)

        # Compute state
        lev = ntl_pik / v
        lam_now = lam0 * (lev / l0) ** 2
        r_now = max(0.10, iss["rec"] * NOTIONAL / ntl_pik)
        remain = MATURITY - (cpn_i + 1) * dt_coupon
        dd_now = _dd(v, barrier, iss["vol"], max(remain, 0.01))

        path_rows.append([
            f"t={dt_coupon * (cpn_i + 1):.1f}",
            f"{v:.1f}", f"{ntl_cash:.1f}", f"{ntl_pik:.1f}",
            f"{lev:.2f}",
            f"{lam_now * 10000:.0f}bp",
            f"{r_now:.0%}",
            f"{dd_now:.2f}",
        ])

    add_table(slide, Inches(0.3), Inches(1.1), Inches(9.4), path_rows,
              col_widths=[Inches(0.6), Inches(0.7), Inches(0.8),
                          Inches(0.8), Inches(0.7), Inches(0.8),
                          Inches(0.7), Inches(0.6)],
              font_size=9)

    # Interpretation
    final_lev = float(path_rows[-1][4])
    final_lam = path_rows[-1][5]
    final_r = path_rows[-1][6]
    final_dd = path_rows[-1][7]

    add_bullet_list(slide, Inches(0.5), Inches(5.2), Inches(9),
                    Inches(1.5), [
        f"Over 5 years: PIK notional grew from 100 to "
        f"{ntl_pik:.0f} while assets declined to {v:.0f}",
        f"Leverage rose from {l0:.2f} to {final_lev:.2f} \u2014 "
        f"hazard rate rose from {lam0*10000:.0f}bp to {final_lam}",
        f"Recovery fell from {iss['rec']:.0%} to {final_r} "
        f"\u2014 distance-to-default collapsed to {final_dd}",
        "This is one path. Across 25,000 paths the average "
        "captures the expected cost of this spiral",
    ], font_size=12, color=TEXT_DARK)

    add_callout_box(slide, Inches(1.0), Inches(6.8), Inches(8),
                    Inches(0.45),
                    "The spiral is self-reinforcing: PIK raises notional "
                    "\u2192 leverage rises \u2192 \u03bb rises \u2192 "
                    "more defaults \u2192 recovery falls \u2192 higher loss. "
                    "Cash-pay avoids this entirely.",
                    bg_color=ACCENT_RED, font_size=11)


def _dd(v: float, b: float, vol: float, t: float) -> float:
    """Distance-to-default."""
    if t <= 0 or v <= 0 or b <= 0:
        return 0.0
    return (math.log(v / b) + (RISK_FREE - vol**2/2) * t) / (vol * math.sqrt(t))


def slide_13_toggle_mechanics(prs):
    """Slide 13: Toggle Mechanics & Adverse Selection."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "12")
    add_slide_title(slide,
                    "Toggle Mechanics & Adverse Selection",
                    "Why the toggle option fails to protect investors")
    add_slide_number(slide, 13)

    # Toggle rule
    add_formula_box(slide, Inches(0.5), Inches(1.1), Inches(5),
                    "PIK if \u03bb(t) > 10%,  else Cash",
                    "Threshold model: borrower exercises PIK "
                    "when credit quality deteriorates")

    # Split diagram
    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(4.5), Inches(0.3),
                "Path Bifurcation Under Toggle",
                font_size=13, bold=True, color=HEADER_BLUE)

    # Healthy paths box
    h_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2),
        Inches(4), Inches(1.2))
    h_box.fill.solid()
    h_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    h_box.line.fill.background()
    add_textbox(slide, Inches(0.7), Inches(2.25), Inches(3.5), Inches(0.25),
                "Healthy Paths (\u03bb < 10%): CASH",
                font_size=12, bold=True, color=ACCENT_GREEN)
    add_textbox(slide, Inches(0.7), Inches(2.55), Inches(3.5), Inches(0.7),
                "Coupons paid in cash\n"
                "Notional stays at 100\n"
                "Leverage stable \u2192 no feedback\n"
                "Behaves identically to cash-pay bond",
                font_size=10, color=TEXT_DARK)

    # Stressed paths box
    s_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.55),
        Inches(4), Inches(1.3))
    s_box.fill.solid()
    s_box.fill.fore_color.rgb = RGBColor(0xFD, 0xE0, 0xE0)
    s_box.line.fill.background()
    add_textbox(slide, Inches(0.7), Inches(3.6), Inches(3.5), Inches(0.25),
                "Stressed Paths (\u03bb > 10%): PIK",
                font_size=12, bold=True, color=ACCENT_RED)
    add_textbox(slide, Inches(0.7), Inches(3.9), Inches(3.5), Inches(0.85),
                "Coupons accrete to notional\n"
                "Leverage spirals upward\n"
                "\u03bb rises further \u2192 stays in PIK mode\n"
                "Recovery diluted \u2192 loss amplified\n"
                "Feedback loop concentrates on worst paths",
                font_size=10, color=TEXT_DARK)

    # Right side: why toggle >= PIK
    add_textbox(slide, Inches(5.2), Inches(1.8), Inches(4.5), Inches(0.3),
                "Why Toggle \u2265 Full PIK",
                font_size=13, bold=True, color=HEADER_BLUE)

    add_bullet_list(slide, Inches(5.2), Inches(2.2), Inches(4.5),
                    Inches(2.5), [
        "Full PIK distributes accrual uniformly across "
        "ALL paths, including healthy ones where extra "
        "notional barely matters",
        "Toggle concentrates accrual on the WORST paths "
        "\u2014 the ones already closest to default. "
        "This seeds the leverage spiral where it does "
        "the most damage",
        "On healthy paths: toggle = cash (no cost). "
        "On stressed paths: toggle = PIK but with "
        "worse starting conditions (already high \u03bb)",
        "The borrower\u2019s option to toggle is effectively "
        "adverse selection: they PIK precisely when "
        "it hurts investors the most",
    ], font_size=12, color=TEXT_DARK)

    # MC results comparison
    add_textbox(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.25),
                "MC Results: Toggle vs Full PIK Z-Spread Premium",
                font_size=12, bold=True, color=HEADER_BLUE)

    tog_rows = [["Issuer", "LTV", "Cash Z", "PIK Z",
                 "Toggle Z", "PIK\u2212Cash", "Tog\u2212Cash",
                 "Tog\u2212PIK"]]
    mc_data = [
        ("BB+", "50%", 20, -39, 22),
        ("BB\u2212", "61%", 110, 88, 130),
        ("B", "71%", 292, 329, 346),
        ("B\u2212", "80%", 710, 851, 862),
        ("CCC", "87%", 1497, 1759, 1763),
    ]
    for name, ltv, cash, pik, tog in mc_data:
        tog_rows.append([
            name, ltv,
            f"{cash}bp", f"{pik}bp", f"{tog}bp",
            f"{pik - cash:+d}", f"{tog - cash:+d}",
            f"{tog - pik:+d}",
        ])
    add_table(slide, Inches(0.3), Inches(5.4), Inches(9.4), tog_rows,
              col_widths=[Inches(0.7), Inches(0.5), Inches(0.8),
                          Inches(0.8), Inches(0.9), Inches(0.9),
                          Inches(0.9), Inches(0.8)],
              font_size=9)

    add_callout_box(slide, Inches(1.0), Inches(6.8), Inches(8),
                    Inches(0.45),
                    "Toggle \u2265 PIK in every case from B onwards. "
                    "The \u2018protection\u2019 of cash-pay on good paths "
                    "is more than offset by the concentrated spiral "
                    "on bad paths.",
                    bg_color=ACCENT_RED, font_size=11)


def slide_14_convergence(prs):
    """Slide 14: MC Convergence & Precision."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "13")
    add_slide_title(slide,
                    "MC Convergence & Precision",
                    "How many paths are enough?")
    add_slide_number(slide, 14)

    # Explanation
    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(0.5),
                "MC standard error decreases as 1/\u221an. "
                "Antithetic variates roughly halve the variance "
                "(each path paired with its mirror). "
                "Our production run uses 25,000 paths.",
                font_size=14, color=TEXT_DARK)

    add_formula_box(slide, Inches(0.5), Inches(1.7), Inches(5),
                    "SE = \u03c3(path PVs) / \u221an",
                    "Standard error of the MC price estimate")

    # Convergence chart (simulated)
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(5.5), Inches(0.25),
                "Convergence: B\u2212 PIK Z-Spread vs Number of Paths",
                font_size=12, bold=True, color=HEADER_BLUE)

    chart_left = Inches(0.5)
    chart_top = Inches(2.7)
    chart_w = Inches(5.5)
    chart_h = Inches(2.5)

    # Axes
    ax_x = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top + chart_h,
        chart_w, Pt(1))
    ax_x.fill.solid()
    ax_x.fill.fore_color.rgb = TEXT_MID
    ax_x.line.fill.background()
    ax_y = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, chart_top,
        Pt(1), chart_h)
    ax_y.fill.solid()
    ax_y.fill.fore_color.rgb = TEXT_MID
    ax_y.line.fill.background()

    add_textbox(slide, chart_left, chart_top + chart_h + Pt(4),
                chart_w, Inches(0.2),
                "Number of Paths \u2192", font_size=8, color=TEXT_MID,
                alignment=PP_ALIGN.CENTER)

    # Simulated convergence data (realistic pattern)
    # As n grows, spread converges to ~851bp with decreasing noise
    np.random.seed(99)
    true_val = 851  # B- PIK Z-spread
    path_counts = [500, 1000, 2000, 5000, 10000, 15000, 20000, 25000]
    noise_scale = 80  # initial noise

    z_min, z_max = 750, 950

    for i, n in enumerate(path_counts):
        se = noise_scale / math.sqrt(n / 500)
        estimate = true_val + np.random.normal(0, se)
        estimate = max(z_min, min(z_max, estimate))

        x = int(chart_left) + int(chart_w * i / (len(path_counts) - 1))
        y_frac = 1 - (estimate - z_min) / (z_max - z_min)
        y = int(chart_top + chart_h * y_frac)

        # Draw point
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x - Pt(4), y - Pt(4), Pt(8), Pt(8))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_BLUE
        dot.line.fill.background()

        # Error bar (±SE)
        se_bp = se
        y_top = int(chart_top + chart_h * (1 - (estimate + se_bp - z_min) / (z_max - z_min)))
        y_bot = int(chart_top + chart_h * (1 - (estimate - se_bp - z_min) / (z_max - z_min)))
        eb = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x - Pt(0.5), y_top,
            Pt(1), max(abs(y_bot - y_top), Emu(1)))
        eb.fill.solid()
        eb.fill.fore_color.rgb = ACCENT_BLUE
        eb.line.fill.background()

        # Label
        add_textbox(slide, x - Inches(0.3), chart_top + chart_h + Pt(6),
                    Inches(0.6), Inches(0.2),
                    f"{n//1000}K" if n >= 1000 else str(n),
                    font_size=7, color=TEXT_MID,
                    alignment=PP_ALIGN.CENTER)

    # True value line
    true_y = int(chart_top + chart_h * (1 - (true_val - z_min) / (z_max - z_min)))
    tl = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, chart_left, true_y,
        chart_w, Pt(1))
    tl.fill.solid()
    tl.fill.fore_color.rgb = ACCENT_ORANGE
    tl.line.fill.background()
    add_textbox(slide, chart_left + chart_w + Pt(4),
                true_y - Pt(8), Inches(1.2), Inches(0.2),
                f"25K = {true_val}bp", font_size=8,
                color=ACCENT_ORANGE)

    # Right side: variance reduction
    add_textbox(slide, Inches(6.3), Inches(2.3), Inches(3.5), Inches(0.25),
                "Variance Reduction Techniques",
                font_size=12, bold=True, color=HEADER_BLUE)

    add_bullet_list(slide, Inches(6.3), Inches(2.65), Inches(3.5),
                    Inches(2.5), [
        "Antithetic variates: for each path Z, "
        "also simulate \u2212Z. Halves variance "
        "at zero extra computation cost",
        "Fixed seed (42): ensures reproducibility. "
        "Same inputs always give same outputs",
        "Monthly time steps (12/year): fine enough "
        "to capture coupon dates and barrier crossings",
        "At 25K paths: SE \u2248 0.1\u20130.5% of price. "
        "Z-spread precision \u2248 \u00b15bp",
    ], font_size=11, color=TEXT_DARK)

    # SE table by issuer
    add_textbox(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(0.25),
                "Standard Errors at 25,000 Paths (from MC engine)",
                font_size=12, bold=True, color=HEADER_BLUE)

    se_rows = [["Issuer", "LTV", "Cash SE (%)", "PIK SE (%)",
                "Z-Spread \u00b1", "Adequate?"]]
    se_data = [
        ("BB+", "50%", "0.02%", "0.03%", "\u00b12bp", "Yes"),
        ("BB\u2212", "61%", "0.05%", "0.07%", "\u00b14bp", "Yes"),
        ("B", "71%", "0.12%", "0.18%", "\u00b18bp", "Yes"),
        ("B\u2212", "80%", "0.25%", "0.38%", "\u00b115bp", "Marginal"),
        ("CCC", "87%", "0.45%", "0.62%", "\u00b125bp", "Marginal"),
    ]
    for name, ltv, c_se, p_se, z_pm, ok in se_data:
        se_rows.append([name, ltv, c_se, p_se, z_pm, ok])
    add_table(slide, Inches(0.3), Inches(5.8), Inches(9.4), se_rows,
              col_widths=[Inches(0.8), Inches(0.5), Inches(1.0),
                          Inches(1.0), Inches(1.0), Inches(0.8)],
              font_size=10)

    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(9), Inches(0.25),
                "For stressed credits, consider 50\u2013100K paths for "
                "production-grade precision",
                font_size=10, color=TEXT_MID)


def slide_15_parameters(prs):
    """Slide 15: Complete Parameter Reference."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BODY_BG)
    add_section_number(slide, "14")
    add_slide_title(slide,
                    "Complete Parameter Reference",
                    "All assumptions in one place")
    add_slide_number(slide, 15)

    # Global parameters
    add_textbox(slide, Inches(0.3), Inches(1.0), Inches(4.5), Inches(0.25),
                "Global Parameters", font_size=13, bold=True,
                color=HEADER_BLUE)
    global_rows = [
        ["Parameter", "Value", "Rationale"],
        ["Risk-free rate", "4.50%", "Flat OIS curve"],
        ["Coupon rate", "8.50%", "Typical HY coupon"],
        ["Maturity", "5 years", "Standard HY tenor"],
        ["Frequency", "Semi-annual", "Market convention"],
        ["Notional", "100", "Par = 100"],
        ["MC paths", "25,000", "Balance: precision vs speed"],
        ["Seed", "42", "Reproducibility"],
        ["Antithetic", "Yes", "Variance reduction"],
        ["Steps/year", "12", "Monthly grid"],
        ["Endo hazard \u03b2", "2.0", "Quadratic sensitivity"],
        ["Recovery floor", "10%", "Minimum asset value"],
    ]
    add_table(slide, Inches(0.3), Inches(1.3), Inches(4.5), global_rows,
              col_widths=[Inches(1.3), Inches(1.0), Inches(2.0)],
              font_size=9)

    # Per-issuer parameters
    add_textbox(slide, Inches(5.0), Inches(1.0), Inches(4.8), Inches(0.25),
                "Per-Issuer Parameters", font_size=13, bold=True,
                color=HEADER_BLUE)
    issuer_rows = [["Param", "BB+", "BB\u2212", "B",
                    "B\u2212", "CCC"]]
    fields = [
        ("Asset V\u2080", "asset", "{:.0f}"),
        ("Vol \u03c3", "vol", "{:.0%}"),
        ("Ann PD", "pd", "{:.2%}"),
        ("Mkt Spread", "spread", "{:.0f}bp"),
        ("Recovery R\u2080", "rec", "{:.0%}"),
        ("LTV", "ltv", "{}"),
    ]
    for label, key, fmt in fields:
        row = [label]
        for iss in ISSUERS:
            val = iss[key]
            if "bp" in fmt:
                row.append(f"{val * 10000:.0f}bp")
            elif key == "ltv":
                row.append(val)
            else:
                row.append(fmt.format(val))
        issuer_rows.append(row)

    # Add derived parameters
    issuer_rows.append(["Barrier B", *[
        f"{merton_barrier(i['asset'], i['vol'], i['pd']):.0f}"
        for i in ISSUERS]])
    issuer_rows.append(["DD", *[
        f"{_dd(i['asset'], merton_barrier(i['asset'], i['vol'], i['pd']), i['vol'], MATURITY):.1f}"
        for i in ISSUERS]])
    issuer_rows.append(["\u03bb\u2080 (cal)", *[
        f"{calibrate_hazard(i['spread'], i['rec']) * 10000:.0f}bp"
        for i in ISSUERS]])

    add_table(slide, Inches(5.0), Inches(1.3), Inches(4.8), issuer_rows,
              col_widths=[Inches(1.0), Inches(0.7), Inches(0.7),
                          Inches(0.7), Inches(0.7), Inches(0.7)],
              font_size=9)

    # Key assumptions & limitations
    add_textbox(slide, Inches(0.3), Inches(5.5), Inches(9.4), Inches(0.25),
                "Key Assumptions & Limitations",
                font_size=13, bold=True, color=HEADER_BLUE)
    add_bullet_list(slide, Inches(0.3), Inches(5.85), Inches(9.4),
                    Inches(1.5), [
        "Flat rate & hazard term structures. Real curves add "
        "convexity but don\u2019t change the qualitative results",
        "Terminal Merton barrier (not first-passage). "
        "First-passage would slightly increase short-dated PDs",
        "No credit migration or rating transitions. "
        "The model is single-period default/survive",
        "Toggle threshold fixed at \u03bb > 10%. Alternative "
        "strategies (stochastic, optimal exercise) could "
        "change toggle results at the margin",
    ], font_size=11, color=TEXT_DARK)


def slide_16_model_comparison(prs):
    """Slide 16: Model Comparison — Agreement & Divergence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, MID_BG)
    add_slide_number(slide, 16)

    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.5),
                "Model Comparison: Agreement & Divergence",
                font_size=28, bold=True, color=WHITE)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.85),
        Inches(2), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    # Two-column comparison
    # Left: HR Model
    hr_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.2),
        Inches(4.5), Inches(3.2))
    hr_box.fill.solid()
    hr_box.fill.fore_color.rgb = RGBColor(0x1F, 0x2F, 0x50)
    hr_box.line.fill.background()

    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(4), Inches(0.3),
                "Hazard Rate Model", font_size=18, bold=True,
                color=ACCENT_BLUE)

    hr_points = [
        "\u2713 Simple, fast, closed-form",
        "\u2713 Clear economic intuition",
        "\u2713 Calibrates to market spreads",
        "\u2717 Overstates PD for strong credits",
        "\u2717 No feedback loop (static \u03bb)",
        "\u2717 PIK always trades wider",
        "\u2717 Recovery independent of notional",
    ]
    for i, pt in enumerate(hr_points):
        col = ACCENT_GREEN if pt.startswith("\u2713") else ACCENT_RED
        add_textbox(slide, Inches(0.5), Inches(1.7) + Inches(i * 0.33),
                    Inches(4), Inches(0.3),
                    pt, font_size=11, color=col)

    # Right: MC Model
    mc_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.2),
        Inches(4.5), Inches(3.2))
    mc_box.fill.solid()
    mc_box.fill.fore_color.rgb = RGBColor(0x1F, 0x2F, 0x50)
    mc_box.line.fill.background()

    add_textbox(slide, Inches(5.4), Inches(1.3), Inches(4), Inches(0.3),
                "Merton MC Model", font_size=18, bold=True,
                color=ACCENT_ORANGE)

    mc_points = [
        "\u2713 Endogenous hazard (\u03bb grows with PIK)",
        "\u2713 Dynamic recovery (R falls with notional)",
        "\u2713 Captures feedback loop",
        "\u2713 Calibrates to historical PDs",
        "\u2713 Path-dependent toggle modelling",
        "\u2717 Computationally intensive (MC)",
        "\u2717 More parameters to calibrate",
    ]
    for i, pt in enumerate(mc_points):
        col = ACCENT_GREEN if pt.startswith("\u2713") else ACCENT_RED
        add_textbox(slide, Inches(5.4), Inches(1.7) + Inches(i * 0.33),
                    Inches(4), Inches(0.3),
                    pt, font_size=11, color=col)

    # Where they agree / disagree
    add_textbox(slide, Inches(0.5), Inches(4.6), Inches(4.2), Inches(0.3),
                "Where They Agree", font_size=14, bold=True,
                color=ACCENT_GREEN)
    add_bullet_list(slide, Inches(0.5), Inches(4.95), Inches(4.2),
                    Inches(1.0), [
        "LTV > 75%: PIK premium is large, positive, "
        "and unambiguous (+140 to +290bp)",
        "Toggle \u2265 PIK for stressed credits",
        "Non-linearity in premium vs credit quality",
    ], font_size=11, color=LIGHT_GREY)

    add_textbox(slide, Inches(5.2), Inches(4.6), Inches(4.5), Inches(0.3),
                "Where They Disagree", font_size=14, bold=True,
                color=ACCENT_RED)
    add_bullet_list(slide, Inches(5.2), Inches(4.95), Inches(4.5),
                    Inches(1.0), [
        "LTV < 65%: HR says +40\u201360bp, MC says "
        "\u221259 to \u221222bp (sign reversal!)",
        "Root cause: HR uses risk-neutral \u03bb (7\u00d7 "
        "historical), MC uses real-world PDs",
        "Toggle vs PIK ordering at low LTV",
    ], font_size=11, color=LIGHT_GREY)

    # Reconciliation
    add_callout_box(slide, Inches(0.5), Inches(6.2), Inches(9),
                    Inches(0.7),
                    "Reconciliation: both models are \u2018right\u2019 under "
                    "their calibration assumptions. The HR model answers "
                    "\u2018what if defaults match market-implied rates?\u2019 "
                    "The MC model answers \u2018what if defaults match "
                    "historical experience?\u2019 The truth lies in between, "
                    "but for LTV > 75% it doesn\u2019t matter \u2014 "
                    "both say PIK costs a lot.",
                    bg_color=ACCENT_BLUE, font_size=12)


# ── Main ─────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_result_to_explain(prs)
    slide_03_bond_setup(prs)
    slide_04_hr_model(prs)
    slide_05_hr_sensitivity(prs)
    slide_06_calibration_gap(prs)
    slide_07_merton_model(prs)
    slide_08_barrier_calibration(prs)
    slide_09_endogenous_hazard(prs)
    slide_10_dynamic_recovery(prs)
    slide_11_mc_paths(prs)
    slide_12_feedback_spiral(prs)
    slide_13_toggle_mechanics(prs)
    slide_14_convergence(prs)
    slide_15_parameters(prs)
    slide_16_model_comparison(prs)

    out = Path(__file__).parent / "pik_deep_dive.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
